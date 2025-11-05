import streamlit as st
import io
import google.generativeai as genai
from PIL import Image
import datetime
import os
from pymongo import MongoClient
from bson import ObjectId
import json
import hashlib
from google.genai import types
import PyPDF2
from pptx import Presentation
import docx
import openai
from typing import List, Dict, Tuple
import hashlib
import pandas as pd
import re

# Configura√ß√£o inicial
st.set_page_config(
    layout="wide",
    page_title="Agente Social",
    page_icon="ü§ñ"
)

import os
import PyPDF2
import pdfplumber
from pathlib import Path

# --- FUN√á√ïES AUXILIARES MELHORADAS ---

def criar_prompt_validacao_preciso(texto, nome_arquivo, contexto_agente):
    """Cria um prompt de valida√ß√£o muito mais preciso para evitar falsos positivos"""
    
    prompt = f"""
{contexto_agente}


###BEGIN TEXTO PARA VALIDA√á√ÉO###
**Arquivo:** {nome_arquivo}
**Conte√∫do:**
{texto[:12000]}
###END TEXTO PARA VALIDA√á√ÉO###

## FORMATO DE RESPOSTA OBRIGAT√ìRIO:



### ‚úÖ CONFORMIDADE COM DIRETRIZES
- [Itens que est√£o alinhados com as diretrizes de branding]



**INCONSIST√äNCIAS COM BRANDING:**
- [S√≥ liste desvios REAIS das diretrizes de branding]

### üí° TEXTO REVISADO
- [Sugest√µes para aprimorar]

### üìä STATUS FINAL
**Documento:** [Aprovado/Necessita ajustes/Reprovado]
**Principais a√ß√µes necess√°rias:** [Lista resumida]

"""
    return prompt

def analisar_documento_por_slides(doc, contexto_agente):
    """Analisa documento slide por slide com alta precis√£o"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Analisando slide {i+1}..."):
            try:
                prompt_slide = f"""
{contexto_agente}

## AN√ÅLISE POR SLIDE - PRECIS√ÉO ABSOLUTA

###BEGIN TEXTO PARA VALIDA√á√ÉO###
**SLIDE {i+1}:**
{slide['conteudo'][:2000]}
###END TEXTO PARA VALIDA√á√ÉO###


**AN√ÅLISE DO SLIDE {i+1}:**

### ‚úÖ Pontos Fortes:
[O que est√° bom neste slide]

### ‚ö†Ô∏è Problemas REAIS:
- [Lista CURTA de problemas]

### üí° Sugest√µes Espec√≠ficas:
[Melhorias para ESTE slide espec√≠fico]

Considere que slides que s√£o introdut√≥rios ou apenas de t√≠tulos n√£o precisam de tanto rigor de branding

**STATUS:** [‚úîÔ∏è Aprovado / ‚ö†Ô∏è Ajustes Menores / ‚ùå Problemas S√©rios]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '‚ùå' in resposta.text or '‚ö†Ô∏è' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"‚ùå Erro na an√°lise do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relat√≥rio consolidado
    relatorio = f"# üìä RELAT√ìRIO DE VALIDA√á√ÉO - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Altera√ß√µes:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de aten√ß√£o
    slides_com_problemas = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_problemas:
        relatorio += "## üö® SLIDES QUE PRECISAM DE ATEN√á√ÉO:\n\n"
        for resultado in slides_com_problemas:
            relatorio += f"### üìã Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## üìà RESUMO EXECUTIVO\n\n"
    if slides_com_problemas:
        relatorio += f"**‚ö†Ô∏è {len(slides_com_problemas)} slide(s) necessitam de ajustes**\n"
        relatorio += f"**‚úÖ {len(doc['slides']) - len(slides_com_problemas)} slide(s) est√£o adequados**\n"
    else:
        relatorio += "**üéâ Todos os slides est√£o em conformidade com as diretrizes!**\n"
    
    return relatorio

def extract_text_from_pdf_com_slides(arquivo_pdf):
    """Extrai texto de PDF com informa√ß√£o de p√°ginas"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'p√°gina'
            })
        
        texto_completo = "\n\n".join([f"--- P√ÅGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extra√ß√£o PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    """Extrai texto de PPTX com informa√ß√£o de slides"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extra√ß√£o PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    """Extrai texto de arquivos TXT e DOCX"""
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo n√£o suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extra√ß√£o: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file using multiple methods for better coverage
    """
    text = ""

    # Method 1: Try with pdfplumber (better for some PDFs)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to PyPDF2 if pdfplumber didn't extract much text
    if len(text.strip()) < 100:  # If very little text was extracted
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            print(f"PyPDF2 also failed for {pdf_path}: {e}")

    return text
    

# --- Sistema de Autentica√ß√£o MELHORADO ---
def make_hashes(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

def check_hashes(password, hashed_text):
    return make_hashes(password) == hashed_text

# Dados de usu√°rio (em produ√ß√£o, isso deve vir de um banco de dados seguro)
users_db = {
    "admin": {
        "password": make_hashes("senha1234"),
        "squad": "admin",
        "nome": "Administrador"
    }
}

# Conex√£o MongoDB
client = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
db = client['agentes_personalizados']
collection_agentes = db['agentes']
collection_conversas = db['conversas']
collection_usuarios = db['usuarios']  # Nova cole√ß√£o para usu√°rios

# --- FUN√á√ïES DE CADASTRO E LOGIN ---
def criar_usuario(email, senha, nome, squad):
    """Cria um novo usu√°rio no banco de dados"""
    try:
        # Verificar se usu√°rio j√° existe
        if collection_usuarios.find_one({"email": email}):
            return False, "Usu√°rio j√° existe"
        
        # Criar hash da senha
        senha_hash = make_hashes(senha)
        
        novo_usuario = {
            "email": email,
            "senha": senha_hash,
            "nome": nome,
            "squad": squad,
            "data_criacao": datetime.datetime.now(),
            "ultimo_login": None,
            "ativo": True
        }
        
        result = collection_usuarios.insert_one(novo_usuario)
        return True, "Usu√°rio criado com sucesso"
        
    except Exception as e:
        return False, f"Erro ao criar usu√°rio: {str(e)}"

def verificar_login(email, senha):
    """Verifica as credenciais do usu√°rio"""
    try:
        # Primeiro verificar no banco de dados
        usuario = collection_usuarios.find_one({"email": email, "ativo": True})
        
        if usuario:
            if check_hashes(senha, usuario["senha"]):
                # Atualizar √∫ltimo login
                collection_usuarios.update_one(
                    {"_id": usuario["_id"]},
                    {"$set": {"ultimo_login": datetime.datetime.now()}}
                )
                return True, usuario, "Login bem-sucedido"
            else:
                return False, None, "Senha incorreta"
        
        # Fallback para usu√°rios hardcoded (apenas para admin)
        if email in users_db:
            user_data = users_db[email]
            if check_hashes(senha, user_data["password"]):
                usuario_fallback = {
                    "email": email,
                    "nome": user_data["nome"],
                    "squad": user_data["squad"],
                    "_id": "admin"
                }
                return True, usuario_fallback, "Login bem-sucedido"
            else:
                return False, None, "Senha incorreta"
        
        return False, None, "Usu√°rio n√£o encontrado"
        
    except Exception as e:
        return False, None, f"Erro no login: {str(e)}"

def get_current_user():
    """Retorna o usu√°rio atual da sess√£o"""
    return st.session_state.get('user', {})

def get_current_squad():
    """Retorna o squad do usu√°rio atual"""
    user = get_current_user()
    return user.get('squad', 'unknown')

def login():
    """Formul√°rio de login e cadastro"""
    st.title("üîí Agente Social - Login")
    
    tab_login, tab_cadastro = st.tabs(["Login", "Cadastro"])
    
    with tab_login:
        with st.form("login_form"):
            email = st.text_input("Email")
            password = st.text_input("Senha", type="password")
            submit_button = st.form_submit_button("Login")
            
            if submit_button:
                if email and password:
                    sucesso, usuario, mensagem = verificar_login(email, password)
                    if sucesso:
                        st.session_state.logged_in = True
                        st.session_state.user = usuario
                        st.success("Login realizado com sucesso!")
                        st.rerun()
                    else:
                        st.error(mensagem)
                else:
                    st.error("Por favor, preencha todos os campos")
    
    with tab_cadastro:
        with st.form("cadastro_form"):
            st.subheader("Criar Nova Conta")
            
            nome = st.text_input("Nome Completo")
            email = st.text_input("Email")
            squad = st.selectbox(
                "Selecione seu Squad:",
                ["Syngenta", "SME", "Enterprise"],
                help="Escolha o squad ao qual voc√™ pertence"
            )
            senha = st.text_input("Senha", type="password")
            confirmar_senha = st.text_input("Confirmar Senha", type="password")
            
            submit_cadastro = st.form_submit_button("Criar Conta")
            
            if submit_cadastro:
                if not all([nome, email, squad, senha, confirmar_senha]):
                    st.error("Por favor, preencha todos os campos")
                elif senha != confirmar_senha:
                    st.error("As senhas n√£o coincidem")
                elif len(senha) < 6:
                    st.error("A senha deve ter pelo menos 6 caracteres")
                else:
                    sucesso, mensagem = criar_usuario(email, senha, nome, squad)
                    if sucesso:
                        st.success("Conta criada com sucesso! Fa√ßa login para continuar.")
                    else:
                        st.error(mensagem)

# Verificar se o usu√°rio est√° logado
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if not st.session_state.logged_in:
    login()
    st.stop()

# --- CONFIGURA√á√ïES AP√ìS LOGIN ---
gemini_api_key = os.getenv("GEM_API_KEY")
if not gemini_api_key:
    st.error("GEMINI_API_KEY n√£o encontrada nas vari√°veis de ambiente")
    st.stop()

genai.configure(api_key=gemini_api_key)
modelo_vision = genai.GenerativeModel("gemini-2.5-flash", generation_config={"temperature": 0.1})
modelo_texto = genai.GenerativeModel("gemini-2.5-flash")

# Configura√ß√£o da API do Perplexity
perp_api_key = os.getenv("PERP_API_KEY")
if not perp_api_key:
    st.error("PERP_API_KEY n√£o encontrada nas vari√°veis de ambiente")

# --- Configura√ß√£o de Autentica√ß√£o de Administrador ---
def check_admin_password():
    """Retorna True para usu√°rios admin sem verifica√ß√£o de senha."""
    return st.session_state.user.get('squad') == "admin"

# --- FUN√á√ïES CRUD PARA AGENTES (MODIFICADAS PARA SQUADS) ---
def criar_agente(nome, system_prompt, base_conhecimento, comments, planejamento, categoria, squad_permitido, agente_mae_id=None, herdar_elementos=None):
    """Cria um novo agente no MongoDB com squad permitido"""
    agente = {
        "nome": nome,
        "system_prompt": system_prompt,
        "base_conhecimento": base_conhecimento,
        "comments": comments,
        "planejamento": planejamento,
        "categoria": categoria,
        "squad_permitido": squad_permitido,  # Novo campo
        "agente_mae_id": agente_mae_id,
        "herdar_elementos": herdar_elementos or [],
        "data_criacao": datetime.datetime.now(),
        "ativo": True,
        "criado_por": get_current_user().get('email', 'unknown'),
        "criado_por_squad": get_current_squad()  # Novo campo
    }
    result = collection_agentes.insert_one(agente)
    return result.inserted_id

def listar_agentes():
    """Retorna todos os agentes ativos que o usu√°rio atual pode ver"""
    current_squad = get_current_squad()
    
    # Admin v√™ todos os agentes
    if current_squad == "admin":
        return list(collection_agentes.find({"ativo": True}).sort("data_criacao", -1))
    
    # Usu√°rios normais veem apenas agentes do seu squad ou squad "Todos"
    return list(collection_agentes.find({
        "ativo": True,
        "$or": [
            {"squad_permitido": current_squad},
            {"squad_permitido": "Todos"},
            {"criado_por_squad": current_squad}  # Usu√°rio pode ver seus pr√≥prios agentes
        ]
    }).sort("data_criacao", -1))

def listar_agentes_para_heranca(agente_atual_id=None):
    """Retorna todos os agentes ativos que podem ser usados como m√£e (com filtro de squad)"""
    current_squad = get_current_squad()
    
    query = {"ativo": True}
    
    # Filtro por squad
    if current_squad != "admin":
        query["$or"] = [
            {"squad_permitido": current_squad},
            {"squad_permitido": "Todos"},
            {"criado_por_squad": current_squad}
        ]
    
    if agente_atual_id:
        # Excluir o pr√≥prio agente da lista de op√ß√µes para evitar auto-heran√ßa
        if isinstance(agente_atual_id, str):
            agente_atual_id = ObjectId(agente_atual_id)
        query["_id"] = {"$ne": agente_atual_id}
    
    return list(collection_agentes.find(query).sort("data_criacao", -1))

def obter_agente(agente_id):
    """Obt√©m um agente espec√≠fico pelo ID com verifica√ß√£o de permiss√£o por squad"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    agente = collection_agentes.find_one({"_id": agente_id})
    
    # Verificar permiss√£o baseada no squad
    if agente and agente.get('ativo', True):
        current_squad = get_current_squad()
        
        # Admin pode ver tudo
        if current_squad == "admin":
            return agente
        
        # Usu√°rios normais s√≥ podem ver agentes do seu squad ou "Todos"
        squad_permitido = agente.get('squad_permitido')
        criado_por_squad = agente.get('criado_por_squad')
        
        if squad_permitido == current_squad or squad_permitido == "Todos" or criado_por_squad == current_squad:
            return agente
    
    return None

def atualizar_agente(agente_id, nome, system_prompt, base_conhecimento, comments, planejamento, categoria, squad_permitido, agente_mae_id=None, herdar_elementos=None):
    """Atualiza um agente existente com verifica√ß√£o de permiss√£o"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    # Verificar se o usu√°rio tem permiss√£o para editar este agente
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente n√£o encontrado ou sem permiss√£o de edi√ß√£o")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {
            "$set": {
                "nome": nome,
                "system_prompt": system_prompt,
                "base_conhecimento": base_conhecimento,
                "comments": comments,
                "planejamento": planejamento,
                "categoria": categoria,
                "squad_permitido": squad_permitido,  # Novo campo
                "agente_mae_id": agente_mae_id,
                "herdar_elementos": herdar_elementos or [],
                "data_atualizacao": datetime.datetime.now()
            }
        }
    )

def desativar_agente(agente_id):
    """Desativa um agente (soft delete) com verifica√ß√£o de permiss√£o"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    # Verificar se o usu√°rio tem permiss√£o para desativar este agente
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente n√£o encontrado ou sem permiss√£o para desativar")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {"$set": {"ativo": False, "data_desativacao": datetime.datetime.now()}}
    )

def obter_agente_com_heranca(agente_id):
    """Obt√©m um agente com os elementos herdados aplicados"""
    agente = obter_agente(agente_id)
    if not agente or not agente.get('agente_mae_id'):
        return agente
    
    agente_mae = obter_agente(agente['agente_mae_id'])
    if not agente_mae:
        return agente
    
    elementos_herdar = agente.get('herdar_elementos', [])
    agente_completo = agente.copy()
    
    for elemento in elementos_herdar:
        if elemento == 'system_prompt' and not agente_completo.get('system_prompt'):
            agente_completo['system_prompt'] = agente_mae.get('system_prompt', '')
        elif elemento == 'base_conhecimento' and not agente_completo.get('base_conhecimento'):
            agente_completo['base_conhecimento'] = agente_mae.get('base_conhecimento', '')
        elif elemento == 'comments' and not agente_completo.get('comments'):
            agente_completo['comments'] = agente_mae.get('comments', '')
        elif elemento == 'planejamento' and not agente_completo.get('planejamento'):
            agente_completo['planejamento'] = agente_mae.get('planejamento', '')
    
    return agente_completo

def salvar_conversa(agente_id, mensagens, segmentos_utilizados=None):
    """Salva uma conversa no hist√≥rico"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    conversa = {
        "agente_id": agente_id,
        "mensagens": mensagens,
        "segmentos_utilizados": segmentos_utilizados,
        "data_criacao": datetime.datetime.now()
    }
    return collection_conversas.insert_one(conversa)

def obter_conversas(agente_id, limite=10):
    """Obt√©m o hist√≥rico de conversas de um agente"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return list(collection_conversas.find(
        {"agente_id": agente_id}
    ).sort("data_criacao", -1).limit(limite))

# --- Fun√ß√£o para construir contexto com segmentos selecionados ---
def construir_contexto(agente, segmentos_selecionados, historico_mensagens=None):
    """Constr√≥i o contexto com base nos segmentos selecionados"""
    contexto = ""
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto += f"### INSTRU√á√ïES DO SISTEMA ###\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto += f"### BASE DE CONHECIMENTO ###\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto += f"### COMENT√ÅRIOS DO CLIENTE ###\n{agente['comments']}\n\n"
    
    if "planejamento" in segmentos_selecionados and agente.get('planejamento'):
        contexto += f"### PLANEJAMENTO ###\n{agente['planejamento']}\n\n"
    
    # Adicionar hist√≥rico se fornecido
    if historico_mensagens:
        contexto += "### HIST√ìRICO DA CONVERSA ###\n"
        for msg in historico_mensagens:
            contexto += f"{msg['role']}: {msg['content']}\n"
        contexto += "\n"
    
    contexto += "### RESPOSTA ATUAL ###\nassistant:"
    
    return contexto

# --- MODIFICA√á√ÉO: SELECTBOX PARA SELE√á√ÉO DE AGENTE ---
def selecionar_agente_interface():
    """Interface para sele√ß√£o de agente usando selectbox"""
    st.title("ü§ñ Agente Social")
    
    # Carregar agentes dispon√≠veis
    agentes = listar_agentes()
    
    if not agentes:
        st.error("‚ùå Nenhum agente dispon√≠vel. Crie um agente primeiro na aba de Gerenciamento.")
        return None
    
    # Preparar op√ß√µes para o selectbox
    opcoes_agentes = []
    for agente in agentes:
        agente_completo = obter_agente_com_heranca(agente['_id'])
        if agente_completo:  # S√≥ adiciona se tiver permiss√£o
            descricao = f"{agente['nome']} - {agente.get('categoria', 'Social')}"
            if agente.get('agente_mae_id'):
                descricao += " üîó"
            # Adicionar indicador de squad
            squad_permitido = agente.get('squad_permitido', 'Todos')
            descricao += f" üë•{squad_permitido}"
            opcoes_agentes.append((descricao, agente_completo))
    
    if opcoes_agentes:
        # Selectbox para sele√ß√£o de agente
        agente_selecionado_desc = st.selectbox(
            "Selecione uma base de conhecimento para usar o sistema:",
            options=[op[0] for op in opcoes_agentes],
            index=0,
            key="selectbox_agente_principal"
        )
        
        # Encontrar o agente completo correspondente
        agente_completo = None
        for desc, agente in opcoes_agentes:
            if desc == agente_selecionado_desc:
                agente_completo = agente
                break
        
        if agente_completo and st.button("‚úÖ Confirmar Sele√ß√£o", key="confirmar_agente"):
            st.session_state.agente_selecionado = agente_completo
            st.session_state.messages = []
            st.session_state.segmentos_selecionados = ["system_prompt", "base_conhecimento", "comments", "planejamento"]
            st.success(f"‚úÖ Agente '{agente_completo['nome']}' selecionado!")
            st.rerun()
        
        return agente_completo
    else:
        st.info("Nenhum agente dispon√≠vel com as permiss√µes atuais.")
        return None

# --- Verificar se o agente j√° foi selecionado ---
if "agente_selecionado" not in st.session_state:
    st.session_state.agente_selecionado = None

# Se n√£o h√° agente selecionado, mostrar interface de sele√ß√£o
if not st.session_state.agente_selecionado:
    selecionar_agente_interface()
    st.stop()

# --- INTERFACE PRINCIPAL (apenas se agente estiver selecionado) ---
agente_selecionado = st.session_state.agente_selecionado

def is_syn_agent(agent_name):
    """Verifica se o agente √© da baseado no nome"""
    return agent_name and any(keyword in agent_name.upper() for keyword in ['SYN'])

PRODUCT_DESCRIPTIONS = {
    "FORTENZA": "Tratamento de sementes inseticida, focado no Cerrado e posicionado para controle do complexo de lagartas e outras pragas iniciais. Comunica√ß√£o focada no mercado 'on farm' (tratamento feito na fazenda).",
    "ALADE": "Fungicida para controle de doen√ßas em soja, frequentemente posicionado em programa com Mitrion para controle de podrid√µes de vagens e gr√£os.",
    "VERDAVIS": "Inseticida e acaricida composto por PLINAZOLIN¬Æ technology (nova mol√©cula, novo grupo qu√≠mico, modo de a√ß√£o in√©dito) + lambda-cialotrina. KBFs: + mais choque, + mais espectro e + mais dias de controle.",
    "ENGEO PLENO S": "Inseticida de tradi√ß√£o, refer√™ncia no controle de percevejos. Mote: 'Nunca foi sorte. Sempre foi Engeo Pleno S'.",
    "MEGAFOL": "Bioativador da Syn Biologicals. Origem 100% natural (extratos vegetais e de algas Ascophyllum nodosum). Desenvolvido para garantir que a planta alcance todo seu potencial produtivo.",
    "MIRAVIS DUO": "Fungicida da fam√≠lia Miravis. Traz ADEPIDYN technology (novo ingrediente ativo, novo grupo qu√≠mico). Focado no controle de manchas foliares.",
    "AVICTA COMPLETO": "Oferta comercial de tratamento industrial de sementes (TSI). Composto por inseticida, fungicida e nematicida.",
    "MITRION": "Fungicida para controle de doen√ßas em soja, frequentemente posicionado em programa com Alade.",
    "AXIAL": "Herbicida para trigo. Composto por um novo ingrediente ativo. Foco no controle do azev√©m.",
    "CERTANO": "Bionematicida e biofungicida. Composto pela bact√©ria Bacillus velezensis. Controla nematoides e fungos de solo.",
    "MANEJO LIMPO": "Programa da Syn para manejo integrado de plantas daninhas.",
    "ELESTAL NEO": "Fungicida para controle de doen√ßas em soja e algod√£o.",
    "FRONDEO": "Inseticida para cana-de-a√ß√∫car com foco no controle da broca da cana.",
    "FORTENZA ELITE": "Oferta comercial de TSI. Solu√ß√£o robusta contre pragas, doen√ßas e nematoides do Cerrado.",
    "REVERB": "Produto para manejo de doen√ßas em soja e milho com a√ß√£o prolongada ou de espectro amplo.",
    "YIELDON": "Produto focado em maximizar a produtividade das lavouras.",
    "ORONDIS FLEXI": "Fungicida com flexibilidade de uso para controle de requeima, m√≠ldios e manchas.",
    "RIZOLIQ LLI": "Inoculante ou produto para tratamento de sementes que atua na rizosfera.",
    "ARVATICO": "Fungicida ou inseticida com a√ß√£o espec√≠fica para controle de doen√ßas foliares ou pragas.",
    "VERDADERO": "Produto relacionado √† sa√∫de do solo ou nutri√ß√£o vegetal.",
    "MIRAVIS": "Fungicida da fam√≠lia Miravis para controle de doen√ßas.",
    "MIRAVIS PRO": "Fungicida premium da fam√≠lia Miravis para controle avan√ßado de doen√ßas.",
    "INSTIVO": "Lagarticida posicionado como especialista no controle de lagartas do g√™nero Spodoptera.",
    "CYPRESS": "Fungicida posicionado para √∫ltimas aplica√ß√µes na soja, consolidando o manejo de doen√ßas.",
    "CALARIS": "Herbicida composto por atrazina + mesotriona para controle de plantas daninhas no milho.",
    "SPONTA": "Inseticida para algod√£o com PLINAZOLIN¬Æ technology para controle de bicudo e outras pragas.",
    "INFLUX": "Inseticida lagarticida premium para controle de todas as lagartas, especialmente helicoverpa.",
    "JOINER": "Inseticida acaricida com tecnologia PLINAZOLIN para culturas hortifr√∫ti.",
    "DUAL GOLD": "Herbicida para manejo de plantas daninhas.",
}

def extract_product_info(text: str) -> Tuple[str, str, str]:
    """Extrai informa√ß√µes do produto do texto da c√©lula"""
    if not text or not text.strip():
        return None, None, None
    
    text = str(text).strip()
    
    # Remover emojis e marcadores
    clean_text = re.sub(r'[üîµüü†üü¢üî¥üü£üîÉüì≤]', '', text).strip()
    
    # Padr√µes para extra√ß√£o
    patterns = {
        'product': r'\b([A-Z][A-Za-z\s]+(?:PRO|S|NEO|LLI|ELITE|COMPLETO|DUO|FLEXI|PLENO|XTRA)?)\b',
        'culture': r'\b(soja|milho|algod√£o|cana|trigo|HF|caf√©|citrus|batata|mel√£o|uva|tomate|multi)\b',
        'action': r'\b(depoimento|resultados|s√©rie|refor√ßo|controle|lan√ßamento|import√¢ncia|jornada|conceito|v√≠deo|a√ß√£o|diferenciais|awareness|problem√°tica|gloss√°rio|manejo|aplica√ß√£o|posicionamento)\b'
    }
    
    product_match = re.search(patterns['product'], clean_text, re.IGNORECASE)
    culture_match = re.search(patterns['culture'], clean_text, re.IGNORECASE)
    action_match = re.search(patterns['action'], clean_text, re.IGNORECASE)
    
    product = product_match.group(1).strip().upper() if product_match else None
    culture = culture_match.group(0).lower() if culture_match else "multi"
    action = action_match.group(0).lower() if action_match else "conscientiza√ß√£o"
    
    return product, culture, action

def generate_context(content, product_name, culture, action, data_input, formato_principal):
    """Gera o texto de contexto discursivo usando LLM"""
    if not gemini_api_key:
        return "API key do Gemini n√£o configurada. Contexto n√£o dispon√≠vel."
    
    # Determinar m√™s em portugu√™s
    meses = {
        1: "janeiro", 2: "fevereiro", 3: "mar√ßo", 4: "abril",
        5: "maio", 6: "junho", 7: "julho", 8: "agosto",
        9: "setembro", 10: "outubro", 11: "novembro", 12: "dezembro"
    }
    mes = meses[data_input.month]
    
    prompt = f"""
    Como redator especializado em agroneg√≥cio da Syn, elabore um texto contextual discursivo de 3-4 par√°grafos para uma pauta de conte√∫do.

    Informa√ß√µes da pauta:
    - Produto: {product_name}
    - Cultura: {culture}
    - A√ß√£o/tema: {action}
    - M√™s de publica√ß√£o: {mes}
    - Formato principal: {formato_principal}
    - Conte√∫do original: {content}

    Descri√ß√£o do produto: {PRODUCT_DESCRIPTIONS.get(product_name, 'Produto agr√≠cola')}

    Instru√ß√µes:
    - Escreva em formato discursivo e fluido, com 3-4 par√°grafos bem estruturados
    - Mantenha tom t√©cnico mas acess√≠vel, adequado para produtores rurais
    - Contextualize a import√¢ncia do tema para a cultura e √©poca do ano
    - Explique por que este conte√∫do √© relevante neste momento
    - Inclua considera√ß√µes sobre o p√∫blico-alvo e objetivos da comunica√ß√£o
    - N√£o repita literalmente a descri√ß√£o do produto, mas a incorpore naturalmente no texto
    - Use linguagem persuasiva mas factual, baseada em dados t√©cnicos

    Formato: Texto corrido em portugu√™s brasileiro
    """
    
    try:
        response = modelo_texto.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Erro ao gerar contexto: {str(e)}"

def generate_platform_strategy(product_name, culture, action, content):
    """Gera estrat√©gia por plataforma usando Gemini"""
    if not gemini_api_key:
        return "API key do Gemini n√£o configurada. Estrat√©gias por plataforma n√£o dispon√≠veis."
    
    prompt = f"""
    Como especialista em m√≠dias sociais para o agroneg√≥cio, crie uma estrat√©gia de conte√∫do detalhada:

    PRODUTO: {product_name}
    CULTURA: {culture}
    A√á√ÉO: {action}
    CONTE√öDO ORIGINAL: {content}
    DESCRI√á√ÉO DO PRODUTO: {PRODUCT_DESCRIPTIONS.get(product_name, 'Produto agr√≠cola')}

    FORNECER ESTRAT√âGIA PARA:
    - Instagram (Feed, Reels, Stories)
    - Facebook 
    - LinkedIn
    - WhatsApp Business
    - YouTube
    - Portal Mais Agro (blog)

    INCLUIR PARA CADA PLATAFORMA:
    1. Tipo de conte√∫do recomendado
    2. Formato ideal (v√≠deo, carrossel, est√°tico, etc.)
    3. Tom de voz apropriado
    4. CTA espec√≠fico
    5. Melhores pr√°ticas

    Formato: Texto claro com se√ß√µes bem definidas
    """
    
    try:
        response = modelo_texto.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Erro ao gerar estrat√©gia: {str(e)}"

def generate_briefing(content, product_name, culture, action, data_input, formato_principal):
    """Gera um briefing completo em formato de texto puro"""
    description = PRODUCT_DESCRIPTIONS.get(product_name, "Descri√ß√£o do produto n√£o dispon√≠vel.")
    context = generate_context(content, product_name, culture, action, data_input, formato_principal)
    platform_strategy = generate_platform_strategy(product_name, culture, action, content)
    
    briefing = f"""
BRIEFING DE CONTE√öDO - {product_name} - {culture.upper()} - {action.upper()}

CONTEXTO E OBJETIVO
{context}

DESCRI√á√ÉO DO PRODUTO
{description}

ESTRAT√âGIA POR PLATAFORMA
{platform_strategy}

FORMATOS SUGERIDOS
- Instagram: Reels + Stories + Feed post
- Facebook: Carrossel + Link post
- LinkedIn: Artigo + Post informativo
- WhatsApp: Card informativo + Link
- YouTube: Shorts + V√≠deo explicativo
- Portal Mais Agro: Blog post + Webstories

CONTATOS E OBSERVA√á√ïES
- Validar com especialista t√©cnico
- Checar disponibilidade de imagens/v√≠deos
- Incluir CTA para portal Mais Agro
- Seguir guidelines de marca
- Revisar compliance regulat√≥rio

DATA PREVISTA: {data_input.strftime('%d/%m/%Y')}
FORMATO PRINCIPAL: {formato_principal}
"""
    return briefing

# --- Interface Principal ---
st.sidebar.title(f"ü§ñ Bem-vindo, {get_current_user().get('nome', 'Usu√°rio')}!")
st.sidebar.info(f"**Squad:** {get_current_squad()}")
st.sidebar.info(f"**Agente selecionado:** {agente_selecionado['nome']}")

# Bot√£o de logout na sidebar
if st.sidebar.button("üö™ Sair", key="logout_btn"):
    for key in ["logged_in", "user", "admin_password_correct", "admin_user", "agente_selecionado"]:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()

# Bot√£o para trocar agente
if st.sidebar.button("üîÑ Trocar Agente", key="trocar_agente_global"):
    st.session_state.agente_selecionado = None
    st.session_state.messages = []
    st.rerun()

# --- SELECTBOX PARA TROCAR AGENTE ACIMA DAS ABAS ---
st.title("ü§ñ Agente Social")

# Carregar agentes dispon√≠veis
agentes = listar_agentes()

if agentes:
    # Preparar op√ß√µes para o selectbox
    opcoes_agentes = []
    for agente in agentes:
        agente_completo = obter_agente_com_heranca(agente['_id'])
        if agente_completo:  # S√≥ adiciona se tiver permiss√£o
            descricao = f"{agente['nome']} - {agente.get('categoria', 'Social')}"
            if agente.get('agente_mae_id'):
                descricao += " üîó"
            # Adicionar indicador de squad
            squad_permitido = agente.get('squad_permitido', 'Todos')
            descricao += f" üë•{squad_permitido}"
            opcoes_agentes.append((descricao, agente_completo))
    
    if opcoes_agentes:
        # Encontrar o √≠ndice atual
        indice_atual = 0
        for i, (desc, agente) in enumerate(opcoes_agentes):
            if agente['_id'] == st.session_state.agente_selecionado['_id']:
                indice_atual = i
                break
        
        # Selectbox para trocar agente
        col1, col2 = st.columns([3, 1])
        with col1:
            novo_agente_desc = st.selectbox(
                "Selecionar Agente:",
                options=[op[0] for op in opcoes_agentes],
                index=indice_atual,
                key="selectbox_trocar_agente"
            )
        with col2:
            if st.button("üîÑ Trocar", key="botao_trocar_agente"):
                # Encontrar o agente completo correspondente
                for desc, agente in opcoes_agentes:
                    if desc == novo_agente_desc:
                        st.session_state.agente_selecionado = agente
                        st.session_state.messages = []
                        st.success(f"‚úÖ Agente alterado para '{agente['nome']}'!")
                        st.rerun()
                        break
    else:
        st.info("Nenhum agente dispon√≠vel com as permiss√µes atuais.")

# Menu de abas - DETERMINAR QUAIS ABAS MOSTRAR
abas_base = [
    "üí¨ Chat", 
    "‚öôÔ∏è Gerenciar Agentes", 
    "‚úÖ Valida√ß√£o Unificada",
    "‚ú® Gera√ß√£o de Conte√∫do",
    "üìù Revis√£o Ortogr√°fica",
    "Monitoramento de Redes"
]

if is_syn_agent(agente_selecionado['nome']):
    abas_base.append("üìã Briefing")

# Criar abas dinamicamente
tabs = st.tabs(abas_base)

# Mapear abas para suas respectivas funcionalidades
tab_mapping = {}
for i, aba in enumerate(abas_base):
    tab_mapping[aba] = tabs[i]

# --- ABA: CHAT ---
with tab_mapping["üí¨ Chat"]:
    st.header("üí¨ Chat com Agente")
    
    # Inicializar session_state se n√£o existir
    if 'messages' not in st.session_state:
        st.session_state.messages = []
    if 'segmentos_selecionados' not in st.session_state:
        st.session_state.segmentos_selecionados = []
    if 'show_historico' not in st.session_state:
        st.session_state.show_historico = False
    
    agente = st.session_state.agente_selecionado
    st.subheader(f"Conversando com: {agente['nome']}")
    
    # Controles de navega√ß√£o no topo
    col1, col2, col3 = st.columns([1, 1, 1])
    
    with col1:
        if st.button("üìö Carregar Hist√≥rico", key="carregar_historico"):
            st.session_state.show_historico = not st.session_state.show_historico
            st.rerun()
    
    with col2:
        if st.button("üîÑ Limpar Chat", key="limpar_chat"):
            st.session_state.messages = []
            if hasattr(st.session_state, 'historico_contexto'):
                st.session_state.historico_contexto = []
            st.success("Chat limpo!")
            st.rerun()
    
    with col3:
        if st.button("üîÅ Trocar Agente", key="trocar_agente_chat"):
            st.session_state.agente_selecionado = None
            st.session_state.messages = []
            st.session_state.historico_contexto = []
            st.rerun()
    
    # Mostrar se h√° hist√≥rico carregado
    if hasattr(st.session_state, 'historico_contexto') and st.session_state.historico_contexto:
        st.info(f"üìñ Usando hist√≥rico anterior com {len(st.session_state.historico_contexto)} mensagens como contexto")
    
    # Modal para sele√ß√£o de hist√≥rico
    if st.session_state.show_historico:
        with st.expander("üìö Selecionar Hist√≥rico de Conversa", expanded=True):
            conversas_anteriores = obter_conversas(agente['_id'])
            
            if conversas_anteriores:
                for i, conversa in enumerate(conversas_anteriores[:10]):  # √öltimas 10 conversas
                    col_hist1, col_hist2, col_hist3 = st.columns([3, 1, 1])
                    
                    with col_hist1:
                        # CORRE√á√ÉO: Usar get() para evitar KeyError
                        data_display = conversa.get('data_formatada', conversa.get('data', 'Data desconhecida'))
                        mensagens_count = len(conversa.get('mensagens', []))
                        st.write(f"**{data_display}** - {mensagens_count} mensagens")
                    
                    with col_hist2:
                        if st.button("üëÄ Visualizar", key=f"ver_{i}"):
                            st.session_state.conversa_visualizada = conversa.get('mensagens', [])
                    
                    with col_hist3:
                        if st.button("üì• Usar", key=f"usar_{i}"):
                            st.session_state.messages = conversa.get('mensagens', [])
                            st.session_state.historico_contexto = conversa.get('mensagens', [])
                            st.session_state.show_historico = False
                            st.success(f"‚úÖ Hist√≥rico carregado: {len(conversa.get('mensagens', []))} mensagens")
                            st.rerun()
                
                # Visualizar conversa selecionada
                if hasattr(st.session_state, 'conversa_visualizada'):
                    st.subheader("üëÄ Visualiza√ß√£o do Hist√≥rico")
                    for msg in st.session_state.conversa_visualizada[-6:]:  # √öltimas 6 mensagens
                        with st.chat_message(msg.get("role", "user")):
                            st.markdown(msg.get("content", ""))
                    
                    if st.button("Fechar Visualiza√ß√£o", key="fechar_visualizacao"):
                        st.session_state.conversa_visualizada = None
                        st.rerun()
            else:
                st.info("Nenhuma conversa anterior encontrada")
    
    # Mostrar informa√ß√µes de heran√ßa se aplic√°vel
    if 'agente_mae_id' in agente and agente['agente_mae_id']:
        agente_original = obter_agente(agente['_id'])
        if agente_original and agente_original.get('herdar_elementos'):
            st.info(f"üîó Este agente herda {len(agente_original['herdar_elementos'])} elementos do agente m√£e")
    
    # Controles de segmentos na sidebar do chat
    st.sidebar.subheader("üîß Configura√ß√µes do Agente")
    st.sidebar.write("Selecione quais bases de conhecimento usar:")
    
    segmentos_disponiveis = {
        "Prompt do Sistema": "system_prompt",
        "Brand Guidelines": "base_conhecimento", 
        "Coment√°rios do Cliente": "comments",
        "Planejamento": "planejamento"
    }
    
    segmentos_selecionados = []
    for nome, chave in segmentos_disponiveis.items():
        if st.sidebar.checkbox(nome, value=chave in st.session_state.segmentos_selecionados, key=f"seg_{chave}"):
            segmentos_selecionados.append(chave)
    
    st.session_state.segmentos_selecionados = segmentos_selecionados
    
    # Exibir status dos segmentos
    if segmentos_selecionados:
        st.sidebar.success(f"‚úÖ Usando {len(segmentos_selecionados)} segmento(s)")
    else:
        st.sidebar.warning("‚ö†Ô∏è Nenhum segmento selecionado")
    
    # Indicador de posi√ß√£o na conversa
    if len(st.session_state.messages) > 4:
        st.caption(f"üìÑ Conversa com {len(st.session_state.messages)} mensagens")
    
    # CORRE√á√ÉO: Exibir hist√≥rico de mensagens DENTRO do contexto correto
    # Verificar se messages existe e √© iter√°vel
    if hasattr(st.session_state, 'messages') and st.session_state.messages:
        for message in st.session_state.messages:
            # Verificar se message √© um dicion√°rio e tem a chave 'role'
            if isinstance(message, dict) and "role" in message:
                with st.chat_message(message["role"]):
                    st.markdown(message.get("content", ""))
            else:
                # Se a estrutura n√£o for a esperada, pular esta mensagem
                continue
    else:
        # Se n√£o houver mensagens, mostrar estado vazio
        st.info("üí¨ Inicie uma conversa digitando uma mensagem abaixo!")
    
    # Input do usu√°rio
    if prompt := st.chat_input("Digite sua mensagem..."):
        # Adicionar mensagem do usu√°rio ao hist√≥rico
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Construir contexto com segmentos selecionados
        contexto = construir_contexto(
            agente, 
            st.session_state.segmentos_selecionados, 
            st.session_state.messages
        )
        
        # Gerar resposta
        with st.chat_message("assistant"):
            with st.spinner('Pensando...'):
                try:
                    resposta = modelo_texto.generate_content(contexto)
                    st.markdown(resposta.text)
                    
                    # Adicionar ao hist√≥rico
                    st.session_state.messages.append({"role": "assistant", "content": resposta.text})
                    
                    # Salvar conversa com segmentos utilizados
                    salvar_conversa(
                        agente['_id'], 
                        st.session_state.messages,
                        st.session_state.segmentos_selecionados
                    )
                    
                except Exception as e:
                    st.error(f"Erro ao gerar resposta: {str(e)}")

# --- ABA: GERENCIAMENTO DE AGENTES (MODIFICADA PARA SQUADS) ---
with tab_mapping["‚öôÔ∏è Gerenciar Agentes"]:
    st.header("Gerenciamento de Agentes")
    
    # Verificar autentica√ß√£o apenas para gerenciamento
    current_user = get_current_user()
    current_squad = get_current_squad()
    
    if current_squad not in ["admin", "Syngenta", "SME", "Enterprise"]:
        st.warning("Acesso restrito a usu√°rios autorizados")
    else:
        # Para admin, verificar senha adicional
        if current_squad == "admin":
            if not check_admin_password():
                st.warning("Digite a senha de administrador")
            else:
                st.write(f'Bem-vindo administrador!')
        else:
            st.write(f'Bem-vindo {current_user.get("nome", "Usu√°rio")} do squad {current_squad}!')
            
        # Subabas para gerenciamento
        sub_tab1, sub_tab2, sub_tab3 = st.tabs(["Criar Agente", "Editar Agente", "Gerenciar Agentes"])
        
        with sub_tab1:
            st.subheader("Criar Novo Agente")
            
            with st.form("form_criar_agente"):
                nome_agente = st.text_input("Nome do Agente:")
                
                # Sele√ß√£o de categoria - AGORA COM MONITORAMENTO
                categoria = st.selectbox(
                    "Categoria:",
                    ["Social", "SEO", "Conte√∫do", "Monitoramento"],
                    help="Organize o agente por √°rea de atua√ß√£o"
                )
                
                # NOVO: Sele√ß√£o de squad permitido
                squad_permitido = st.selectbox(
                    "Squad Permitido:",
                    ["Todos", "Syngenta", "SME", "Enterprise"],
                    help="Selecione qual squad pode ver e usar este agente"
                )
                
                # Configura√ß√µes espec√≠ficas para agentes de monitoramento
                if categoria == "Monitoramento":
                    st.info("üîç **Agente de Monitoramento**: Este agente ser√° usado apenas na aba de Monitoramento de Redes e ter√° uma estrutura simplificada.")
                    
                    # Para monitoramento, apenas base de conhecimento
                    base_conhecimento = st.text_area(
                        "Base de Conhecimento para Monitoramento:", 
                        height=300,
                        placeholder="""Cole aqui a base de conhecimento espec√≠fica para monitoramento de redes sociais.

PERSONALIDADE: Especialista t√©cnico do agroneg√≥cio com habilidade social - "Especialista que fala como gente"

TOM DE VOZ:
- T√©cnico, confi√°vel e seguro, mas acess√≠vel
- Evita exageros e promessas vazias
- Sempre embasado em fatos e ci√™ncia
- Frases curtas e diretas, mais simp√°ticas
- Toque de leveza e ironia pontual quando o contexto permite

PRODUTOS SYN:
- Fortenza: Tratamento de sementes inseticida para Cerrado
- Verdatis: Inseticida com tecnologia PLINAZOLIN
- Megafol: Bioativador natural
- Miravis Duo: Fungicida para controle de manchas foliares

DIRETRIZES:
- N√ÉO inventar informa√ß√µes t√©cnicas
- Sempre basear respostas em fatos
- Manter tom profissional mas acess√≠vel
- Adaptar resposta ao tipo de pergunta""",
                        help="Esta base ser√° usada exclusivamente para monitoramento de redes sociais"
                    )
                    
                    # Campos espec√≠ficos ocultos para monitoramento
                    system_prompt = ""
                    comments = ""
                    planejamento = ""
                    criar_como_filho = False
                    agente_mae_id = None
                    herdar_elementos = []
                    
                else:
                    # Para outras categorias, manter estrutura original
                    criar_como_filho = st.checkbox("Criar como agente filho (herdar elementos)")
                    
                    agente_mae_id = None
                    herdar_elementos = []
                    
                    if criar_como_filho:
                        # Listar TODOS os agentes dispon√≠veis para heran√ßa (exceto monitoramento)
                        agentes_mae = listar_agentes_para_heranca()
                        agentes_mae = [agente for agente in agentes_mae if agente.get('categoria') != 'Monitoramento']
                        
                        if agentes_mae:
                            agente_mae_options = {f"{agente['nome']} ({agente.get('categoria', 'Social')})": agente['_id'] for agente in agentes_mae}
                            agente_mae_selecionado = st.selectbox(
                                "Agente M√£e:",
                                list(agente_mae_options.keys()),
                                help="Selecione o agente do qual este agente ir√° herdar elementos"
                            )
                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                            
                            st.subheader("Elementos para Herdar")
                            herdar_elementos = st.multiselect(
                                "Selecione os elementos a herdar do agente m√£e:",
                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                help="Estes elementos ser√£o herdados do agente m√£e se n√£o preenchidos abaixo"
                            )
                        else:
                            st.info("Nenhum agente dispon√≠vel para heran√ßa. Crie primeiro um agente m√£e.")
                    
                    system_prompt = st.text_area("Prompt de Sistema:", height=150, 
                                                placeholder="Ex: Voc√™ √© um assistente especializado em...",
                                                help="Deixe vazio se for herdar do agente m√£e")
                    base_conhecimento = st.text_area("Brand Guidelines:", height=200,
                                                   placeholder="Cole aqui informa√ß√µes, diretrizes, dados...",
                                                   help="Deixe vazio se for herdar do agente m√£e")
                    comments = st.text_area("Coment√°rios do cliente:", height=200,
                                                   placeholder="Cole aqui os coment√°rios de ajuste do cliente (Se houver)",
                                                   help="Deixe vazio se for herdar do agente m√£e")
                    planejamento = st.text_area("Planejamento:", height=200,
                                               placeholder="Estrat√©gias, planejamentos, cronogramas...",
                                               help="Deixe vazio se for herdar do agente m√£e")
                
                submitted = st.form_submit_button("Criar Agente")
                if submitted:
                    if nome_agente:
                        agente_id = criar_agente(
                            nome_agente, 
                            system_prompt, 
                            base_conhecimento, 
                            comments, 
                            planejamento,
                            categoria,
                            squad_permitido,  # Novo campo
                            agente_mae_id if criar_como_filho else None,
                            herdar_elementos if criar_como_filho else []
                        )
                        st.success(f"Agente '{nome_agente}' criado com sucesso na categoria {categoria} para o squad {squad_permitido}!")
                    else:
                        st.error("Nome √© obrigat√≥rio!")
        
        with sub_tab2:
            st.subheader("Editar Agente Existente")
            
            agentes = listar_agentes()
            if agentes:
                agente_options = {agente['nome']: agente for agente in agentes}
                agente_selecionado_nome = st.selectbox("Selecione o agente para editar:", 
                                                     list(agente_options.keys()))
                
                if agente_selecionado_nome:
                    agente = agente_options[agente_selecionado_nome]
                    
                    with st.form("form_editar_agente"):
                        novo_nome = st.text_input("Nome do Agente:", value=agente['nome'])
                        
                        # Categoria - AGORA COM MONITORAMENTO
                        categorias_disponiveis = ["Social", "SEO", "Conte√∫do", "Monitoramento"]
                        if agente.get('categoria') in categorias_disponiveis:
                            index_categoria = categorias_disponiveis.index(agente.get('categoria', 'Social'))
                        else:
                            index_categoria = 0
                            
                        nova_categoria = st.selectbox(
                            "Categoria:",
                            categorias_disponiveis,
                            index=index_categoria,
                            help="Organize o agente por √°rea de atua√ß√£o"
                        )
                        
                        # NOVO: Squad permitido
                        squads_disponiveis = ["Todos", "Syngenta", "SME", "Enterprise"]
                        squad_atual = agente.get('squad_permitido', 'Todos')
                        if squad_atual in squads_disponiveis:
                            index_squad = squads_disponiveis.index(squad_atual)
                        else:
                            index_squad = 0
                            
                        novo_squad_permitido = st.selectbox(
                            "Squad Permitido:",
                            squads_disponiveis,
                            index=index_squad,
                            help="Selecione qual squad pode ver e usar este agente"
                        )
                        
                        # Interface diferente para agentes de monitoramento
                        if nova_categoria == "Monitoramento":
                            st.info("üîç **Agente de Monitoramento**: Este agente ser√° usado apenas na aba de Monitoramento de Redes.")
                            
                            # Para monitoramento, apenas base de conhecimento
                            nova_base = st.text_area(
                                "Base de Conhecimento para Monitoramento:", 
                                value=agente.get('base_conhecimento', ''),
                                height=300,
                                help="Esta base ser√° usada exclusivamente para monitoramento de redes sociais"
                            )
                            
                            # Campos espec√≠ficos ocultos para monitoramento
                            novo_prompt = ""
                            nova_comment = ""
                            novo_planejamento = ""
                            agente_mae_id = None
                            herdar_elementos = []
                            
                            # Remover heran√ßa se existir
                            if agente.get('agente_mae_id'):
                                st.warning("‚ö†Ô∏è Agentes de monitoramento n√£o suportam heran√ßa. A heran√ßa ser√° removida.")
                            
                        else:
                            # Para outras categorias, manter estrutura original
                            
                            # Informa√ß√µes de heran√ßa (apenas se n√£o for monitoramento)
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.info(f"üîó Este agente √© filho de: {agente_mae['nome']}")
                                    st.write(f"Elementos herdados: {', '.join(agente.get('herdar_elementos', []))}")
                            
                            # Op√ß√£o para tornar independente
                            if agente.get('agente_mae_id'):
                                tornar_independente = st.checkbox("Tornar agente independente (remover heran√ßa)")
                                if tornar_independente:
                                    agente_mae_id = None
                                    herdar_elementos = []
                                else:
                                    agente_mae_id = agente.get('agente_mae_id')
                                    herdar_elementos = agente.get('herdar_elementos', [])
                            else:
                                agente_mae_id = None
                                herdar_elementos = []
                                # Op√ß√£o para adicionar heran√ßa
                                adicionar_heranca = st.checkbox("Adicionar heran√ßa de agente m√£e")
                                if adicionar_heranca:
                                    # Listar TODOS os agentes dispon√≠veis para heran√ßa (excluindo o pr√≥prio e monitoramento)
                                    agentes_mae = listar_agentes_para_heranca(agente['_id'])
                                    agentes_mae = [agente_mae for agente_mae in agentes_mae if agente_mae.get('categoria') != 'Monitoramento']
                                    
                                    if agentes_mae:
                                        agente_mae_options = {f"{agente_mae['nome']} ({agente_mae.get('categoria', 'Social')})": agente_mae['_id'] for agente_mae in agentes_mae}
                                        if agente_mae_options:
                                            agente_mae_selecionado = st.selectbox(
                                                "Agente M√£e:",
                                                list(agente_mae_options.keys()),
                                                help="Selecione o agente do qual este agente ir√° herdar elementos"
                                            )
                                            agente_mae_id = agente_mae_options[agente_mae_selecionado]
                                            herdar_elementos = st.multiselect(
                                                "Elementos para herdar:",
                                                ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                                default=herdar_elementos
                                            )
                                        else:
                                            st.info("Nenhum agente dispon√≠vel para heran√ßa.")
                                    else:
                                        st.info("Nenhum agente dispon√≠vel para heran√ßa.")
                            
                            novo_prompt = st.text_area("Prompt de Sistema:", value=agente['system_prompt'], height=150)
                            nova_base = st.text_area("Brand Guidelines:", value=agente.get('base_conhecimento', ''), height=200)
                            nova_comment = st.text_area("Coment√°rios:", value=agente.get('comments', ''), height=200)
                            novo_planejamento = st.text_area("Planejamento:", value=agente.get('planejamento', ''), height=200)
                        
                        submitted = st.form_submit_button("Atualizar Agente")
                        if submitted:
                            if novo_nome:
                                atualizar_agente(
                                    agente['_id'], 
                                    novo_nome, 
                                    novo_prompt, 
                                    nova_base, 
                                    nova_comment, 
                                    novo_planejamento,
                                    nova_categoria,
                                    novo_squad_permitido,  # Novo campo
                                    agente_mae_id,
                                    herdar_elementos
                                )
                                st.success(f"Agente '{novo_nome}' atualizado com sucesso!")
                                st.rerun()
                            else:
                                st.error("Nome √© obrigat√≥rio!")
            else:
                st.info("Nenhum agente criado ainda.")
        
        with sub_tab3:
            st.subheader("Gerenciar Agentes")
            
            # Mostrar informa√ß√µes do usu√°rio atual
            current_squad = get_current_squad()
            if current_squad == "admin":
                st.info("üëë Modo Administrador: Visualizando todos os agentes do sistema")
            else:
                st.info(f"üë§ Visualizando agentes do squad {current_squad} e squad 'Todos'")
            
            # Filtros por categoria - AGORA COM MONITORAMENTO
            categorias = ["Todos", "Social", "SEO", "Conte√∫do", "Monitoramento"]
            categoria_filtro = st.selectbox("Filtrar por categoria:", categorias)
            
            agentes = listar_agentes()
            
            # Aplicar filtro
            if categoria_filtro != "Todos":
                agentes = [agente for agente in agentes if agente.get('categoria') == categoria_filtro]
            
            if agentes:
                for i, agente in enumerate(agentes):
                    with st.expander(f"{agente['nome']} - {agente.get('categoria', 'Social')} - Squad: {agente.get('squad_permitido', 'Todos')} - Criado em {agente['data_criacao'].strftime('%d/%m/%Y')}"):
                        
                        # Mostrar propriet√°rio se for admin
                        owner_info = ""
                        if current_squad == "admin" and agente.get('criado_por'):
                            owner_info = f" | üë§ {agente['criado_por']}"
                            st.write(f"**Propriet√°rio:** {agente['criado_por']}")
                            st.write(f"**Squad do Criador:** {agente.get('criado_por_squad', 'N/A')}")
                        
                        # Mostrar informa√ß√µes espec√≠ficas por categoria
                        if agente.get('categoria') == 'Monitoramento':
                            st.info("üîç **Agente de Monitoramento** - Usado apenas na aba de Monitoramento de Redes")
                            
                            if agente.get('base_conhecimento'):
                                st.write(f"**Base de Conhecimento:** {agente['base_conhecimento'][:200]}...")
                            else:
                                st.warning("‚ö†Ô∏è Base de conhecimento n√£o configurada")
                            
                            # Agentes de monitoramento n√£o mostram outros campos
                            st.write("**System Prompt:** (N√£o utilizado em monitoramento)")
                            st.write("**Coment√°rios:** (N√£o utilizado em monitoramento)")
                            st.write("**Planejamento:** (N√£o utilizado em monitoramento)")
                            
                        else:
                            # Para outras categorias, mostrar estrutura completa
                            if agente.get('agente_mae_id'):
                                agente_mae = obter_agente(agente['agente_mae_id'])
                                if agente_mae:
                                    st.write(f"**üîó Herda de:** {agente_mae['nome']}")
                                    st.write(f"**Elementos herdados:** {', '.join(agente.get('herdar_elementos', []))}")
                            
                            st.write(f"**Prompt de Sistema:** {agente['system_prompt'][:100]}..." if agente['system_prompt'] else "**Prompt de Sistema:** (herdado ou vazio)")
                            if agente.get('base_conhecimento'):
                                st.write(f"**Brand Guidelines:** {agente['base_conhecimento'][:200]}...")
                            if agente.get('comments'):
                                st.write(f"**Coment√°rios do cliente:** {agente['comments'][:200]}...")
                            if agente.get('planejamento'):
                                st.write(f"**Planejamento:** {agente['planejamento'][:200]}...")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("Selecionar para Chat", key=f"select_{i}"):
                                agente_completo = obter_agente_com_heranca(agente['_id'])
                                st.session_state.agente_selecionado = agente_completo
                                st.session_state.messages = []
                                st.success(f"Agente '{agente['nome']}' selecionado!")
                                st.rerun()
                        with col2:
                            if st.button("Desativar", key=f"delete_{i}"):
                                desativar_agente(agente['_id'])
                                st.success(f"Agente '{agente['nome']}' desativado!")
                                st.rerun()
            else:
                st.info("Nenhum agente encontrado para esta categoria.")

if "üìã Briefing" in tab_mapping:
    with tab_mapping["üìã Briefing"]:
        st.header("üìã Gerador de Briefings - SYN")
        st.markdown("Digite o conte√∫do da c√©lula do calend√°rio para gerar um briefing completo no padr√£o SYN.")
        
        # Abas para diferentes modos de opera√ß√£o
        tab1, tab2 = st.tabs(["Briefing Individual", "Processamento em Lote (CSV)"])
        
        with tab1:
            st.markdown("### Digite o conte√∫do da c√©lula do calend√°rio")

            content_input = st.text_area(
                "Conte√∫do da c√©lula:",
                placeholder="Ex: megafol - s√©rie - potencial m√°ximo, todo o tempo",
                height=100,
                help="Cole aqui o conte√∫do exato da c√©lula do calend√°rio do Sheets",
                key="individual_content"
            )

            # Campos opcionais para ajuste
            col1, col2 = st.columns(2)

            with col1:
                data_input = st.date_input("Data prevista:", value=datetime.datetime.now(), key="individual_date")

            with col2:
                formato_principal = st.selectbox(
                    "Formato principal:",
                    ["Reels + capa", "Carrossel + stories", "Blog + redes", "V√≠deo + stories", "Multiplataforma"],
                    key="individual_format"
                )

            generate_btn = st.button("Gerar Briefing Individual", type="primary", key="individual_btn")

            # Processamento e exibi√ß√£o do briefing individual
            if generate_btn and content_input:
                with st.spinner("Analisando conte√∫do e gerando briefing..."):
                    # Extrair informa√ß√µes do produto
                    product, culture, action = extract_product_info(content_input)
                    
                    if product and product in PRODUCT_DESCRIPTIONS:
                        # Gerar briefing completo
                        briefing = generate_briefing(content_input, product, culture, action, data_input, formato_principal)
                        
                        # Exibir briefing
                        st.markdown("## Briefing Gerado")
                        st.text(briefing)
                        
                        # Bot√£o de download
                        st.download_button(
                            label="Baixar Briefing",
                            data=briefing,
                            file_name=f"briefing_{product}_{data_input.strftime('%Y%m%d')}.txt",
                            mime="text/plain",
                            key="individual_download"
                        )
                        
                        # Informa√ß√µes extras
                        with st.expander("Informa√ß√µes Extra√≠das"):
                            st.write(f"Produto: {product}")
                            st.write(f"Cultura: {culture}")
                            st.write(f"A√ß√£o: {action}")
                            st.write(f"Data: {data_input.strftime('%d/%m/%Y')}")
                            st.write(f"Formato principal: {formato_principal}")
                            st.write(f"Descri√ß√£o: {PRODUCT_DESCRIPTIONS[product]}")
                            
                    elif product:
                        st.warning(f"Produto '{product}' n√£o encontrado no dicion√°rio. Verifique a grafia.")
                        st.info("Produtos dispon√≠veis: " + ", ".join(list(PRODUCT_DESCRIPTIONS.keys())[:10]) + "...")
                    else:
                        st.error("N√£o foi poss√≠vel identificar um produto no conte√∫do. Tente formatos como:")
                        st.code("""
                        megafol - s√©rie - potencial m√°ximo, todo o tempo
                        verdavis - soja - depoimento produtor
                        engeo pleno s - milho - controle percevejo
                        miravis duo - algod√£o - refor√ßo preventivo
                        """)

        with tab2:
            st.markdown("### Processamento em Lote via CSV")
            
            st.info("""
            Fa√ßa upload de um arquivo CSV exportado do Google Sheets.
            O sistema ir√° processar cada linha a partir da segunda linha (ignorando cabe√ßalhos)
            e gerar briefings apenas para as linhas que cont√™m produtos reconhecidos.
            """)
            
            uploaded_file = st.file_uploader(
                "Escolha o arquivo CSV", 
                type=['csv'],
                help="Selecione o arquivo CSV exportado do Google Sheets"
            )
            
            if uploaded_file is not None:
                try:
                    # Ler o CSV
                    df = pd.read_csv(uploaded_file)
                    st.success(f"CSV carregado com sucesso! {len(df)} linhas encontradas.")
                    
                    # Mostrar pr√©via do arquivo
                    with st.expander("Visualizar primeiras linhas do CSV"):
                        st.dataframe(df.head())
                    
                    # Configura√ß√µes para processamento em lote
                    st.markdown("### Configura√ß√µes do Processamento em Lote")
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        data_padrao = st.date_input(
                            "Data padr√£o para todos os briefings:",
                            value=datetime.datetime.now(),
                            key="batch_date"
                        )
                    
                    with col2:
                        formato_padrao = st.selectbox(
                            "Formato principal padr√£o:",
                            ["Reels + capa", "Carrossel + stories", "Blog + redes", "V√≠deo + stories", "Multiplataforma"],
                            key="batch_format"
                        )
                    
                    # Identificar coluna com conte√∫do
                    colunas = df.columns.tolist()
                    coluna_conteudo = st.selectbox(
                        "Selecione a coluna que cont√©m o conte√∫do das c√©lulas:",
                        colunas,
                        help="Selecione a coluna que cont√©m os textos das c√©lulas do calend√°rio"
                    )
                    
                    processar_lote = st.button("Processar CSV e Gerar Briefings", type="primary", key="batch_btn")
                    
                    if processar_lote:
                        briefings_gerados = []
                        linhas_processadas = 0
                        linhas_com_produto = 0
                        
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        for index, row in df.iterrows():
                            linhas_processadas += 1
                            progress_bar.progress(linhas_processadas / len(df))
                            status_text.text(f"Processando linha {linhas_processadas} de {len(df)}...")
                            
                            # Pular a primeira linha (cabe√ßalhos)
                            if index == 0:
                                continue
                            
                            # Obter conte√∫do da c√©lula
                            content = str(row[coluna_conteudo]) if pd.notna(row[coluna_conteudo]) else ""
                            
                            if content:
                                # Extrair informa√ß√µes do produto
                                product, culture, action = extract_product_info(content)
                                
                                if product and product in PRODUCT_DESCRIPTIONS:
                                    linhas_com_produto += 1
                                    # Gerar briefing
                                    briefing = generate_briefing(
                                        content, 
                                        product, 
                                        culture, 
                                        action, 
                                        data_padrao, 
                                        formato_padrao
                                    )
                                    
                                    briefings_gerados.append({
                                        'linha': index + 1,
                                        'produto': product,
                                        'conteudo': content,
                                        'briefing': briefing,
                                        'arquivo': f"briefing_{product}_{index+1}.txt"
                                    })
                        
                        progress_bar.empty()
                        status_text.empty()
                        
                        # Resultados do processamento
                        st.success(f"Processamento conclu√≠do! {linhas_com_produto} briefings gerados de {linhas_processadas-1} linhas processadas.")
                        
                        if briefings_gerados:
                            # Exibir resumo
                            st.markdown("### Briefings Gerados")
                            resumo_df = pd.DataFrame([{
                                'Linha': b['linha'],
                                'Produto': b['produto'],
                                'Conte√∫do': b['conteudo'][:50] + '...' if len(b['conteudo']) > 50 else b['conteudo']
                            } for b in briefings_gerados])
                            
                            st.dataframe(resumo_df)
                            
                            # Criar arquivo ZIP com todos os briefings
                            import zipfile
                            from io import BytesIO
                            
                            zip_buffer = BytesIO()
                            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                                for briefing_info in briefings_gerados:
                                    zip_file.writestr(
                                        briefing_info['arquivo'], 
                                        briefing_info['briefing']
                                    )
                            
                            zip_buffer.seek(0)
                            
                            # Bot√£o para download do ZIP
                            st.download_button(
                                label="üì• Baixar Todos os Briefings (ZIP)",
                                data=zip_buffer,
                                file_name="briefings_syn.zip",
                                mime="application/zip",
                                key="batch_download_zip"
                            )
                            
                            # Tamb√©m permitir download individual
                            st.markdown("---")
                            st.markdown("### Download Individual")
                            
                            for briefing_info in briefings_gerados:
                                col1, col2 = st.columns([3, 1])
                                with col1:
                                    st.text(f"Linha {briefing_info['linha']}: {briefing_info['produto']} - {briefing_info['conteudo'][:30]}...")
                                with col2:
                                    st.download_button(
                                        label="üìÑ Baixar",
                                        data=briefing_info['briefing'],
                                        file_name=briefing_info['arquivo'],
                                        mime="text/plain",
                                        key=f"download_{briefing_info['linha']}"
                                    )
                        else:
                            st.warning("Nenhum briefing foi gerado. Verifique se o CSV cont√©m produtos reconhecidos.")
                            st.info("Produtos reconhecidos: " + ", ".join(list(PRODUCT_DESCRIPTIONS.keys())[:15]) + "...")
                            
                except Exception as e:
                    st.error(f"Erro ao processar o arquivo CSV: {str(e)}")

        # Se√ß√£o de exemplos
        with st.expander("Exemplos de Conte√∫do", expanded=True):
            st.markdown("""
            Formatos Reconhecidos:

            Padr√£o: PRODUTO - CULTURA - A√á√ÉO ou PRODUTO - A√á√ÉO

            Exemplos:
            - megafol - s√©rie - potencial m√°ximo, todo o tempo
            - verdavis - milho - resultados do produto
            - engeo pleno s - soja - resultados GTEC
            - miravis duo - algod√£o - depoimento produtor
            - axial - trigo - refor√ßo p√≥s-emergente
            - manejo limpo - import√¢ncia manejo antecipado
            - certano HF - a jornada de certano
            - elestal neo - soja - depoimento de produtor
            - fortenza - a jornada da semente mais forte - EP 01
            - reverb - v√≠deo conceito
            """)

        # Lista de produtos reconhecidos
        with st.expander("Produtos Reconhecidos"):
            col1, col2, col3 = st.columns(3)
            products = list(PRODUCT_DESCRIPTIONS.keys())
            
            with col1:
                for product in products[:10]:
                    st.write(f"‚Ä¢ {product}")
            
            with col2:
                for product in products[10:20]:
                    st.write(f"‚Ä¢ {product}")
            
            with col3:
                for product in products[20:]:
                    st.write(f"‚Ä¢ {product}")

        # Rodap√©
        st.markdown("---")
        st.caption("Ferramenta de gera√ß√£o autom√°tica de briefings - Padr√£o SYN. Digite o conte√∫do da c√©lula do calend√°rio para gerar briefings completos.")

def criar_analisadores_especialistas(contexto_agente, contexto_global):
    """Cria prompts especializados para cada √°rea de an√°lise"""
    
    analisadores = {
        'ortografia': {
            'nome': 'üî§ Especialista em Ortografia e Gram√°tica',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM ORTOGRAFIA E GRAM√ÅTICA PORTUGU√äS BR

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos ortogr√°ficos e gramaticais.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Ortografia** - Erros de escrita
2. **Gram√°tica** - Concord√¢ncia, reg√™ncia, coloca√ß√£o
3. **Pontua√ß√£o** - Uso de v√≠rgulas, pontos, etc.
4. **Acentua√ß√£o** - Erros de acentua√ß√£o
5. **Padr√£o Culto** - Conformidade com norma culta

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üî§ RELAT√ìRIO ORTOGR√ÅFICO

### ‚úÖ ACERTOS
- [Itens corretos]

### ‚ùå ERROS IDENTIFICADOS
- [Lista espec√≠fica de erros com corre√ß√µes]

### üìä SCORE ORTOGR√ÅFICO: [X/10]

### üí° SUGEST√ïES DE MELHORIA
- [Recomenda√ß√µes espec√≠ficas]
"""
        },
        'lexico': {
            'nome': 'üìö Especialista em L√©xico e Vocabul√°rio',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM L√âXICO E VOCABUL√ÅRIO

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos lexicais e de vocabul√°rio.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Variedade Lexical** - Riqueza de vocabul√°rio
2. **Precis√£o Sem√¢ntica** - Uso adequado das palavras
3. **Repeti√ß√£o** - Palavras ou express√µes repetidas em excesso
4. **Jarg√µes** - Uso inadequado de termos t√©cnicos
5. **Clareza** - Facilidade de compreens√£o

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üìö RELAT√ìRIO LEXICAL

### ‚úÖ VOCABUL√ÅRIO ADEQUADO
- [Pontos fortes do vocabul√°rio]

### ‚ö†Ô∏è ASPECTOS A MELHORAR
- [Problemas lexicais identificados]

### üîÑ SUGEST√ïES DE SIN√îNIMOS
- [Palavras para substituir]

### üìä SCORE LEXICAL: [X/10]
"""
        },
        'branding': {
            'nome': 'üé® Especialista em Branding e Identidade',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM BRANDING E IDENTIDADE

**Sua tarefa:** Analisar EXCLUSIVAMENTE conformidade com diretrizes de branding.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Tom de Voz** - Alinhamento com personalidade da marca
2. **Mensagem Central** - Consist√™ncia da mensagem
3. **Valores da Marca** - Reflexo dos valores organizacionais
4. **P√∫blico-Alvo** - Adequa√ß√£o ao p√∫blico pretendido
5. **Diferencia√ß√£o** - Elementos √∫nicos da marca

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üé® RELAT√ìRIO DE BRANDING

### ‚úÖ ALINHAMENTOS
- [Elementos que seguem as diretrizes]

### ‚ùå DESVIOS IDENTIFICADOS
- [Elementos fora do padr√£o da marca]

### üìä SCORE BRANDING: [X/10]

### üí° RECOMENDA√á√ïES ESTRAT√âGICAS
- [Sugest√µes para melhor alinhamento]
"""
        
        
        }
    }
    
    return analisadores

def executar_analise_especializada(texto, nome_arquivo, analisadores):
    """Executa an√°lise com m√∫ltiplos especialistas"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN TEXTO PARA AN√ÅLISE###
**Arquivo:** {nome_arquivo}
**Conte√∫do:**
{texto[:8000]}
###END TEXTO PARA AN√ÅLISE###

Por favor, forne√ßa sua an√°lise no formato solicitado.
"""
                
                resposta = modelo_texto.generate_content(prompt_completo)
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': resposta.text,
                    'score': extrair_score(resposta.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"‚ùå Erro na an√°lise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def extrair_score(texto_analise):
    """Extrai score num√©rico do texto de an√°lise"""
    import re
    padrao = r'SCORE.*?\[(\d+)(?:/10)?\]'
    correspondencias = re.findall(padrao, texto_analise, re.IGNORECASE)
    if correspondencias:
        return int(correspondencias[0])
    return 5  # Score padr√£o se n√£o encontrar

def gerar_relatorio_consolidado(resultados_especialistas, nome_arquivo):
    """Gera relat√≥rio consolidado a partir das an√°lises especializadas"""
    
    # Calcular score m√©dio
    scores = [resultado['score'] for resultado in resultados_especialistas.values() if resultado['score'] > 0]
    score_medio = sum(scores) / len(scores) if scores else 0
    
    # Determinar status geral
    if score_medio >= 8:
        status = "‚úÖ APROVADO"
        cor_status = "green"
    elif score_medio >= 6:
        status = "‚ö†Ô∏è AJUSTES MENORES"
        cor_status = "orange"
    else:
        status = "‚ùå REPROVADO"
        cor_status = "red"
    
    relatorio = f"""
# üìä RELAT√ìRIO CONSOLIDADO DE VALIDA√á√ÉO

**Documento:** {nome_arquivo}
**Status Geral:** <span style='color:{cor_status}'>{status}</span>
**Score M√©dio:** {score_medio:.1f}/10
**Data da An√°lise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## üéñÔ∏è SCORES POR √ÅREA
"""
    
    # Adicionar scores individuais
    for area, resultado in resultados_especialistas.items():
        emoji = "‚úÖ" if resultado['score'] >= 8 else "‚ö†Ô∏è" if resultado['score'] >= 6 else "‚ùå"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## üìã AN√ÅLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar an√°lises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## üöÄ RESUMO EXECUTIVO

**Status Final:** {status}
**Score Geral:** {score_medio:.1f}/10

### üéØ PR√ìXIMOS PASSOS RECOMENDADOS:
"""
    
    # Recomenda√ß√µes baseadas nos scores
    areas_baixas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] < 6]
    if areas_baixas:
        relatorio += f"- **Prioridade:** Focar em {', '.join(areas_baixas)}\n"
    
    areas_medianas = [area for area, resultado in resultados_especialistas.items() if 6 <= resultado['score'] < 8]
    if areas_medianas:
        relatorio += f"- **Otimiza√ß√£o:** Melhorar {', '.join(areas_medianas)}\n"
    
    relatorio += "- **Manuten√ß√£o:** Manter as √°reas com scores altos\n"
    
    return relatorio, score_medio, status

# --- FUN√á√ïES ORIGINAIS MANTIDAS ---

def criar_prompt_validacao_preciso(texto, nome_arquivo, contexto_agente):
    """Cria um prompt de valida√ß√£o muito mais preciso para evitar falsos positivos"""
    
    prompt = f"""
{contexto_agente}

###BEGIN TEXTO PARA VALIDA√á√ÉO###
**Arquivo:** {nome_arquivo}
**Conte√∫do:**
{texto[:12000]}
###END TEXTO PARA VALIDA√á√ÉO###

## FORMATO DE RESPOSTA OBRIGAT√ìRIO:

### ‚úÖ CONFORMIDADE COM DIRETRIZES
- [Itens que est√£o alinhados com as diretrizes de branding]

**INCONSIST√äNCIAS COM BRANDING:**
- [S√≥ liste desvios REAIS das diretrizes de branding]

### üí° TEXTO REVISADO
- [Sugest√µes para aprimorar]

### üìä STATUS FINAL
**Documento:** [Aprovado/Necessita ajustes/Reprovado]
**Principais a√ß√µes necess√°rias:** [Lista resumida]
"""
    return prompt

def analisar_documento_por_slides(doc, contexto_agente):
    """Analisa documento slide por slide com alta precis√£o"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Analisando slide {i+1}..."):
            try:
                prompt_slide = f"""
{contexto_agente}

## AN√ÅLISE POR SLIDE - PRECIS√ÉO ABSOLUTA

###BEGIN TEXTO PARA VALIDA√á√ÉO###
**SLIDE {i+1}:**
{slide['conteudo'][:2000]}
###END TEXTO PARA VALIDA√á√ÉO###

**AN√ÅLISE DO SLIDE {i+1}:**

### ‚úÖ Pontos Fortes:
[O que est√° bom neste slide]

### ‚ö†Ô∏è Problemas REAIS:
- [Lista CURTA de problemas]

### üí° Sugest√µes Espec√≠ficas:
[Melhorias para ESTE slide espec√≠fico]

Considere que slides que s√£o introdut√≥rios ou apenas de t√≠tulos n√£o precisam de tanto rigor de branding

**STATUS:** [‚úîÔ∏è Aprovado / ‚ö†Ô∏è Ajustes Menores / ‚ùå Problemas S√©rios]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '‚ùå' in resposta.text or '‚ö†Ô∏è' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"‚ùå Erro na an√°lise do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relat√≥rio consolidado
    relatorio = f"# üìä RELAT√ìRIO DE VALIDA√á√ÉO - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Altera√ß√µes:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de aten√ß√£o
    slides_com_problemas = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_problemas:
        relatorio += "## üö® SLIDES QUE PRECISAM DE ATEN√á√ÉO:\n\n"
        for resultado in slides_com_problemas:
            relatorio += f"### üìã Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## üìà RESUMO EXECUTIVO\n\n"
    if slides_com_problemas:
        relatorio += f"**‚ö†Ô∏è {len(slides_com_problemas)} slide(s) necessitam de ajustes**\n"
        relatorio += f"**‚úÖ {len(doc['slides']) - len(slides_com_problemas)} slide(s) est√£o adequados**\n"
    else:
        relatorio += "**üéâ Todos os slides est√£o em conformidade com as diretrizes!**\n"
    
    return relatorio

def extract_text_from_pdf_com_slides(arquivo_pdf):
    """Extrai texto de PDF com informa√ß√£o de p√°ginas"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'p√°gina'
            })
        
        texto_completo = "\n\n".join([f"--- P√ÅGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extra√ß√£o PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    """Extrai texto de PPTX com informa√ß√£o de slides"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extra√ß√£o PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    """Extrai texto de arquivos TXT e DOCX"""
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo n√£o suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extra√ß√£o: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file using multiple methods for better coverage
    """
    text = ""

    # Method 1: Try with pdfplumber (better for some PDFs)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to PyPDF2 if pdfplumber didn't extract much text
    if len(text.strip()) < 100:  # If very little text was extracted
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            print(f"PyPDF2 also failed for {pdf_path}: {e}")

    return text

def criar_analisadores_imagem(contexto_agente, contexto_global):
    """Cria analisadores especializados para imagens"""
    
    analisadores = {
        'composicao_visual': {
            'nome': 'üé® Especialista em Composi√ß√£o Visual',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM COMPOSI√á√ÉO VISUAL

**Sua tarefa:** Analisar EXCLUSIVAMENTE a composi√ß√£o visual da imagem.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Balanceamento** - Distribui√ß√£o equilibrada dos elementos
2. **Hierarquia Visual** - Foco e pontos de aten√ß√£o
3. **Espa√ßamento** - Uso adequado do espa√ßo
4. **Propor√ß√µes** - Rela√ß√£o entre elementos visuais
5. **Harmonia** - Conjunto visual coeso

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üé® RELAT√ìRIO DE COMPOSI√á√ÉO VISUAL

### ‚úÖ PONTOS FORTES DA COMPOSI√á√ÉO
- [Elementos bem compostos]

### ‚ö†Ô∏è PROBLEMAS DE COMPOSI√á√ÉO
- [Issues de organiza√ß√£o visual]

### üìä SCORE COMPOSI√á√ÉO: [X/10]

### üí° SUGEST√ïES DE MELHORIA VISUAL
- [Recomenda√ß√µes para melhor composi√ß√£o]
"""
        },
        'cores_branding': {
            'nome': 'üåà Especialista em Cores e Branding',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM CORES E BRANDING

**Sua tarefa:** Analisar EXCLUSIVAMENTE cores e alinhamento com branding.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Paleta de Cores** - Cores utilizadas na imagem
2. **Contraste** - Legibilidade e visibilidade
3. **Consist√™ncia** - Coer√™ncia com identidade visual
4. **Psicologia das Cores** - Efeito emocional das cores
5. **Acessibilidade** - Visibilidade para diferentes usu√°rios

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üåà RELAT√ìRIO DE CORES E BRANDING

### ‚úÖ CORES ALINHADAS
- [Cores que seguem as diretrizes]

### ‚ùå PROBLEMAS DE COR
- [Cores fora do padr√£o]

### üìä SCORE CORES: [X/10]

### üéØ RECOMENDA√á√ïES DE COR
- [Sugest√µes para paleta de cores]
"""
        },
        'tipografia_texto': {
            'nome': 'üî§ Especialista em Tipografia e Texto',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM TIPOGRAFIA E TEXTO

**Sua tarefa:** Analisar EXCLUSIVAMENTE tipografia e elementos textuais.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Legibilidade** - Facilidade de leitura do texto
2. **Hierarquia Tipogr√°fica** - Tamanhos e pesos de fonte
3. **Alinhamento** - Organiza√ß√£o do texto na imagem
4. **Consist√™ncia** - Uso uniforme de fontes
5. **Mensagem Textual** - Conte√∫do das palavras

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üî§ RELAT√ìRIO DE TIPOGRAFIA

### ‚úÖ ACERTOS TIPOGR√ÅFICOS
- [Elementos textuais bem executados]

### ‚ö†Ô∏è PROBLEMAS DE TEXTO
- [Issues com tipografia e texto]

### üìä SCORE TIPOGRAFIA: [X/10]

### ‚úèÔ∏è SUGEST√ïES TIPOGR√ÅFICAS
- [Melhorias para texto e fontes]
"""
        },
        'elementos_marca': {
            'nome': 'üè∑Ô∏è Especialista em Elementos de Marca',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM ELEMENTOS DE MARCA

**Sua tarefa:** Analisar EXCLUSIVAMENTE elementos de identidade visual da marca.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Logo e Identidade** - Uso correto da marca
2. **Elementos Gr√°ficos** - √çcones, padr√µes, ilustra√ß√µes
3. **Fotografia** - Estilo e tratamento de imagens
4. **Consist√™ncia Visual** - Coer√™ncia com guidelines
5. **Diferencia√ß√£o** - Elementos √∫nicos da marca

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üè∑Ô∏è RELAT√ìRIO DE ELEMENTOS DE MARCA

### ‚úÖ ELEMENTOS CORRETOS
- [Elementos alinhados com a marca]

### ‚ùå ELEMENTOS INCORRETOS
- [Elementos fora do padr√£o]

### üìä SCORE MARCA: [X/10]

### üé® RECOMENDA√á√ïES DE MARCA
- [Sugest√µes para identidade visual]
"""
        },
        'impacto_comunicacao': {
            'nome': 'üéØ Especialista em Impacto e Comunica√ß√£o',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM IMPACTO E COMUNICA√á√ÉO

**Sua tarefa:** Analisar EXCLUSIVAMENTE impacto visual e comunica√ß√£o.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Mensagem Central** - Clareza da comunica√ß√£o
2. **Apelo Emocional** - Conex√£o com o p√∫blico
3. **Chamada para A√ß√£o** - Efetividade persuasiva
4. **Originalidade** - Diferencia√ß√£o criativa
5. **Memorabilidade** - Capacidade de ser lembrado

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üéØ RELAT√ìRIO DE IMPACTO

### ‚úÖ PONTOS DE IMPACTO
- [Elementos comunicativos eficazes]

### üìâ OPORTUNIDADES DE MELHORIA
- [√Åreas para aumentar impacto]

### üìä SCORE IMPACTO: [X/10]

### üöÄ ESTRAT√âGIAS DE COMUNICA√á√ÉO
- [T√©cnicas para melhor comunica√ß√£o]
"""
        }
    }
    
    return analisadores

def criar_analisadores_video(contexto_agente, contexto_global, contexto_video_especifico):
    """Cria analisadores especializados para v√≠deos"""
    
    analisadores = {
        'narrativa_estrutura': {
            'nome': 'üìñ Especialista em Narrativa e Estrutura',
            'prompt': f"""
{contexto_agente}
{contexto_global}
{contexto_video_especifico}

## FUN√á√ÉO: ESPECIALISTA EM NARRATIVA E ESTRUTURA

**Sua tarefa:** Analisar EXCLUSIVAMENTE a estrutura narrativa do v√≠deo.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Arco Narrativo** - Desenvolvimento da hist√≥ria
2. **Ritmo** - Velocidade e fluidez da narrativa
3. **Estrutura** - Organiza√ß√£o do conte√∫do
4. **Transi√ß√µes** - Conex√£o entre cenas/ideias
5. **Cl√≠max e Resolu√ß√£o** - Ponto alto e conclus√£o

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üìñ RELAT√ìRIO DE NARRATIVA

### ‚úÖ PONTOS FORTES DA NARRATIVA
- [Elementos narrativos bem executados]

### ‚ö†Ô∏è PROBLEMAS DE ESTRUTURA
- [Issues na organiza√ß√£o do conte√∫do]

### üìä SCORE NARRATIVA: [X/10]

### üí° SUGEST√ïES NARRATIVAS
- [Melhorias para estrutura e ritmo]
"""
        },
        'qualidade_audio': {
            'nome': 'üîä Especialista em Qualidade de √Åudio',
            'prompt': f"""
{contexto_agente}
{contexto_global}
{contexto_video_especifico}

## FUN√á√ÉO: ESPECIALISTA EM QUALIDADE DE √ÅUDIO

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos de √°udio do v√≠deo.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Clareza Vocal** - Inteligibilidade da fala
2. **Qualidade T√©cnica** - Ru√≠do, distor√ß√£o, equil√≠brio
3. **Trilha Sonora** - M√∫sica e efeitos sonoros
4. **Sincroniza√ß√£o** - Rela√ß√£o √°udio-v√≠deo
5. **Mixagem** - Balanceamento de elementos sonoros

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üîä RELAT√ìRIO DE √ÅUDIO

### ‚úÖ ACERTOS DE √ÅUDIO
- [Elementos sonoros bem executados]

### ‚ùå PROBLEMAS DE √ÅUDIO
- [Issues t√©cnicos e de qualidade]

### üìä SCORE √ÅUDIO: [X/10]

### üéß RECOMENDA√á√ïES DE √ÅUDIO
- [Sugest√µes para melhor qualidade sonora]
"""
        },
        'visual_cinematografia': {
            'nome': 'üé• Especialista em Visual e Cinematografia',
            'prompt': f"""
{contexto_agente}
{contexto_global}
{contexto_video_especifico}

## FUN√á√ÉO: ESPECIALISTA EM VISUAL

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos visuais do v√≠deo.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Enquadramento** - Composi√ß√£o de cenas
2. **Ilumina√ß√£o** - Uso da luz e sombras
3. **Cores e Grading** - Tratamento de cor
4. **Movimento de C√¢mera** - Din√¢mica visual
5. **Qualidade de Imagem** - Resolu√ß√£o e nitidez

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üé• RELAT√ìRIO VISUAL

### ‚úÖ PONTOS FORTES VISUAIS
- [Elementos visuais bem executados]

### ‚ö†Ô∏è PROBLEMAS VISUAIS
- [Issues de qualidade visual]

### üìä SCORE VISUAL: [X/10]

### üåü SUGEST√ïES VISUAIS
- [Melhorias para cinematografia]
"""
        },
        'branding_consistencia': {
            'nome': 'üè¢ Especialista em Branding e Consist√™ncia',
            'prompt': f"""
{contexto_agente}
{contexto_global}
{contexto_video_especifico}

## FUN√á√ÉO: ESPECIALISTA EM BRANDING E CONSIST√äNCIA

**Sua tarefa:** Analisar EXCLUSIVAMENTE alinhamento com branding.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Identidade Visual** - Cores, logos, elementos da marca
2. **Tom de Voz** - Personalidade da comunica√ß√£o
3. **Mensagem Central** - Alinhamento com valores
4. **Consist√™ncia** - Uniformidade ao longo do v√≠deo
5. **P√∫blico-Alvo** - Adequa√ß√£o ao destinat√°rio

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üè¢ RELAT√ìRIO DE BRANDING

### ‚úÖ ALINHAMENTOS DE MARCA
- [Elementos que seguem as diretrizes]

### ‚ùå DESVIOS DE MARCA
- [Elementos fora do padr√£o]

### üìä SCORE BRANDING: [X/10]

### üéØ RECOMENDA√á√ïES DE MARCA
- [Sugest√µes para melhor alinhamento]
"""
        }
        }
    
    
    return analisadores

def executar_analise_imagem_especializada(uploaded_image, nome_imagem, analisadores):
    """Executa an√°lise especializada para imagens com m√∫ltiplos especialistas"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN IMAGEM PARA AN√ÅLISE###
**Arquivo:** {nome_imagem}
**An√°lise solicitada para:** {config['nome']}
###END IMAGEM PARA AN√ÅLISE###

Por favor, forne√ßa sua an√°lise especializada no formato solicitado.
"""
                
                # Processar imagem com o especialista espec√≠fico
                response = modelo_vision.generate_content([
                    prompt_completo,
                    {"mime_type": "image/jpeg", "data": uploaded_image.getvalue()}
                ])
                
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': response.text,
                    'score': extrair_score(response.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"‚ùå Erro na an√°lise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def executar_analise_video_especializada(uploaded_video, nome_video, analisadores):
    """Executa an√°lise especializada para v√≠deos com m√∫ltiplos especialistas"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN V√çDEO PARA AN√ÅLISE###
**Arquivo:** {nome_video}
**An√°lise solicitada para:** {config['nome']}
###END V√çDEO PARA AN√ÅLISE###

Por favor, forne√ßa sua an√°lise especializada no formato solicitado.
"""
                
                # Processar v√≠deo com o especialista espec√≠fico
                video_bytes = uploaded_video.getvalue()
                
                if len(video_bytes) < 200 * 1024 * 1024:
                    response = modelo_vision.generate_content([
                        prompt_completo,
                        {"mime_type": uploaded_video.type, "data": video_bytes}
                    ])
                else:
                    response = modelo_vision.generate_content([
                        prompt_completo,
                        {"mime_type": uploaded_video.type, "data": video_bytes}
                    ])
                
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': response.text,
                    'score': extrair_score(response.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"‚ùå Erro na an√°lise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_imagem_consolidado(resultados_especialistas, nome_imagem, dimensoes):
    """Gera relat√≥rio consolidado para imagens"""
    
    # Calcular score m√©dio
    scores = [resultado['score'] for resultado in resultados_especialistas.values() if resultado['score'] > 0]
    score_medio = sum(scores) / len(scores) if scores else 0
    
    # Determinar status geral
    if score_medio >= 8:
        status = "‚úÖ APROVADO"
        cor_status = "green"
    elif score_medio >= 6:
        status = "‚ö†Ô∏è AJUSTES MENORES"
        cor_status = "orange"
    else:
        status = "‚ùå REPROVADO"
        cor_status = "red"
    
    relatorio = f"""
# üñºÔ∏è RELAT√ìRIO CONSOLIDADO DE IMAGEM

**Arquivo:** {nome_imagem}
**Dimens√µes:** {dimensoes}
**Status Geral:** <span style='color:{cor_status}'>{status}</span>
**Score M√©dio:** {score_medio:.1f}/10
**Data da An√°lise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## üéñÔ∏è SCORES POR √ÅREA ESPECIALIZADA
"""
    
    # Adicionar scores individuais
    for area, resultado in resultados_especialistas.items():
        emoji = "‚úÖ" if resultado['score'] >= 8 else "‚ö†Ô∏è" if resultado['score'] >= 6 else "‚ùå"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## üìã AN√ÅLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar an√°lises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## üöÄ RESUMO EXECUTIVO - IMAGEM

**Status Final:** {status}
**Score Geral:** {score_medio:.1f}/10

### üéØ PR√ìXIMOS PASSOS RECOMENDADOS:
"""
    
    # Recomenda√ß√µes baseadas nos scores
    areas_baixas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] < 6]
    if areas_baixas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_baixas]
        relatorio += f"- **Prioridade M√°xima:** Focar em {', '.join(nomes_areas)}\n"
    
    areas_medianas = [area for area, resultado in resultados_especialistas.items() if 6 <= resultado['score'] < 8]
    if areas_medianas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_medianas]
        relatorio += f"- **Otimiza√ß√£o Necess√°ria:** Melhorar {', '.join(nomes_areas)}\n"
    
    areas_altas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] >= 8]
    if areas_altas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_altas]
        relatorio += f"- **Manuten√ß√£o:** Manter a excel√™ncia em {', '.join(nomes_areas)}\n"
    
    return relatorio, score_medio, status

def gerar_relatorio_video_consolidado(resultados_especialistas, nome_video, tipo_video):
    """Gera relat√≥rio consolidado para v√≠deos"""
    
    # Calcular score m√©dio
    scores = [resultado['score'] for resultado in resultados_especialistas.values() if resultado['score'] > 0]
    score_medio = sum(scores) / len(scores) if scores else 0
    
    # Determinar status geral
    if score_medio >= 8:
        status = "‚úÖ APROVADO"
        cor_status = "green"
    elif score_medio >= 6:
        status = "‚ö†Ô∏è AJUSTES MENORES"
        cor_status = "orange"
    else:
        status = "‚ùå REPROVADO"
        cor_status = "red"
    
    relatorio = f"""
# üé¨ RELAT√ìRIO CONSOLIDADO DE V√çDEO

**Arquivo:** {nome_video}
**Formato:** {tipo_video}
**Status Geral:** <span style='color:{cor_status}'>{status}</span>
**Score M√©dio:** {score_medio:.1f}/10
**Data da An√°lise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## üéñÔ∏è SCORES POR √ÅREA ESPECIALIZADA
"""
    
    # Adicionar scores individuais
    for area, resultado in resultados_especialistas.items():
        emoji = "‚úÖ" if resultado['score'] >= 8 else "‚ö†Ô∏è" if resultado['score'] >= 6 else "‚ùå"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## üìã AN√ÅLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar an√°lises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## üöÄ RESUMO EXECUTIVO - V√çDEO

**Status Final:** {status}
**Score Geral:** {score_medio:.1f}/10

### üéØ PR√ìXIMOS PASSOS RECOMENDADOS:
"""
    
    # Recomenda√ß√µes baseadas nos scores
    areas_baixas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] < 6]
    if areas_baixas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_baixas]
        relatorio += f"- **Prioridade M√°xima:** Focar em {', '.join(nomes_areas)}\n"
    
    areas_medianas = [area for area, resultado in resultados_especialistas.items() if 6 <= resultado['score'] < 8]
    if areas_medianas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_medianas]
        relatorio += f"- **Otimiza√ß√£o Necess√°ria:** Melhorar {', '.join(nomes_areas)}\n"
    
    areas_altas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] >= 8]
    if areas_altas:
        nomes_areas = [resultados_especialistas[area]['nome'] for area in areas_altas]
        relatorio += f"- **Manuten√ß√£o:** Manter a excel√™ncia em {', '.join(nomes_areas)}\n"
    
    return relatorio, score_medio, status

# --- FUN√á√ïES DE AN√ÅLISE DE TEXTO (MANTIDAS) ---

def criar_analisadores_texto(contexto_agente, contexto_global):
    """Cria prompts especializados para cada √°rea de an√°lise de texto"""
    
    analisadores = {
        'ortografia': {
            'nome': 'üî§ Especialista em Ortografia e Gram√°tica',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM ORTOGRAFIA E GRAM√ÅTICA PORTUGU√äS BR

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos ortogr√°ficos e gramaticais.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Ortografia** - Erros de escrita
2. **Gram√°tica** - Concord√¢ncia, reg√™ncia, coloca√ß√£o
3. **Pontua√ß√£o** - Uso de v√≠rgulas, pontos, etc.
4. **Acentua√ß√£o** - Erros de acentua√ß√£o
5. **Padr√£o Culto** - Conformidade com norma culta

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üî§ RELAT√ìRIO ORTOGR√ÅFICO

### ‚úÖ ACERTOS
- [Itens corretos]

### ‚ùå ERROS IDENTIFICADOS
- [Lista espec√≠fica de erros com corre√ß√µes]

### üìä SCORE ORTOGR√ÅFICO: [X/10]

### üí° SUGEST√ïES DE MELHORIA
- [Recomenda√ß√µes espec√≠ficas]
"""
        },
        'lexico': {
            'nome': 'üìö Especialista em L√©xico e Vocabul√°rio',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM L√âXICO E VOCABUL√ÅRIO

**Sua tarefa:** Analisar EXCLUSIVAMENTE aspectos lexicais e de vocabul√°rio.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Variedade Lexical** - Riqueza de vocabul√°rio
2. **Precis√£o Sem√¢ntica** - Uso adequado das palavras
3. **Repeti√ß√£o** - Palavras ou express√µes repetidas em excesso
4. **Jarg√µes** - Uso inadequado de termos t√©cnicos
5. **Clareza** - Facilidade de compreens√£o

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üìö RELAT√ìRIO LEXICAL

### ‚úÖ VOCABUL√ÅRIO ADEQUADO
- [Pontos fortes do vocabul√°rio]

### ‚ö†Ô∏è ASPECTOS A MELHORAR
- [Problemas lexicais identificados]

### üîÑ SUGEST√ïES DE SIN√îNIMOS
- [Palavras para substituir]

### üìä SCORE LEXICAL: [X/10]
"""
        },
        'branding': {
            'nome': 'üé® Especialista em Branding e Identidade',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM BRANDING E IDENTIDADE

**Sua tarefa:** Analisar EXCLUSIVAMENTE conformidade com diretrizes de branding.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Tom de Voz** - Alinhamento com personalidade da marca
2. **Mensagem Central** - Consist√™ncia da mensagem
3. **Valores da Marca** - Reflexo dos valores organizacionais
4. **P√∫blico-Alvo** - Adequa√ß√£o ao p√∫blico pretendido
5. **Diferencia√ß√£o** - Elementos √∫nicos da marca

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üé® RELAT√ìRIO DE BRANDING

### ‚úÖ ALINHAMENTOS
- [Elementos que seguem as diretrizes]

### ‚ùå DESVIOS IDENTIFICADOS
- [Elementos fora do padr√£o da marca]

### üìä SCORE BRANDING: [X/10]

### üí° RECOMENDA√á√ïES ESTRAT√âGICAS
- [Sugest√µes para melhor alinhamento]
"""
        },
        'estrutura': {
            'nome': 'üìã Especialista em Estrutura e Formata√ß√£o',
            'prompt': f"""
{contexto_agente}
{contexto_global}

## FUN√á√ÉO: ESPECIALISTA EM ESTRUTURA E FORMATA√á√ÉO

**Sua tarefa:** Analisar EXCLUSIVAMENTE estrutura e organiza√ß√£o do conte√∫do.

### CRIT√âRIOS DE AN√ÅLISE:
1. **Organiza√ß√£o** - Estrutura l√≥gica e sequ√™ncia
2. **Hierarquia** - Uso adequado de t√≠tulos e subt√≠tulos
3. **Coes√£o** - Liga√ß√£o entre ideias e par√°grafos
4. **Formata√ß√£o** - Consist√™ncia visual
5. **Objetividade** - Clareza na apresenta√ß√£o das ideias

### FORMATO DE RESPOSTA OBRIGAT√ìRIO:

## üìã RELAT√ìRIO ESTRUTURAL

### ‚úÖ ESTRUTURA ADEQUADA
- [Elementos bem organizados]

### ‚ö†Ô∏è PROBLEMAS ESTRUTURAIS
- [Issues de organiza√ß√£o identificados]

### üìä SCORE ESTRUTURAL: [X/10]

### üèóÔ∏è SUGEST√ïES DE REORGANIZA√á√ÉO
- [Melhorias na estrutura]
"""
        }
        
    }
    
    return analisadores

def executar_analise_texto_especializada(texto, nome_arquivo, analisadores):
    """Executa an√°lise com m√∫ltiplos especialistas para texto"""
    
    resultados = {}
    
    for area, config in analisadores.items():
        with st.spinner(f"Executando {config['nome']}..."):
            try:
                prompt_completo = f"""
{config['prompt']}

###BEGIN TEXTO PARA AN√ÅLISE###
**Arquivo:** {nome_arquivo}
**Conte√∫do:**
{texto[:8000]}
###END TEXTO PARA AN√ÅLISE###

Por favor, forne√ßa sua an√°lise no formato solicitado.
"""
                
                resposta = modelo_texto.generate_content(prompt_completo)
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': resposta.text,
                    'score': extrair_score(resposta.text)
                }
                
            except Exception as e:
                resultados[area] = {
                    'nome': config['nome'],
                    'analise': f"‚ùå Erro na an√°lise: {str(e)}",
                    'score': 0
                }
    
    return resultados

def gerar_relatorio_texto_consolidado(resultados_especialistas, nome_arquivo):
    """Gera relat√≥rio consolidado a partir das an√°lises especializadas de texto"""
    
    # Calcular score m√©dio
    scores = [resultado['score'] for resultado in resultados_especialistas.values() if resultado['score'] > 0]
    score_medio = sum(scores) / len(scores) if scores else 0
    
    # Determinar status geral
    if score_medio >= 8:
        status = "‚úÖ APROVADO"
        cor_status = "green"
    elif score_medio >= 6:
        status = "‚ö†Ô∏è AJUSTES MENORES"
        cor_status = "orange"
    else:
        status = "‚ùå REPROVADO"
        cor_status = "red"
    
    relatorio = f"""
# üìä RELAT√ìRIO CONSOLIDADO DE VALIDA√á√ÉO

**Documento:** {nome_arquivo}
**Status Geral:** <span style='color:{cor_status}'>{status}</span>
**Score M√©dio:** {score_medio:.1f}/10
**Data da An√°lise:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}

## üéñÔ∏è SCORES POR √ÅREA
"""
    
    # Adicionar scores individuais
    for area, resultado in resultados_especialistas.items():
        emoji = "‚úÖ" if resultado['score'] >= 8 else "‚ö†Ô∏è" if resultado['score'] >= 6 else "‚ùå"
        relatorio += f"- {emoji} **{resultado['nome']}:** {resultado['score']}/10\n"
    
    relatorio += "\n## üìã AN√ÅLISES DETALHADAS POR ESPECIALISTA\n"
    
    # Adicionar an√°lises detalhadas
    for area, resultado in resultados_especialistas.items():
        relatorio += f"\n### {resultado['nome']}\n"
        relatorio += f"{resultado['analise']}\n"
        relatorio += "---\n"
    
    # Resumo executivo
    relatorio += f"""
## üöÄ RESUMO EXECUTIVO

**Status Final:** {status}
**Score Geral:** {score_medio:.1f}/10

### üéØ PR√ìXIMOS PASSOS RECOMENDADOS:
"""
    
    # Recomenda√ß√µes baseadas nos scores
    areas_baixas = [area for area, resultado in resultados_especialistas.items() if resultado['score'] < 6]
    if areas_baixas:
        relatorio += f"- **Prioridade:** Focar em {', '.join(areas_baixas)}\n"
    
    areas_medianas = [area for area, resultado in resultados_especialistas.items() if 6 <= resultado['score'] < 8]
    if areas_medianas:
        relatorio += f"- **Otimiza√ß√£o:** Melhorar {', '.join(areas_medianas)}\n"
    
    relatorio += "- **Manuten√ß√£o:** Manter as √°reas com scores altos\n"
    
    return relatorio, score_medio, status

def extrair_score(texto_analise):
    """Extrai score num√©rico do texto de an√°lise"""
    import re
    padrao = r'SCORE.*?\[(\d+)(?:/10)?\]'
    correspondencias = re.findall(padrao, texto_analise, re.IGNORECASE)
    if correspondencias:
        return int(correspondencias[0])
    return 5  # Score padr√£o se n√£o encontrar

# --- FUN√á√ïES ORIGINAIS MANTIDAS ---

def criar_prompt_validacao_preciso(texto, nome_arquivo, contexto_agente):
    """Cria um prompt de valida√ß√£o muito mais preciso para evitar falsos positivos"""
    
    prompt = f"""
{contexto_agente}

###BEGIN TEXTO PARA VALIDA√á√ÉO###
**Arquivo:** {nome_arquivo}
**Conte√∫do:**
{texto[:12000]}
###END TEXTO PARA VALIDA√á√ÉO###

## FORMATO DE RESPOSTA OBRIGAT√ìRIO:

### ‚úÖ CONFORMIDADE COM DIRETRIZES
- [Itens que est√£o alinhados com as diretrizes de branding]

**INCONSIST√äNCIAS COM BRANDING:**
- [S√≥ liste desvios REAIS das diretrizes de branding]

### üí° TEXTO REVISADO
- [Sugest√µes para aprimorar]

### üìä STATUS FINAL
**Documento:** [Aprovado/Necessita ajustes/Reprovado]
**Principais a√ß√µes necess√°rias:** [Lista resumida]
"""
    return prompt

def analisar_documento_por_slides(doc, contexto_agente):
    """Analisa documento slide por slide com alta precis√£o"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Analisando slide {i+1}..."):
            try:
                prompt_slide = f"""
{contexto_agente}

## AN√ÅLISE POR SLIDE - PRECIS√ÉO ABSOLUTA

###BEGIN TEXTO PARA VALIDA√á√ÉO###
**SLIDE {i+1}:**
{slide['conteudo'][:2000]}
###END TEXTO PARA VALIDA√á√ÉO###

**AN√ÅLISE DO SLIDE {i+1}:**

### ‚úÖ Pontos Fortes:
[O que est√° bom neste slide]

### ‚ö†Ô∏è Problemas REAIS:
- [Lista CURTA de problemas]

### üí° Sugest√µes Espec√≠ficas:
[Melhorias para ESTE slide espec√≠fico]

Considere que slides que s√£o introdut√≥rios ou apenas de t√≠tulos n√£o precisam de tanto rigor de branding

**STATUS:** [‚úîÔ∏è Aprovado / ‚ö†Ô∏è Ajustes Menores / ‚ùå Problemas S√©rios]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '‚ùå' in resposta.text or '‚ö†Ô∏è' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"‚ùå Erro na an√°lise do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relat√≥rio consolidado
    relatorio = f"# üìä RELAT√ìRIO DE VALIDA√á√ÉO - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Altera√ß√µes:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de aten√ß√£o
    slides_com_problemas = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_problemas:
        relatorio += "## üö® SLIDES QUE PRECISAM DE ATEN√á√ÉO:\n\n"
        for resultado in slides_com_problemas:
            relatorio += f"### üìã Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## üìà RESUMO EXECUTIVO\n\n"
    if slides_com_problemas:
        relatorio += f"**‚ö†Ô∏è {len(slides_com_problemas)} slide(s) necessitam de ajustes**\n"
        relatorio += f"**‚úÖ {len(doc['slides']) - len(slides_com_problemas)} slide(s) est√£o adequados**\n"
    else:
        relatorio += "**üéâ Todos os slides est√£o em conformidade com as diretrizes!**\n"
    
    return relatorio

def extract_text_from_pdf_com_slides(arquivo_pdf):
    """Extrai texto de PDF com informa√ß√£o de p√°ginas"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo_pdf)
        slides_info = []
        
        for pagina_num, pagina in enumerate(pdf_reader.pages):
            texto = pagina.extract_text()
            slides_info.append({
                'numero': pagina_num + 1,
                'conteudo': texto,
                'tipo': 'p√°gina'
            })
        
        texto_completo = "\n\n".join([f"--- P√ÅGINA {s['numero']} ---\n{s['conteudo']}" for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extra√ß√£o PDF: {str(e)}", []

def extract_text_from_pptx_com_slides(arquivo_pptx):
    """Extrai texto de PPTX com informa√ß√£o de slides"""
    try:
        from pptx import Presentation
        import io
        
        prs = Presentation(io.BytesIO(arquivo_pptx.read()))
        slides_info = []
        
        for slide_num, slide in enumerate(prs.slides):
            texto_slide = f"--- SLIDE {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto_slide += shape.text + "\n"
            
            slides_info.append({
                'numero': slide_num + 1,
                'conteudo': texto_slide,
                'tipo': 'slide'
            })
        
        texto_completo = "\n\n".join([s['conteudo'] for s in slides_info])
        return texto_completo, slides_info
        
    except Exception as e:
        return f"Erro na extra√ß√£o PPTX: {str(e)}", []

def extrair_texto_arquivo(arquivo):
    """Extrai texto de arquivos TXT e DOCX"""
    try:
        if arquivo.type == "text/plain":
            return str(arquivo.read(), "utf-8")
        elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        else:
            return f"Tipo n√£o suportado: {arquivo.type}"
    except Exception as e:
        return f"Erro na extra√ß√£o: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """
    Extract text from a PDF file using multiple methods for better coverage
    """
    text = ""

    # Method 1: Try with pdfplumber (better for some PDFs)
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
    except Exception as e:
        print(f"pdfplumber failed for {pdf_path}: {e}")

    # Method 2: Fallback to PyPDF2 if pdfplumber didn't extract much text
    if len(text.strip()) < 100:  # If very little text was extracted
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text 
        except Exception as e:
            print(f"PyPDF2 also failed for {pdf_path}: {e}")

    return text

# --- INICIALIZA√á√ÉO DE SESSION_STATE ---
if 'analise_especializada_texto' not in st.session_state:
    st.session_state.analise_especializada_texto = True

if 'analise_especializada_imagem' not in st.session_state:
    st.session_state.analise_especializada_imagem = True

if 'analise_especializada_video' not in st.session_state:
    st.session_state.analise_especializada_video = True

if 'analisadores_selecionados_texto' not in st.session_state:
    st.session_state.analisadores_selecionados_texto = ['ortografia', 'lexico', 'branding']

if 'analisadores_selecionados_imagem' not in st.session_state:
    st.session_state.analisadores_selecionados_imagem = ['composicao_visual', 'cores_branding', 'tipografia_texto', 'elementos_marca']

if 'analisadores_selecionados_video' not in st.session_state:
    st.session_state.analisadores_selecionados_video = ['narrativa_estrutura', 'qualidade_audio', 'visual_cinematografia', 'branding_consistencia']

if 'analise_detalhada' not in st.session_state:
    st.session_state.analise_detalhada = True

if 'validacao_triggered' not in st.session_state:
    st.session_state.validacao_triggered = False

if 'todos_textos' not in st.session_state:
    st.session_state.todos_textos = []

if 'resultados_analise_imagem' not in st.session_state:
    st.session_state.resultados_analise_imagem = []

if 'resultados_analise_video' not in st.session_state:
    st.session_state.resultados_analise_video = []

# --- ABA: VALIDA√á√ÉO UNIFICADA ---
with tab_mapping["‚úÖ Valida√ß√£o Unificada"]:
    st.header("‚úÖ Valida√ß√£o Unificada de Conte√∫do")
    
    if not st.session_state.get('agente_selecionado'):
        st.info("Selecione um agente primeiro na aba de Chat")
    else:
        agente = st.session_state.agente_selecionado
        st.subheader(f"Valida√ß√£o com: {agente.get('nome', 'Agente')}")
        
        # Container de contexto global
        st.markdown("---")
        st.subheader("üéØ Contexto para An√°lise")
        
        contexto_global = st.text_area(
            "**‚úçÔ∏è Contexto adicional para todas as an√°lises:**", 
            height=120, 
            key="contexto_global_validacao",
            placeholder="Forne√ßa contexto adicional que ser√° aplicado a TODAS as an√°lises (texto, documentos, imagens e v√≠deos)..."
        )
        
        # Subabas para diferentes tipos de valida√ß√£o
        subtab_imagem, subtab_texto, subtab_video = st.tabs(["üñºÔ∏è Valida√ß√£o de Imagem", "üìÑ Valida√ß√£o de Documentos", "üé¨ Valida√ß√£o de V√≠deo"])
        
        with subtab_texto:
            st.subheader("üìÑ Valida√ß√£o de Documentos e Texto")
            
            # Bot√£o para limpar an√°lises de texto
            if st.button("üóëÔ∏è Limpar An√°lises de Texto", key="limpar_analises_texto"):
                st.session_state.validacao_triggered = False
                st.session_state.todos_textos = []
                st.rerun()
            
            # Container principal com duas colunas
            col_entrada, col_saida = st.columns([1, 1])
            
            with col_entrada:
                st.markdown("### üì• Entrada de Conte√∫do")
                
                # Op√ß√£o 1: Texto direto
                texto_input = st.text_area(
                    "**‚úçÔ∏è Digite o texto para valida√ß√£o:**", 
                    height=150, 
                    key="texto_validacao",
                    placeholder="Cole aqui o texto que deseja validar..."
                )
                
                # Op√ß√£o 2: Upload de m√∫ltiplos arquivos
                st.markdown("### üìé Ou carregue arquivos")
                
                arquivos_documentos = st.file_uploader(
                    "**Documentos suportados:** PDF, PPTX, TXT, DOCX",
                    type=['pdf', 'pptx', 'txt', 'docx'],
                    accept_multiple_files=True,
                    key="arquivos_documentos_validacao"
                )
                
                # Configura√ß√µes de an√°lise
                with st.expander("‚öôÔ∏è Configura√ß√µes de An√°lise de Texto"):
                    analise_especializada = st.checkbox(
                        "An√°lise especializada por √°reas (recomendado)",
                        value=st.session_state.analise_especializada_texto,
                        help="Usa m√∫ltiplos especialistas para an√°lise mais precisa"
                    )
                    
                    analisadores_selecionados = st.multiselect(
                        "Especialistas de texto a incluir:",
                        options=['ortografia', 'lexico', 'branding', 'estrutura', 'engajamento'],
                        default=st.session_state.analisadores_selecionados_texto,
                        format_func=lambda x: {
                            'ortografia': 'üî§ Ortografia e Gram√°tica',
                            'lexico': 'üìö L√©xico e Vocabul√°rio', 
                            'branding': 'üé® Branding e Identidade',
                            'estrutura': 'üìã Estrutura e Formata√ß√£o',
                            'engajamento': 'üéØ Engajamento e Persuas√£o'
                        }[x]
                    )
                    
                    analise_detalhada = st.checkbox(
                        "An√°lise detalhada por slide/p√°gina",
                        value=st.session_state.analise_detalhada
                    )
                
                # Bot√£o de valida√ß√£o
                if st.button("‚úÖ Validar Conte√∫do de Texto", type="primary", key="validate_documents", use_container_width=True):
                    st.session_state.validacao_triggered = True
                    st.session_state.analise_especializada_texto = analise_especializada
                    st.session_state.analise_detalhada = analise_detalhada
                    st.session_state.analisadores_selecionados_texto = analisadores_selecionados
            
            with col_saida:
                st.markdown("### üìä Resultados de Texto")
                
                if st.session_state.validacao_triggered:
                    # Processar todos os conte√∫dos
                    todos_textos = []
                    arquivos_processados = []
                    
                    # Adicionar texto manual se existir
                    if texto_input and texto_input.strip():
                        todos_textos.append({
                            'nome': 'Texto_Manual',
                            'conteudo': texto_input,
                            'tipo': 'texto_direto',
                            'tamanho': len(texto_input),
                            'slides': []
                        })
                    
                    # Processar arquivos uploadados
                    if arquivos_documentos:
                        for arquivo in arquivos_documentos:
                            with st.spinner(f"Processando {arquivo.name}..."):
                                try:
                                    if arquivo.type == "application/pdf":
                                        texto_extraido, slides_info = extract_text_from_pdf_com_slides(arquivo)
                                    elif arquivo.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                        texto_extraido, slides_info = extract_text_from_pptx_com_slides(arquivo)
                                    elif arquivo.type in ["text/plain", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
                                        texto_extraido = extrair_texto_arquivo(arquivo)
                                        slides_info = []
                                    else:
                                        st.warning(f"Tipo de arquivo n√£o suportado: {arquivo.name}")
                                        continue
                                    
                                    if texto_extraido and texto_extraido.strip():
                                        todos_textos.append({
                                            'nome': arquivo.name,
                                            'conteudo': texto_extraido,
                                            'slides': slides_info,
                                            'tipo': arquivo.type,
                                            'tamanho': len(texto_extraido)
                                        })
                                        arquivos_processados.append(arquivo.name)
                                    
                                except Exception as e:
                                    st.error(f"‚ùå Erro ao processar {arquivo.name}: {str(e)}")
                    
                    # Verificar se h√° conte√∫do para validar
                    if not todos_textos:
                        st.warning("‚ö†Ô∏è Nenhum conte√∫do v√°lido encontrado para valida√ß√£o.")
                    else:
                        st.success(f"‚úÖ {len(todos_textos)} documento(s) processado(s) com sucesso!")
                        
                        # Exibir estat√≠sticas r√°pidas
                        col_docs, col_palavras, col_chars = st.columns(3)
                        with col_docs:
                            st.metric("üìÑ Documentos", len(todos_textos))
                        with col_palavras:
                            total_palavras = sum(len(doc['conteudo'].split()) for doc in todos_textos)
                            st.metric("üìù Palavras", total_palavras)
                        with col_chars:
                            total_chars = sum(doc['tamanho'] for doc in todos_textos)
                            st.metric("üî§ Caracteres", f"{total_chars:,}")
                        
                        # An√°lise individual por documento
                        st.markdown("---")
                        st.subheader("üìã An√°lise Individual por Documento")
                        
                        for doc in todos_textos:
                            with st.expander(f"üìÑ {doc['nome']} - {doc['tamanho']} chars", expanded=True):
                                # Informa√ß√µes b√°sicas do documento
                                col_info1, col_info2 = st.columns(2)
                                with col_info1:
                                    st.write(f"**Tipo:** {doc['tipo']}")
                                    st.write(f"**Tamanho:** {doc['tamanho']} caracteres")
                                with col_info2:
                                    if doc['slides']:
                                        st.write(f"**Slides/P√°ginas:** {len(doc['slides'])}")
                                    else:
                                        st.write("**Estrutura:** Texto simples")
                                
                                # Contexto aplicado
                                if contexto_global and contexto_global.strip():
                                    st.info(f"**üéØ Contexto Aplicado:** {contexto_global}")
                                
                                # An√°lise de branding
                                with st.spinner(f"Analisando {doc['nome']}..."):
                                    try:
                                        # Construir contexto do agente
                                        contexto_agente = ""
                                        if "base_conhecimento" in agente:
                                            contexto_agente = f"""
                                            ###BEGIN DIRETRIZES DE BRANDING DO AGENTE:###
                                            {agente['base_conhecimento']}
                                            ###END DIRETRIZES DE BRANDING DO AGENTE###
                                            """
                                        
                                        # Adicionar contexto global se fornecido
                                        contexto_completo = contexto_agente
                                        if contexto_global and contexto_global.strip():
                                            contexto_completo += f"""
                                            ###BEGIN CONTEXTO ADICIONAL DO USUARIO###
                                            {contexto_global}
                                            ###END CONTEXTO ADICIONAL DO USUARIO###
                                            """
                                        
                                        # Escolher m√©todo de an√°lise
                                        if st.session_state.analise_especializada_texto:
                                            # AN√ÅLISE ESPECIALIZADA POR M√öLTIPLOS ESPECIALISTAS
                                            st.info("üéØ **Executando an√°lise especializada por m√∫ltiplos especialistas...**")
                                            
                                            # Criar analisadores especialistas
                                            analisadores_config = criar_analisadores_texto(contexto_completo, "")
                                            
                                            # Filtrar apenas os selecionados
                                            analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                     if k in st.session_state.analisadores_selecionados_texto}
                                            
                                            # Executar an√°lises especializadas
                                            resultados_especialistas = executar_analise_texto_especializada(
                                                doc['conteudo'], 
                                                doc['nome'], 
                                                analisadores_filtrados
                                            )
                                            
                                            # Gerar relat√≥rio consolidado
                                            relatorio_consolidado, score_medio, status = gerar_relatorio_texto_consolidado(
                                                resultados_especialistas, 
                                                doc['nome']
                                            )
                                            
                                            st.markdown(relatorio_consolidado, unsafe_allow_html=True)
                                            
                                        elif st.session_state.analise_detalhada and doc['slides']:
                                            # An√°lise detalhada por slide (m√©todo antigo)
                                            resultado_analise = analisar_documento_por_slides(doc, contexto_completo)
                                            st.markdown(resultado_analise)
                                        else:
                                            # An√°lise geral do documento (m√©todo antigo)
                                            prompt_analise = criar_prompt_validacao_preciso(doc['conteudo'], doc['nome'], contexto_completo)
                                            resposta = modelo_texto.generate_content(prompt_analise)
                                            st.markdown(resposta.text)
                                        
                                    except Exception as e:
                                        st.error(f"‚ùå Erro na an√°lise de {doc['nome']}: {str(e)}")
                        
                        # Relat√≥rio consolidado
                        st.markdown("---")
                        st.subheader("üìë Relat√≥rio Consolidado de Texto")
                        
                        # Bot√£o para exportar
                        if st.button("üì• Exportar Relat√≥rio Completo de Texto", key="exportar_relatorio_completo"):
                            relatorio = f"""
                            # RELAT√ìRIO DE VALIDA√á√ÉO DE CONTE√öDO DE TEXTO
                            
                            **Agente:** {agente.get('nome', 'N/A')}
                            **Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
                            **Total de Documentos:** {len(todos_textos)}
                            **Contexto Aplicado:** {contexto_global if contexto_global else 'Nenhum contexto adicional'}
                            **M√©todo de An√°lise:** {'Especializada por M√∫ltiplos Especialistas' if st.session_state.analise_especializada_texto else 'Tradicional'}
                            
                            ## DOCUMENTOS ANALISADOS:
                            {chr(10).join([f"{idx+1}. {doc['nome']} ({doc['tipo']}) - {doc['tamanho']} caracteres" for idx, doc in enumerate(todos_textos)])}
                            """
                            
                            st.download_button(
                                "üíæ Baixar Relat√≥rio em TXT",
                                data=relatorio,
                                file_name=f"relatorio_validacao_texto_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain"
                            )
                        
                        # Armazenar na sess√£o
                        st.session_state.todos_textos = todos_textos
                
                else:
                    st.info("Digite texto ou carregue arquivos para validar")

        with subtab_imagem:
            st.subheader("üñºÔ∏è Valida√ß√£o de Imagem")
            
            # Bot√£o para limpar an√°lises de imagem
            if st.button("üóëÔ∏è Limpar An√°lises de Imagem", key="limpar_analises_imagem"):
                st.session_state.resultados_analise_imagem = []
                st.rerun()
            
            uploaded_images = st.file_uploader(
                "Carregue uma ou mais imagens para an√°lise", 
                type=["jpg", "jpeg", "png", "webp"], 
                key="image_upload_validacao",
                accept_multiple_files=True
            )
            
            # Configura√ß√µes de an√°lise de imagem
            with st.expander("‚öôÔ∏è Configura√ß√µes de An√°lise de Imagem"):
                analise_especializada_imagem = st.checkbox(
                    "An√°lise especializada por √°reas (recomendado)",
                    value=st.session_state.analise_especializada_imagem,
                    help="Usa m√∫ltiplos especialistas visuais para an√°lise mais precisa",
                    key="analise_especializada_imagem_check"
                )
                
                analisadores_selecionados_imagem = st.multiselect(
                    "Especialistas de imagem a incluir:",
                    options=['composicao_visual', 'cores_branding', 'tipografia_texto', 'elementos_marca', 'impacto_comunicacao'],
                    default=st.session_state.analisadores_selecionados_imagem,
                    format_func=lambda x: {
                        'composicao_visual': 'üé® Composi√ß√£o Visual',
                        'cores_branding': 'üåà Cores e Branding', 
                        'tipografia_texto': 'üî§ Tipografia e Texto',
                        'elementos_marca': 'üè∑Ô∏è Elementos de Marca',
                        'impacto_comunicacao': 'üéØ Impacto e Comunica√ß√£o'
                    }[x],
                    key="analisadores_imagem_select"
                )
            
            if uploaded_images:
                st.success(f"‚úÖ {len(uploaded_images)} imagem(ns) carregada(s)")
                
                # Bot√£o para validar todas as imagens
                if st.button("üîç Validar Todas as Imagens", type="primary", key="validar_imagens_multiplas"):
                    
                    # Lista para armazenar resultados
                    resultados_analise = []
                    
                    # Loop atrav√©s de cada imagem
                    for idx, uploaded_image in enumerate(uploaded_images):
                        with st.spinner(f'Analisando imagem {idx+1} de {len(uploaded_images)}: {uploaded_image.name}...'):
                            try:
                                # Criar container para cada imagem
                                with st.container():
                                    st.markdown("---")
                                    col_img, col_info = st.columns([2, 1])
                                    
                                    with col_img:
                                        # Exibir imagem
                                        image = Image.open(uploaded_image)
                                        st.image(image, use_container_width=True, caption=f"Imagem {idx+1}: {uploaded_image.name}")
                                    
                                    with col_info:
                                        # Informa√ß√µes da imagem
                                        st.metric("üìê Dimens√µes", f"{image.width} x {image.height}")
                                        st.metric("üìä Formato", uploaded_image.type)
                                        st.metric("üìÅ Tamanho", f"{uploaded_image.size / 1024:.1f} KB")
                                    
                                    # Contexto aplicado
                                    if contexto_global and contexto_global.strip():
                                        st.info(f"**üéØ Contexto Aplicado:** {contexto_global}")
                                    
                                    # An√°lise individual
                                    with st.expander(f"üìã An√°lise Detalhada - Imagem {idx+1}", expanded=True):
                                        try:
                                            # Construir contexto com base de conhecimento do agente
                                            contexto_agente = ""
                                            if "base_conhecimento" in agente:
                                                contexto_agente = f"""
                                                ###BEGIN DIRETRIZES DE BRANDING DO AGENTE:###
                                                {agente['base_conhecimento']}
                                                ###END DIRETRIZES DE BRANDING DO AGENTE###
                                                """
                                            
                                            # Adicionar contexto global se fornecido
                                            contexto_completo = contexto_agente
                                            if contexto_global and contexto_global.strip():
                                                contexto_completo += f"""
                                                ###BEGIN CONTEXTO ADICIONAL DO USUARIO###
                                                {contexto_global}
                                                ###END CONTEXTO ADICIONAL DO USUARIO###
                                                """
                                            
                                            # Escolher m√©todo de an√°lise
                                            if st.session_state.analise_especializada_imagem:
                                                # AN√ÅLISE ESPECIALIZADA POR M√öLTIPLOS ESPECIALISTAS VISUAIS
                                                st.info("üéØ **Executando an√°lise especializada por m√∫ltiplos especialistas visuais...**")
                                                
                                                # Criar analisadores especialistas
                                                analisadores_config = criar_analisadores_imagem(contexto_completo, "")
                                                
                                                # Filtrar apenas os selecionados
                                                analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                         if k in st.session_state.analisadores_selecionados_imagem}
                                                
                                                # Executar an√°lises especializadas
                                                resultados_especialistas = executar_analise_imagem_especializada(
                                                    uploaded_image, 
                                                    uploaded_image.name, 
                                                    analisadores_filtrados
                                                )
                                                
                                                # Gerar relat√≥rio consolidado
                                                relatorio_consolidado, score_medio, status = gerar_relatorio_imagem_consolidado(
                                                    resultados_especialistas, 
                                                    uploaded_image.name,
                                                    f"{image.width}x{image.height}"
                                                )
                                                
                                                st.markdown(relatorio_consolidado, unsafe_allow_html=True)
                                                
                                                # Armazenar resultado
                                                resultados_analise.append({
                                                    'nome': uploaded_image.name,
                                                    'indice': idx,
                                                    'analise': relatorio_consolidado,
                                                    'dimensoes': f"{image.width}x{image.height}",
                                                    'tamanho': uploaded_image.size,
                                                    'score_medio': score_medio,
                                                    'status': status
                                                })
                                                
                                            else:
                                                # An√°lise geral da imagem (m√©todo antigo)
                                                prompt_analise = f"""
                                                {contexto_completo}
                                                
                                                Analise esta imagem e verifique o alinhamento com as diretrizes de branding.
                                                
                                                Forne√ßa a an√°lise em formato claro:
                                                
                                                ## üñºÔ∏è RELAT√ìRIO DE ALINHAMENTO - IMAGEM {idx+1}
                                                
                                                **Arquivo:** {uploaded_image.name}
                                                **Dimens√µes:** {image.width} x {image.height}
                                                
                                                ### üéØ RESUMO DA IMAGEM
                                                [Avalia√ß√£o geral de conformidade visual e textual]
                                                
                                                ### ‚úÖ ELEMENTOS ALINHADOS 
                                                - [Itens visuais e textuais que seguem as diretrizes]
                                                
                                                ### ‚ö†Ô∏è ELEMENTOS FORA DO PADR√ÉO
                                                - [Itens visuais e textuais que n√£o seguem as diretrizes]
                                                
                                                ### üí° RECOMENDA√á√ïES
                                                - [Sugest√µes para melhorar o alinhamento visual e textual]
                                                
                                                ### üé® ASPECTOS T√âCNICOS
                                                - [Composi√ß√£o, cores, tipografia, etc.]
                                                """
                                                
                                                # Processar imagem
                                                response = modelo_vision.generate_content([
                                                    prompt_analise,
                                                    {"mime_type": "image/jpeg", "data": uploaded_image.getvalue()}
                                                ])
                                                
                                                st.markdown(response.text)
                                                
                                                # Armazenar resultado
                                                resultados_analise.append({
                                                    'nome': uploaded_image.name,
                                                    'indice': idx,
                                                    'analise': response.text,
                                                    'dimensoes': f"{image.width}x{image.height}",
                                                    'tamanho': uploaded_image.size
                                                })
                                            
                                        except Exception as e:
                                            st.error(f"‚ùå Erro ao processar imagem {uploaded_image.name}: {str(e)}")
                                
                                # Separador visual entre imagens
                                if idx < len(uploaded_images) - 1:
                                    st.markdown("---")
                                    
                            except Exception as e:
                                st.error(f"‚ùå Erro ao carregar imagem {uploaded_image.name}: {str(e)}")
                    
                    # Armazenar na sess√£o
                    st.session_state.resultados_analise_imagem = resultados_analise
                    
                    # Resumo executivo
                    st.markdown("---")
                    st.subheader("üìã Resumo Executivo de Imagens")
                    
                    if resultados_analise and all('score_medio' in resultado for resultado in resultados_analise):
                        # Calcular estat√≠sticas com scores
                        scores = [resultado['score_medio'] for resultado in resultados_analise if 'score_medio' in resultado]
                        score_medio_geral = sum(scores) / len(scores) if scores else 0
                        
                        col_resumo1, col_resumo2, col_resumo3, col_resumo4 = st.columns(4)
                        with col_resumo1:
                            st.metric("üìä Total de Imagens", len(uploaded_images))
                        with col_resumo2:
                            st.metric("‚úÖ An√°lises Conclu√≠das", len(resultados_analise))
                        with col_resumo3:
                            st.metric("‚≠ê Score M√©dio", f"{score_medio_geral:.1f}/10")
                        with col_resumo4:
                            aprovadas = sum(1 for r in resultados_analise if r.get('status') == '‚úÖ APROVADO')
                            st.metric("üéØ Aprovadas", aprovadas)
                    else:
                        col_resumo1, col_resumo2, col_resumo3 = st.columns(3)
                        with col_resumo1:
                            st.metric("üìä Total de Imagens", len(uploaded_images))
                        with col_resumo2:
                            st.metric("‚úÖ An√°lises Conclu√≠das", len(resultados_analise))
                        with col_resumo3:
                            st.metric("üñºÔ∏è Processadas", len(uploaded_images))
                    
                    # Contexto aplicado no resumo
                    if contexto_global and contexto_global.strip():
                        st.info(f"**üéØ Contexto Aplicado em Todas as An√°lises:** {contexto_global}")
                    
                    # Bot√£o para download do relat√≥rio consolidado
                    if st.button("üì• Exportar Relat√≥rio Completo de Imagens", key="exportar_relatorio_imagens"):
                        relatorio = f"""
                        # RELAT√ìRIO DE VALIDA√á√ÉO DE IMAGENS
                        
                        **Agente:** {agente.get('nome', 'N/A')}
                        **Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
                        **Total de Imagens:** {len(uploaded_images)}
                        **Contexto Aplicado:** {contexto_global if contexto_global else 'Nenhum contexto adicional'}
                        **M√©todo de An√°lise:** {'Especializada por M√∫ltiplos Especialistas' if st.session_state.analise_especializada_imagem else 'Tradicional'}
                        
                        ## RESUMO EXECUTIVO
                        {chr(10).join([f"{idx+1}. {img.name}" for idx, img in enumerate(uploaded_images)])}
                        
                        ## AN√ÅLISES INDIVIDUAIS
                        {chr(10).join([f'### {res["nome"]} {chr(10)}{res["analise"]}' for res in resultados_analise])}
                        """
                        
                        st.download_button(
                            "üíæ Baixar Relat√≥rio em TXT",
                            data=relatorio,
                            file_name=f"relatorio_validacao_imagens_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
            
            # Mostrar an√°lises existentes da sess√£o
            elif st.session_state.resultados_analise_imagem:
                st.info("üìã An√°lises anteriores encontradas. Use o bot√£o 'Limpar An√°lises' para recome√ßar.")
                
                for resultado in st.session_state.resultados_analise_imagem:
                    with st.expander(f"üñºÔ∏è {resultado['nome']} - An√°lise Salva", expanded=False):
                        st.markdown(resultado['analise'])
            
            else:
                st.info("üìÅ Carregue uma ou mais imagens para iniciar a valida√ß√£o de branding")

        with subtab_video:
            st.subheader("üé¨ Valida√ß√£o de V√≠deo")
            
            # Bot√£o para limpar an√°lises de v√≠deo
            if st.button("üóëÔ∏è Limpar An√°lises de V√≠deo", key="limpar_analises_video"):
                st.session_state.resultados_analise_video = []
                st.rerun()
            
            # Container principal
            col_upload, col_config = st.columns([2, 1])
            
            with col_upload:
                uploaded_videos = st.file_uploader(
                    "Carregue um ou mais v√≠deos para an√°lise",
                    type=["mp4", "mpeg", "mov", "avi", "flv", "mpg", "webm", "wmv", "3gpp"],
                    key="video_upload_validacao",
                    accept_multiple_files=True
                )
            
            with col_config:
                st.markdown("### ‚öôÔ∏è Configura√ß√µes de V√≠deo")
                contexto_video_especifico = st.text_area(
                    "**üéØ Contexto espec√≠fico para v√≠deos:**", 
                    height=120, 
                    key="video_context_especifico",
                    placeholder="Contexto adicional espec√≠fico para an√°lise de v√≠deos (opcional)..."
                )
                
                analise_especializada_video = st.checkbox(
                    "An√°lise especializada por √°reas (recomendado)",
                    value=st.session_state.analise_especializada_video,
                    help="Usa m√∫ltiplos especialistas em v√≠deo para an√°lise mais precisa",
                    key="analise_especializada_video_check"
                )
                
                analisadores_selecionados_video = st.multiselect(
                    "Especialistas de v√≠deo a incluir:",
                    options=['narrativa_estrutura', 'qualidade_audio', 'visual_cinematografia', 'branding_consistencia', 'engajamento_eficacia'],
                    default=st.session_state.analisadores_selecionados_video,
                    format_func=lambda x: {
                        'narrativa_estrutura': 'üìñ Narrativa e Estrutura',
                        'qualidade_audio': 'üîä Qualidade de √Åudio', 
                        'visual_cinematografia': 'üé• Visual e Cinematografia',
                        'branding_consistencia': 'üè¢ Branding e Consist√™ncia',
                        'engajamento_eficacia': 'üìà Engajamento e Efic√°cia'
                    }[x],
                    key="analisadores_video_select"
                )
            
            if uploaded_videos:
                st.success(f"‚úÖ {len(uploaded_videos)} v√≠deo(s) carregado(s)")
                
                # Contexto aplicado
                if contexto_global and contexto_global.strip():
                    st.info(f"**üéØ Contexto Global Aplicado:** {contexto_global}")
                if contexto_video_especifico and contexto_video_especifico.strip():
                    st.info(f"**üéØ Contexto Espec√≠fico para V√≠deos:** {contexto_video_especifico}")
                
                # Exibir informa√ß√µes dos v√≠deos
                st.markdown("### üìä Informa√ß√µes dos V√≠deos")
                
                for idx, video in enumerate(uploaded_videos):
                    col_vid, col_info, col_actions = st.columns([2, 2, 1])
                    
                    with col_vid:
                        st.write(f"**{idx+1}. {video.name}**")
                        st.caption(f"Tipo: {video.type} | Tamanho: {video.size / (1024*1024):.1f} MB")
                    
                    with col_info:
                        st.write("üìè Dura√ß√£o: A ser detectada")
                        st.write("üéûÔ∏è Resolu√ß√£o: A ser detectada")
                    
                    with col_actions:
                        if st.button("üîç Preview", key=f"preview_{idx}"):
                            st.video(video, format=f"video/{video.type.split('/')[-1]}")
                
                # Bot√£o para validar todos os v√≠deos
                if st.button("üé¨ Validar Todos os V√≠deos", type="primary", key="validar_videos_multiplas"):
                    
                    resultados_video = []
                    
                    for idx, uploaded_video in enumerate(uploaded_videos):
                        with st.spinner(f'Analisando v√≠deo {idx+1} de {len(uploaded_videos)}: {uploaded_video.name}...'):
                            try:
                                # Container para cada v√≠deo
                                with st.container():
                                    st.markdown("---")
                                    
                                    # Header do v√≠deo
                                    col_header, col_stats = st.columns([3, 1])
                                    
                                    with col_header:
                                        st.subheader(f"üé¨ {uploaded_video.name}")
                                    
                                    with col_stats:
                                        st.metric("üìä Status", "Processando")
                                    
                                    # Contexto aplicado para este v√≠deo
                                    if contexto_global and contexto_global.strip():
                                        st.info(f"**üéØ Contexto Aplicado:** {contexto_global}")
                                    if contexto_video_especifico and contexto_video_especifico.strip():
                                        st.info(f"**üéØ Contexto Espec√≠fico:** {contexto_video_especifico}")
                                    
                                    # Preview do v√≠deo
                                    with st.expander("üëÄ Preview do V√≠deo", expanded=False):
                                        st.video(uploaded_video, format=f"video/{uploaded_video.type.split('/')[-1]}")
                                    
                                    # An√°lise detalhada
                                    with st.expander(f"üìã An√°lise Completa - {uploaded_video.name}", expanded=True):
                                        try:
                                            # Construir contexto com base de conhecimento do agente
                                            contexto_agente = ""
                                            if "base_conhecimento" in agente:
                                                contexto_agente = f"""
                                                ###BEGIN DIRETRIZES DE BRANDING DO AGENTE:###
                                                {agente['base_conhecimento']}
                                                ###END DIRETRIZES DE BRANDING DO AGENTE###
                                                """
                                            
                                            # Adicionar contexto global se fornecido
                                            contexto_completo = contexto_agente
                                            if contexto_global and contexto_global.strip():
                                                contexto_completo += f"""
                                                ###BEGIN CONTEXTO GLOBAL DO USUARIO###
                                                {contexto_global}
                                                ###END CONTEXTO GLOBAL DO USUARIO###
                                                """
                                            
                                            # Adicionar contexto espec√≠fico de v√≠deo se fornecido
                                            if contexto_video_especifico and contexto_video_especifico.strip():
                                                contexto_completo += f"""
                                                ###BEGIN CONTEXTO ESPEC√çFICO PARA V√çDEOS###
                                                {contexto_video_especifico}
                                                ###END CONTEXTO ESPEC√çFICO PARA V√çDEOS###
                                                """
                                            
                                            # Escolher m√©todo de an√°lise
                                            if st.session_state.analise_especializada_video:
                                                # AN√ÅLISE ESPECIALIZADA POR M√öLTIPLOS ESPECIALISTAS DE V√çDEO
                                                st.info("üéØ **Executando an√°lise especializada por m√∫ltiplos especialistas de v√≠deo...**")
                                                
                                                # Criar analisadores especialistas
                                                analisadores_config = criar_analisadores_video(contexto_agente, contexto_global, contexto_video_especifico)
                                                
                                                # Filtrar apenas os selecionados
                                                analisadores_filtrados = {k: v for k, v in analisadores_config.items() 
                                                                         if k in st.session_state.analisadores_selecionados_video}
                                                
                                                # Executar an√°lises especializadas
                                                resultados_especialistas = executar_analise_video_especializada(
                                                    uploaded_video, 
                                                    uploaded_video.name, 
                                                    analisadores_filtrados
                                                )
                                                
                                                # Gerar relat√≥rio consolidado
                                                relatorio_consolidado, score_medio, status = gerar_relatorio_video_consolidado(
                                                    resultados_especialistas, 
                                                    uploaded_video.name,
                                                    uploaded_video.type
                                                )
                                                
                                                st.markdown(relatorio_consolidado, unsafe_allow_html=True)
                                                
                                                # Armazenar resultado
                                                resultados_video.append({
                                                    'nome': uploaded_video.name,
                                                    'indice': idx,
                                                    'analise': relatorio_consolidado,
                                                    'tipo': uploaded_video.type,
                                                    'tamanho': uploaded_video.size,
                                                    'score_medio': score_medio,
                                                    'status': status
                                                })
                                                
                                            else:
                                                # An√°lise geral do v√≠deo (m√©todo antigo)
                                                prompt_analise = f"""
                                                {contexto_completo}
                                                
                                                Analise este v√≠deo considerando:
                                                - Alinhamento com diretrizes de branding
                                                - Qualidade e consist√™ncia visual  
                                                - Mensagem e tom da comunica√ß√£o
                                                - Elementos de √°udio e transcri√ß√£o
                                                - Texto presente nos frames
                                                
                                                Forne√ßa a an√°lise em formato estruturado:
                                                
                                                ## üé¨ RELAT√ìRIO DE ALINHAMENTO - V√çDEO {idx+1}
                                                
                                                **Arquivo:** {uploaded_video.name}
                                                **Formato:** {uploaded_video.type}
                                                
                                                ### üéØ RESUMO EXECUTIVO
                                                [Avalia√ß√£o geral do alinhamento do v√≠deo com as diretrizes]
                                                
                                                ### üîä AN√ÅLISE DE √ÅUDIO
                                                [Transcri√ß√£o e an√°lise do conte√∫do de √°udio, tom, mensagem verbal]
                                                
                                                ### üëÅÔ∏è AN√ÅLISE VISUAL
                                                [An√°lise de elementos visuais, cores, composi√ß√£o, branding visual]

                                                ### üìù TEXTO EM FRAMES
                                                [Identifica√ß√£o e an√°lise de texto presente nos frames]
                                                
                                                ### ‚úÖ PONTOS FORTES
                                                - [Elementos bem alinhados com as diretrizes]
                                                
                                                ### ‚ö†Ô∏è PONTOS DE ATEN√á√ÉO
                                                - [Desvios identificados e timestamps espec√≠ficos]
                                                
                                                ### üí° RECOMENDA√á√ïES
                                                - [Sugest√µes para melhorar o alinhamento]
                                                
                                                ### üïí MOMENTOS CHAVE
                                                [Timestamps importantes com descri√ß√£o: MM:SS]
                                                """
                                                
                                                # Processar v√≠deo usando a API do Gemini
                                                video_bytes = uploaded_video.getvalue()
                                                
                                                if len(video_bytes) < 200 * 1024 * 1024:
                                                    response = modelo_vision.generate_content([
                                                        prompt_analise,
                                                        {"mime_type": uploaded_video.type, "data": video_bytes}
                                                    ])
                                                else:
                                                    st.info("üì§ Uploading v√≠deo para processamento...")
                                                    response = modelo_vision.generate_content([
                                                        prompt_analise,
                                                        {"mime_type": uploaded_video.type, "data": video_bytes}
                                                    ])
                                                
                                                st.markdown(response.text)
                                                
                                                # Armazenar resultado
                                                resultados_video.append({
                                                    'nome': uploaded_video.name,
                                                    'indice': idx,
                                                    'analise': response.text,
                                                    'tipo': uploaded_video.type,
                                                    'tamanho': uploaded_video.size
                                                })
                                            
                                        except Exception as e:
                                            st.error(f"‚ùå Erro ao processar v√≠deo {uploaded_video.name}: {str(e)}")
                                            resultados_video.append({
                                                'nome': uploaded_video.name,
                                                'indice': idx,
                                                'analise': f"Erro na an√°lise: {str(e)}",
                                                'tipo': uploaded_video.type,
                                                'tamanho': uploaded_video.size
                                            })
                                
                                # Separador entre v√≠deos
                                if idx < len(uploaded_videos) - 1:
                                    st.markdown("---")
                                    
                            except Exception as e:
                                st.error(f"‚ùå Erro ao processar v√≠deo {uploaded_video.name}: {str(e)}")
                    
                    # Armazenar resultados na sess√£o
                    st.session_state.resultados_analise_video = resultados_video
                    
                    # Resumo executivo dos v√≠deos
                    st.markdown("---")
                    st.subheader("üìã Resumo Executivo - V√≠deos")
                    
                    if resultados_video and all('score_medio' in resultado for resultado in resultados_video):
                        # Calcular estat√≠sticas com scores
                        scores = [resultado['score_medio'] for resultado in resultados_video if 'score_medio' in resultado]
                        score_medio_geral = sum(scores) / len(scores) if scores else 0
                        
                        col_vid1, col_vid2, col_vid3, col_vid4 = st.columns(4)
                        with col_vid1:
                            st.metric("üé¨ Total de V√≠deos", len(uploaded_videos))
                        with col_vid2:
                            st.metric("‚úÖ An√°lises Conclu√≠das", len(resultados_video))
                        with col_vid3:
                            st.metric("‚≠ê Score M√©dio", f"{score_medio_geral:.1f}/10")
                        with col_vid4:
                            aprovados = sum(1 for r in resultados_video if r.get('status') == '‚úÖ APROVADO')
                            st.metric("üéØ Aprovados", aprovados)
                    else:
                        col_vid1, col_vid2 = st.columns(2)
                        with col_vid1:
                            st.metric("üé¨ Total de V√≠deos", len(uploaded_videos))
                        with col_vid2:
                            st.metric("‚úÖ An√°lises Conclu√≠das", len(resultados_video))
                    
                    # Contexto aplicado no resumo
                    if contexto_global and contexto_global.strip():
                        st.info(f"**üéØ Contexto Global Aplicado:** {contexto_global}")
                    if contexto_video_especifico and contexto_video_especifico.strip():
                        st.info(f"**üéØ Contexto Espec√≠fico Aplicado:** {contexto_video_especifico}")
                    
                    # Bot√£o para download do relat√≥rio
                    if st.button("üì• Exportar Relat√≥rio de V√≠deos", key="exportar_relatorio_videos"):
                        relatorio_videos = f"""
                        # RELAT√ìRIO DE VALIDA√á√ÉO DE V√çDEOS
                        
                        **Agente:** {agente.get('nome', 'N/A')}
                        **Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}
                        **Total de V√≠deos:** {len(uploaded_videos)}
                        **Contexto Global:** {contexto_global if contexto_global else 'Nenhum'}
                        **Contexto Espec√≠fico:** {contexto_video_especifico if contexto_video_especifico else 'Nenhum'}
                        **M√©todo de An√°lise:** {'Especializada por M√∫ltiplos Especialistas' if st.session_state.analise_especializada_video else 'Tradicional'}
                        
                        ## V√çDEOS ANALISADOS:
                        {chr(10).join([f"{idx+1}. {vid.name} ({vid.type}) - {vid.size/(1024*1024):.1f} MB" for idx, vid in enumerate(uploaded_videos)])}
                        
                        ## AN√ÅLISES INDIVIDUAIS:
                        {chr(10).join([f'### {res["nome"]} {chr(10)}{res["analise"]}' for res in resultados_video])}
                        """
                        
                        st.download_button(
                            "üíæ Baixar Relat√≥rio em TXT",
                            data=relatorio_videos,
                            file_name=f"relatorio_validacao_videos_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain"
                        )
            
            # Mostrar an√°lises existentes da sess√£o
            elif st.session_state.resultados_analise_video:
                st.info("üìã An√°lises anteriores encontradas. Use o bot√£o 'Limpar An√°lises' para recome√ßar.")
                
                for resultado in st.session_state.resultados_analise_video:
                    with st.expander(f"üé¨ {resultado['nome']} - An√°lise Salva", expanded=False):
                        st.markdown(resultado['analise'])
            
            else:
                st.info("üé¨ Carregue um ou mais v√≠deos para iniciar a valida√ß√£o")
# --- ABA: GERA√á√ÉO DE CONTE√öDO ---
with tab_mapping["‚ú® Gera√ß√£o de Conte√∫do"]:
    st.header("‚ú® Gera√ß√£o de Conte√∫do com M√∫ltiplos Insumos")
    
    # Conex√£o com MongoDB para briefings
    try:
        client2 = MongoClient("mongodb+srv://gustavoromao3345:RqWFPNOJQfInAW1N@cluster0.5iilj.mongodb.net/auto_doc?retryWrites=true&w=majority&ssl=true&ssl_cert_reqs=CERT_NONE&tlsAllowInvalidCertificates=true")
        db_briefings = client2['briefings_Broto_Tecnologia']
        collection_briefings = db_briefings['briefings']
        mongo_connected_conteudo = True
    except Exception as e:
        st.error(f"Erro na conex√£o com MongoDB: {str(e)}")
        mongo_connected_conteudo = False

    # Fun√ß√£o para extrair texto de diferentes tipos de arquivo
    def extrair_texto_arquivo(arquivo):
        """Extrai texto de diferentes formatos de arquivo"""
        try:
            extensao = arquivo.name.split('.')[-1].lower()
            
            if extensao == 'pdf':
                return extrair_texto_pdf(arquivo)
            elif extensao == 'txt':
                return extrair_texto_txt(arquivo)
            elif extensao in ['pptx', 'ppt']:
                return extrair_texto_pptx(arquivo)
            elif extensao in ['docx', 'doc']:
                return extrair_texto_docx(arquivo)
            else:
                return f"Formato {extensao} n√£o suportado para extra√ß√£o de texto."
                
        except Exception as e:
            return f"Erro ao extrair texto do arquivo {arquivo.name}: {str(e)}"

    def extrair_texto_pdf(arquivo):
        """Extrai texto de arquivos PDF"""
        try:
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(arquivo)
            texto = ""
            for pagina in pdf_reader.pages:
                texto += pagina.extract_text() + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PDF: {str(e)}"

    def extrair_texto_txt(arquivo):
        """Extrai texto de arquivos TXT"""
        try:
            return arquivo.read().decode('utf-8')
        except:
            try:
                return arquivo.read().decode('latin-1')
            except Exception as e:
                return f"Erro na leitura do TXT: {str(e)}"

    def extrair_texto_pptx(arquivo):
        """Extrai texto de arquivos PowerPoint"""
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(arquivo.read()))
            texto = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PowerPoint: {str(e)}"

    def extrair_texto_docx(arquivo):
        """Extrai texto de arquivos Word"""
        try:
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do Word: {str(e)}"

    # Layout principal
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("üìù Fontes de Conte√∫do")
        
        # Op√ß√£o 1: Upload de m√∫ltiplos arquivos
        st.write("üìé Upload de Arquivos (PDF, TXT, PPTX, DOCX):")
        arquivos_upload = st.file_uploader(
            "Selecione um ou mais arquivos:",
            type=['pdf', 'txt', 'pptx', 'ppt', 'docx', 'doc'],
            accept_multiple_files=True,
            help="Arquivos ser√£o convertidos para texto e usados como base para gera√ß√£o de conte√∫do"
        )
        
        # Processar arquivos uploadados
        textos_arquivos = ""
        if arquivos_upload:
            st.success(f"‚úÖ {len(arquivos_upload)} arquivo(s) carregado(s)")
            
            with st.expander("üìã Visualizar Conte√∫do dos Arquivos", expanded=False):
                for i, arquivo in enumerate(arquivos_upload):
                    st.write(f"**{arquivo.name}** ({arquivo.size} bytes)")
                    with st.spinner(f"Processando {arquivo.name}..."):
                        texto_extraido = extrair_texto_arquivo(arquivo)
                        textos_arquivos += f"\n\n--- CONTE√öDO DE {arquivo.name.upper()} ---\n{texto_extraido}"
                        
                        # Mostrar preview
                        if len(texto_extraido) > 500:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido[:500] + "...", 
                                       height=100,
                                       key=f"preview_{i}")
                        else:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido, 
                                       height=100,
                                       key=f"preview_{i}")
        
        
        
        # Op√ß√£o 3: Inserir briefing manualmente
        st.write("‚úçÔ∏è Briefing Manual:")
        briefing_manual = st.text_area("Ou cole o briefing completo aqui:", height=150,
                                      placeholder="""Exemplo:
T√≠tulo: Campanha de Lan√ßamento
Objetivo: Divulgar novo produto
P√∫blico-alvo: Empres√°rios...
Pontos-chave: [lista os principais pontos]""")
        
        # Transcri√ß√£o de √°udio/v√≠deo
        st.write("üé§ Transcri√ß√£o de √Åudio/Video:")
        arquivos_midia = st.file_uploader(
            "√Åudios/V√≠deos para transcri√ß√£o:",
            type=['mp3', 'wav', 'mp4', 'mov', 'avi'],
            accept_multiple_files=True,
            help="Arquivos de m√≠dia ser√£o transcritos automaticamente"
        )
        
        transcricoes_texto = ""
        if arquivos_midia:
            st.info(f"üé¨ {len(arquivos_midia)} arquivo(s) de m√≠dia carregado(s)")
            if st.button("üîÑ Transcrever Todos os Arquivos de M√≠dia"):
                with st.spinner("Transcrevendo arquivos de m√≠dia..."):
                    for arquivo in arquivos_midia:
                        tipo = "audio" if arquivo.type.startswith('audio') else "video"
                        transcricao = transcrever_audio_video(arquivo, tipo)
                        transcricoes_texto += f"\n\n--- TRANSCRI√á√ÉO DE {arquivo.name.upper()} ---\n{transcricao}"
                        st.success(f"‚úÖ {arquivo.name} transcrito!")
    
    with col2:
        st.subheader("‚öôÔ∏è Configura√ß√µes de Gera√ß√£o")
        
        # Op√ß√£o para o usu√°rio escolher entre configura√ß√µes padr√£o ou prompt personalizado
        modo_geracao = st.radio(
            "Modo de Gera√ß√£o:",
            ["Configura√ß√µes Padr√£o", "Prompt Personalizado"],
            help="Escolha entre usar configura√ß√µes pr√©-definidas ou escrever seu pr√≥prio prompt"
        )
        
        if modo_geracao == "Configura√ß√µes Padr√£o":
            # Configura√ß√µes b√°sicas (vers√£o simplificada)
            tipo_conteudo = st.selectbox("Tipo de Conte√∫do:", 
                                       ["Post Social", "Artigo Blog", "Email Marketing", 
                                        "Landing Page", "Script V√≠deo", "Relat√≥rio T√©cnico",
                                        "Press Release", "Newsletter", "Case Study"])
            
            tom_voz = st.selectbox("Tom de Voz:", 
                                  ["Formal", "Informal", "Persuasivo", "Educativo", 
                                   "Inspirador", "T√©cnico", "Jornal√≠stico"])
            
            palavras_chave = st.text_input("Palavras-chave (opcional):",
                                          placeholder="separadas por v√≠rgula")
            
            numero_palavras = st.slider("N√∫mero de Palavras:", 100, 3000, 800)
            
            # Configura√ß√µes avan√ßadas simplificadas
            with st.expander("üîß Configura√ß√µes Avan√ßadas"):
                usar_contexto_agente = st.checkbox("Usar contexto do agente selecionado", 
                                                 value=bool(st.session_state.agente_selecionado))
                
                incluir_cta = st.checkbox("Incluir Call-to-Action", value=True)
                
                formato_saida = st.selectbox("Formato de Sa√≠da:", 
                                           ["Texto Simples", "Markdown", "HTML B√°sico"])
        
        else:  # Prompt Personalizado
            st.info("üí° Escreva seu pr√≥prio prompt de gera√ß√£o. Use {contexto} para incluir automaticamente todas as fontes de conte√∫do.")
            prompt_personalizado = st.text_area(
                "Seu Prompt Personalizado:",
                height=200,
                placeholder="""Exemplo:
Com base no contexto fornecido, crie um artigo detalhado que:

1. Explique os conceitos principais de forma clara
2. Destaque os benef√≠cios para o p√∫blico-alvo
3. Inclua exemplos pr√°ticos de aplica√ß√£o
4. Mantenha um tom {tom} e acess√≠vel

Contexto: {contexto}

Gere o conte√∫do em formato {formato} com aproximadamente {palavras} palavras."""
            )
            
            # Vari√°veis que o usu√°rio pode usar no prompt personalizado
            col_var1, col_var2, col_var3 = st.columns(3)
            with col_var1:
                tom_personalizado = st.selectbox("Tom:", 
                                               ["formal", "informal", "persuasivo", "educativo"], 
                                               key="tom_personalizado")
            with col_var2:
                formato_personalizado = st.selectbox("Formato:", 
                                                   ["texto simples", "markdown", "HTML b√°sico"], 
                                                   key="formato_personalizado")
            with col_var3:
                palavras_personalizado = st.slider("Palavras:", 100, 3000, 800, key="palavras_personalizado")
            
            usar_contexto_agente = st.checkbox("Usar contexto do agente selecionado", 
                                             value=bool(st.session_state.agente_selecionado),
                                             key="contexto_personalizado")

    # √Årea de instru√ß√µes espec√≠ficas (apenas para modo padr√£o)
    if modo_geracao == "Configura√ß√µes Padr√£o":
        st.subheader("üéØ Instru√ß√µes Espec√≠ficas")
        instrucoes_especificas = st.text_area(
            "Diretrizes adicionais para gera√ß√£o:",
            placeholder="""Exemplos:
- Focar nos benef√≠cios para o usu√°rio final
- Incluir estat√≠sticas quando poss√≠vel
- Manter linguagem acess√≠vel
- Evitar jarg√µes t√©cnicos excessivos
- Seguir estrutura: problema ‚Üí solu√ß√£o ‚Üí benef√≠cios""",
            height=100
        )

    # Bot√£o para gerar conte√∫do
    if st.button("üöÄ Gerar Conte√∫do com Todos os Insumos", type="primary", use_container_width=True):
        # Verificar se h√° pelo menos uma fonte de conte√∫do
        tem_conteudo = (arquivos_upload or 
                       briefing_manual or 
                       ('briefing_data' in locals() and briefing_data) or
                       arquivos_midia)
        
        if not tem_conteudo:
            st.error("‚ùå Por favor, forne√ßa pelo menos uma fonte de conte√∫do (arquivos, briefing ou m√≠dia)")
        elif modo_geracao == "Prompt Personalizado" and not prompt_personalizado:
            st.error("‚ùå Por favor, escreva um prompt personalizado para gera√ß√£o")
        else:
            with st.spinner("Processando todos os insumos e gerando conte√∫do..."):
                try:
                    # Construir o contexto combinado de todas as fontes
                    contexto_completo = "## FONTES DE CONTE√öDO COMBINADAS:\n\n"
                    
                    # Adicionar conte√∫do dos arquivos uploadados
                    if textos_arquivos:
                        contexto_completo += "### CONTE√öDO DOS ARQUIVOS:\n" + textos_arquivos + "\n\n"
                    
                    # Adicionar briefing do banco ou manual
                    if briefing_manual:
                        contexto_completo += "### BRIEFING MANUAL:\n" + briefing_manual + "\n\n"
                    elif 'briefing_data' in locals() and briefing_data:
                        contexto_completo += "### BRIEFING DO BANCO:\n" + briefing_data['conteudo'] + "\n\n"
                    
                    # Adicionar transcri√ß√µes
                    if transcricoes_texto:
                        contexto_completo += "### TRANSCRI√á√ïES DE M√çDIA:\n" + transcricoes_texto + "\n\n"
                    
                    # Adicionar contexto do agente se selecionado
                    contexto_agente = ""
                    if usar_contexto_agente and st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                    
                    # Construir prompt final baseado no modo selecionado
                    if modo_geracao == "Configura√ß√µes Padr√£o":
                        prompt_final = f"""
                        {contexto_agente}
                        
                        ## INSTRU√á√ïES PARA GERA√á√ÉO DE CONTE√öDO:
                        
                        **TIPO DE CONTE√öDO:** {tipo_conteudo}
                        **TOM DE VOZ:** {tom_voz}
                        **PALAVRAS-CHAVE:** {palavras_chave if palavras_chave else 'N√£o especificadas'}
                        **N√öMERO DE PALAVRAS:** {numero_palavras} (¬±10%)
                        **INCLUIR CALL-TO-ACTION:** {incluir_cta}
                        
                        **INSTRU√á√ïES ESPEC√çFICAS:**
                        {instrucoes_especificas if instrucoes_especificas else 'Nenhuma instru√ß√£o espec√≠fica fornecida.'}
                        
                        ## FONTES E REFER√äNCIAS:
                        {contexto_completo}
                        
                        ## TAREFA:
                        Com base em TODAS as fontes fornecidas acima, gere um conte√∫do do tipo {tipo_conteudo} que:
                        
                        1. **S√≠ntese Eficiente:** Combine e sintetize informa√ß√µes de todas as fontes
                        2. **Coer√™ncia:** Mantenha consist√™ncia com as informa√ß√µes originais
                        3. **Valor Agregado:** V√° al√©m da simples c√≥pia, agregando insights
                        4. **Engajamento:** Crie conte√∫do que engaje o p√∫blico-alvo
                        5. **Clareza:** Comunique ideias complexas de forma acess√≠vel
                        
                        **FORMATO DE SA√çDA:** {formato_saida}
                        
                        Gere um conte√∫do completo e profissional.
                        """
                    else:  # Prompt Personalizado
                        # Substituir vari√°veis no prompt personalizado
                        prompt_processado = prompt_personalizado.replace("{contexto}", contexto_completo)
                        prompt_processado = prompt_processado.replace("{tom}", tom_personalizado)
                        prompt_processado = prompt_processado.replace("{formato}", formato_personalizado)
                        prompt_processado = prompt_processado.replace("{palavras}", str(palavras_personalizado))
                        
                        prompt_final = f"""
                        {contexto_agente}
                        
                        {prompt_processado}
                        """
                    
                    resposta = modelo_texto.generate_content(prompt_final)
                    
                    # Determinar formato de sa√≠da baseado no modo
                    if modo_geracao == "Configura√ß√µes Padr√£o":
                        formato_output = formato_saida
                    else:
                        formato_output = formato_personalizado
                    
                    # Processar sa√≠da baseada no formato selecionado
                    conteudo_gerado = resposta.text
                    
                    if formato_output == "HTML B√°sico" or formato_output == "HTML b√°sico":
                        # Converter markdown para HTML b√°sico
                        import re
                        conteudo_gerado = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'\*(.*?)\*', r'<em>\1</em>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'### (.*?)\n', r'<h3>\1</h3>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'## (.*?)\n', r'<h2>\1</h2>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'# (.*?)\n', r'<h1>\1</h1>', conteudo_gerado)
                        conteudo_gerado = conteudo_gerado.replace('\n', '<br>')
                    
                    st.subheader("üìÑ Conte√∫do Gerado")
                    
                    if formato_output == "HTML B√°sico" or formato_output == "HTML b√°sico":
                        st.components.v1.html(conteudo_gerado, height=400, scrolling=True)
                    else:
                        st.markdown(conteudo_gerado)
                    
                    # Estat√≠sticas
                    palavras_count = len(conteudo_gerado.split())
                    col_stat1, col_stat2, col_stat3 = st.columns(3)
                    with col_stat1:
                        st.metric("Palavras Geradas", palavras_count)
                    with col_stat2:
                        st.metric("Arquivos Processados", len(arquivos_upload) if arquivos_upload else 0)
                    with col_stat3:
                        st.metric("Fontes Utilizadas", 
                                 (1 if arquivos_upload else 0) + 
                                 (1 if briefing_manual or 'briefing_data' in locals() else 0) +
                                 (1 if transcricoes_texto else 0))
                    
                    # Bot√µes de download
                    extensao = ".html" if "HTML" in formato_output else ".md" if "markdown" in formato_output.lower() else ".txt"
                    
                    st.download_button(
                        f"üíæ Baixar Conte√∫do ({formato_output})",
                        data=conteudo_gerado,
                        file_name=f"conteudo_gerado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}{extensao}",
                        mime="text/html" if "HTML" in formato_output else "text/plain"
                    )
                    
                    # Salvar no hist√≥rico se MongoDB dispon√≠vel
                    if mongo_connected_conteudo:
                        try:
                            from bson import ObjectId
                            historico_data = {
                                "modo_geracao": modo_geracao,
                                "tipo_conteudo": tipo_conteudo if modo_geracao == "Configura√ß√µes Padr√£o" else "Personalizado",
                                "tom_voz": tom_voz if modo_geracao == "Configura√ß√µes Padr√£o" else tom_personalizado,
                                "palavras_chave": palavras_chave if modo_geracao == "Configura√ß√µes Padr√£o" else "Personalizado",
                                "numero_palavras": numero_palavras if modo_geracao == "Configura√ß√µes Padr√£o" else palavras_personalizado,
                                "conteudo_gerado": conteudo_gerado,
                                "fontes_utilizadas": {
                                    "arquivos_upload": [arquivo.name for arquivo in arquivos_upload] if arquivos_upload else [],
                                    "briefing_manual": bool(briefing_manual),
                                    "transcricoes": len(arquivos_midia) if arquivos_midia else 0
                                },
                                "data_criacao": datetime.datetime.now()
                            }
                            db_briefings['historico_geracao'].insert_one(historico_data)
                            st.success("‚úÖ Conte√∫do salvo no hist√≥rico!")
                        except Exception as e:
                            st.warning(f"Conte√∫do gerado, mas n√£o salvo no hist√≥rico: {str(e)}")
                    
                except Exception as e:
                    st.error(f"‚ùå Erro ao gerar conte√∫do: {str(e)}")
                    st.info("üí° Dica: Verifique se os arquivos n√£o est√£o corrompidos e tente novamente.")

    # Se√ß√£o de hist√≥rico r√°pido
    if mongo_connected_conteudo:
        with st.expander("üìö Hist√≥rico de Gera√ß√µes Recentes"):
            try:
                historico = list(db_briefings['historico_geracao'].find().sort("data_criacao", -1).limit(5))
                if historico:
                    for item in historico:
                        st.write(f"**{item['tipo_conte√∫do']}** - {item['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                        st.caption(f"Palavras-chave: {item.get('palavras_chave', 'Nenhuma')} | Tom: {item['tom_voz']}")
                        with st.expander("Ver conte√∫do"):
                            st.write(item['conteudo_gerado'][:500] + "..." if len(item['conteudo_gerado']) > 500 else item['conteudo_gerado'])
                else:
                    st.info("Nenhuma gera√ß√£o no hist√≥rico")
            except Exception as e:
                st.warning(f"Erro ao carregar hist√≥rico: {str(e)}")



# --- FUN√á√ïES DE REVIS√ÉO ORTOGR√ÅFICA ---

def revisar_texto_ortografia(texto, agente, segmentos_selecionados, revisao_estilo=True, manter_estrutura=True, explicar_alteracoes=True):
    """
    Realiza revis√£o ortogr√°fica e gramatical do texto considerando as diretrizes do agente
    usando a API do Gemini
    """
    
    # Construir o contexto do agente
    contexto_agente = "CONTEXTO DO AGENTE PARA REVIS√ÉO:\n\n"
    
    if "system_prompt" in segmentos_selecionados and "system_prompt" in agente:
        contexto_agente += f"DIRETRIZES PRINCIPAIS:\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and "base_conhecimento" in agente:
        contexto_agente += f"BASE DE CONHECIMENTO:\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and "comments" in agente:
        contexto_agente += f"COMENT√ÅRIOS E OBSERVA√á√ïES:\n{agente['comments']}\n\n"
    
    if "planejamento" in segmentos_selecionados and "planejamento" in agente:
        contexto_agente += f"PLANEJAMENTO E ESTRAT√âGIA:\n{agente['planejamento']}\n\n"
    
    # Construir instru√ß√µes baseadas nas configura√ß√µes
    instrucoes_revisao = ""
    
    if revisao_estilo:
        instrucoes_revisao += """
        - Analise e melhore a clareza, coes√£o e coer√™ncia textual
        - Verifique adequa√ß√£o ao tom da marca
        - Elimine v√≠cios de linguagem e redund√¢ncias
        - Simplifique frases muito longas ou complexas
        """
    
    if manter_estrutura:
        instrucoes_revisao += """
        - Mantenha a estrutura geral do texto original
        - Preserve par√°grafos e se√ß√µes quando poss√≠vel
        - Conserve o fluxo l√≥gico do conte√∫do
        """
    
    if explicar_alteracoes:
        instrucoes_revisao += """
        - Inclua justificativa para as principais altera√ß√µes
        - Explique corre√ß√µes gramaticais importantes
        - Destaque melhorias de estilo significativas
        """
    
    # Construir o prompt para revis√£o
    prompt_revisao = f"""
    {contexto_agente}
    
    TEXTO PARA REVIS√ÉO:
    {texto}
    
    INSTRU√á√ïES PARA REVIS√ÉO:
    
    1. **REVIS√ÉO ORTOGR√ÅFICA E GRAMATICAL:**
       - Corrija erros de ortografia, acentua√ß√£o e grafia
       - Verifique concord√¢ncia nominal e verbal
       - Ajuste pontua√ß√£o (v√≠rgulas, pontos, travess√µes)
       - Corrija reg√™ncia verbal e nominal
       - Ajuste coloca√ß√£o pronominal
    
    2. **REVIS√ÉO DE ESTILO E CLAREZA:**
       {instrucoes_revisao}
    
    3. **CONFORMIDADE COM AS DIRETRIZES:**
       - Alinhe o texto ao tom e estilo definidos
       - Mantenha consist√™ncia terminol√≥gica
       - Preserve a estrutura original quando poss√≠vel
       - Adapte ao p√∫blico-alvo definido
    
    FORMATO DA RESPOSTA:
    
    ## üìã TEXTO REVISADO
    [Aqui vai o texto completo revisado, mantendo a estrutura geral quando poss√≠vel]
    
    ## üîç PRINCIPAIS ALTERA√á√ïES REALIZADAS
    [Lista das principais corre√ß√µes realizadas com justificativa]
    
    
    **IMPORTANTE:**
    - Seja detalhado e preciso nas explica√ß√µes
    - Mantenha o formato markdown para f√°cil leitura
    - Inclua exemplos espec√≠ficos quando relevante
    - Foque nas corre√ß√µes ortogr√°ficas e gramaticais
    """
    
    try:
        # Chamar a API do Gemini
        response = modelo_texto.generate_content(prompt_revisao)
        
        if response and response.text:
            return response.text
        else:
            return "‚ùå Erro: N√£o foi poss√≠vel gerar a revis√£o. Tente novamente."
        
    except Exception as e:
        return f"‚ùå Erro durante a revis√£o: {str(e)}"

def revisar_documento_por_slides(doc, agente, segmentos_selecionados, revisao_estilo=True, explicar_alteracoes=True):
    """Revisa documento slide por slide com an√°lise detalhada"""
    
    resultados = []
    
    for i, slide in enumerate(doc['slides']):
        with st.spinner(f"Revisando slide {i+1} de {len(doc['slides'])}..."):
            try:
                # Construir contexto do agente para este slide
                contexto_agente = "CONTEXTO DO AGENTE PARA REVIS√ÉO:\n\n"
                
                if "system_prompt" in segmentos_selecionados and "system_prompt" in agente:
                    contexto_agente += f"DIRETRIZES PRINCIPAIS:\n{agente['system_prompt']}\n\n"
                
                if "base_conhecimento" in segmentos_selecionados and "base_conhecimento" in agente:
                    contexto_agente += f"BASE DE CONHECIMENTO:\n{agente['base_conhecimento']}\n\n"
                
                prompt_slide = f"""
{contexto_agente}

## REVIS√ÉO ORTOGR√ÅFICA - SLIDE {i+1}

**CONTE√öDO DO SLIDE {i+1}:**
{slide['conteudo'][:1500]}

**INSTRU√á√ïES:**
- Fa√ßa uma revis√£o ortogr√°fica e gramatical detalhada
- Corrija erros de portugu√™s, acentua√ß√£o e pontua√ß√£o
- Mantenha o conte√∫do original sempre que poss√≠vel
- { "Inclua sugest√µes de melhoria de estilo" if revisao_estilo else "Foque apenas em corre√ß√µes gramaticais" }
- { "Explique as principais altera√ß√µes" if explicar_alteracoes else "Apenas apresente o texto corrigido" }

**FORMATO DE RESPOSTA:**

### üìã SLIDE {i+1} - TEXTO REVISADO
[Texto corrigido do slide]

### üîç ALTERA√á√ïES REALIZADAS
- [Lista das corre√ß√µes com explica√ß√£o]

### ‚úÖ STATUS
[‚úîÔ∏è Sem erros / ‚ö†Ô∏è Pequenos ajustes / ‚ùå Corre√ß√µes necess√°rias]
"""
                
                resposta = modelo_texto.generate_content(prompt_slide)
                resultados.append({
                    'slide_num': i+1,
                    'analise': resposta.text,
                    'tem_alteracoes': '‚ùå' in resposta.text or '‚ö†Ô∏è' in resposta.text or 'Corre√ß√µes' in resposta.text
                })
                
            except Exception as e:
                resultados.append({
                    'slide_num': i+1,
                    'analise': f"‚ùå Erro na revis√£o do slide: {str(e)}",
                    'tem_alteracoes': False
                })
    
    # Construir relat√≥rio consolidado
    relatorio = f"# üìä RELAT√ìRIO DE REVIS√ÉO ORTOGR√ÅFICA - {doc['nome']}\n\n"
    relatorio += f"**Total de Slides:** {len(doc['slides'])}\n"
    relatorio += f"**Slides com Corre√ß√µes:** {sum(1 for r in resultados if r['tem_alteracoes'])}\n\n"
    
    # Slides que precisam de aten√ß√£o
    slides_com_correcoes = [r for r in resultados if r['tem_alteracoes']]
    if slides_com_correcoes:
        relatorio += "## üö® SLIDES COM CORRE√á√ïES:\n\n"
        for resultado in slides_com_correcoes:
            relatorio += f"### üìã Slide {resultado['slide_num']}\n"
            relatorio += f"{resultado['analise']}\n\n"
    
    # Resumo executivo
    relatorio += "## üìà RESUMO EXECUTIVO\n\n"
    if slides_com_correcoes:
        relatorio += f"**‚ö†Ô∏è {len(slides_com_correcoes)} slide(s) necessitam de corre√ß√µes**\n"
        relatorio += f"**‚úÖ {len(doc['slides']) - len(slides_com_correcoes)} slide(s) est√£o corretos**\n"
        
        # Lista resumida de problemas
        relatorio += "\n**üìù PRINCIPAIS TIPOS DE CORRE√á√ïES:**\n"
        problemas_comuns = []
        for resultado in slides_com_correcoes:
            if "ortogr√°fico" in resultado['analise'].lower():
                problemas_comuns.append("Erros ortogr√°ficos")
            if "pontua√ß√£o" in resultado['analise'].lower():
                problemas_comuns.append("Problemas de pontua√ß√£o")
            if "concord√¢ncia" in resultado['analise'].lower():
                problemas_comuns.append("Erros de concord√¢ncia")
        
        problemas_unicos = list(set(problemas_comuns))
        for problema in problemas_unicos:
            relatorio += f"- {problema}\n"
    else:
        relatorio += "**üéâ Todos os slides est√£o ortograficamente corretos!**\n"
    
    return relatorio

# --- ABA: REVIS√ÉO ORTOGR√ÅFICA ---
with tab_mapping["üìù Revis√£o Ortogr√°fica"]:
    st.header("üìù Revis√£o Ortogr√°fica e Gramatical")
    
    if not st.session_state.agente_selecionado:
        st.info("Selecione um agente primeiro na aba de Chat")
    else:
        agente = st.session_state.agente_selecionado
        st.subheader(f"Revis√£o com: {agente['nome']}")
        
        # Configura√ß√µes de segmentos para revis√£o
        st.sidebar.subheader("üîß Configura√ß√µes de Revis√£o")
        st.sidebar.write("Selecione bases para orientar a revis√£o:")
        
        segmentos_revisao = st.sidebar.multiselect(
            "Bases para revis√£o:",
            options=["system_prompt", "base_conhecimento", "comments", "planejamento"],
            default=st.session_state.segmentos_selecionados,
            key="revisao_segmentos"
        )
        
        # Layout em abas para diferentes m√©todos de entrada
        tab_texto, tab_arquivo = st.tabs(["üìù Texto Direto", "üìé Upload de Arquivos"])
        
        with tab_texto:
            # Layout em colunas para texto direto
            col_original, col_resultado = st.columns(2)
            
            with col_original:
                st.subheader("üìÑ Texto Original")
                
                texto_para_revisao = st.text_area(
                    "Cole o texto que deseja revisar:",
                    height=400,
                    placeholder="Cole aqui o texto que precisa de revis√£o ortogr√°fica e gramatical...",
                    help="O texto ser√° analisado considerando as diretrizes do agente selecionado",
                    key="texto_revisao"
                )
                
                # Estat√≠sticas do texto
                if texto_para_revisao:
                    palavras = len(texto_para_revisao.split())
                    caracteres = len(texto_para_revisao)
                    paragrafos = texto_para_revisao.count('\n\n') + 1
                    
                    col_stats1, col_stats2, col_stats3 = st.columns(3)
                    with col_stats1:
                        st.metric("üìä Palavras", palavras)
                    with col_stats2:
                        st.metric("üî§ Caracteres", caracteres)
                    with col_stats3:
                        st.metric("üìÑ Par√°grafos", paragrafos)
                
                # Configura√ß√µes de revis√£o
                with st.expander("‚öôÔ∏è Configura√ß√µes da Revis√£o"):
                    revisao_estilo = st.checkbox(
                        "Incluir revis√£o de estilo",
                        value=True,
                        help="Analisar clareza, coes√£o e adequa√ß√£o ao tom da marca",
                        key="revisao_estilo"
                    )
                    
                    manter_estrutura = st.checkbox(
                        "Manter estrutura original",
                        value=True,
                        help="Preservar a estrutura geral do texto quando poss√≠vel",
                        key="manter_estrutura"
                    )
                    
                    explicar_alteracoes = st.checkbox(
                        "Explicar altera√ß√µes principais",
                        value=True,
                        help="Incluir justificativa para as mudan√ßas mais importantes",
                        key="explicar_alteracoes"
                    )
            
            with col_resultado:
                st.subheader("üìã Resultado da Revis√£o")
                
                if st.button("üîç Realizar Revis√£o Completa", type="primary", key="revisar_texto"):
                    if not texto_para_revisao.strip():
                        st.warning("‚ö†Ô∏è Por favor, cole o texto que deseja revisar.")
                    else:
                        with st.spinner("üîÑ Analisando texto e realizando revis√£o..."):
                            try:
                                resultado = revisar_texto_ortografia(
                                    texto=texto_para_revisao,
                                    agente=agente,
                                    segmentos_selecionados=segmentos_revisao,
                                    revisao_estilo=revisao_estilo,
                                    manter_estrutura=manter_estrutura,
                                    explicar_alteracoes=explicar_alteracoes
                                )
                                
                                st.markdown(resultado)
                                
                                # Op√ß√µes de download
                                col_dl1, col_dl2, col_dl3 = st.columns(3)
                                
                                with col_dl1:
                                    st.download_button(
                                        "üíæ Baixar Relat√≥rio Completo",
                                        data=resultado,
                                        file_name=f"relatorio_revisao_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                        mime="text/plain",
                                        key="download_revisao_completo"
                                    )
                                
                                with col_dl2:
                                    # Extrair apenas o texto revisado se dispon√≠vel
                                    if "## üìã TEXTO REVISADO" in resultado:
                                        texto_revisado_start = resultado.find("## üìã TEXTO REVISADO")
                                        texto_revisado_end = resultado.find("##", texto_revisado_start + 1)
                                        texto_revisado = resultado[texto_revisado_start:texto_revisado_end] if texto_revisado_end != -1 else resultado[texto_revisado_start:]
                                        
                                        st.download_button(
                                            "üìÑ Baixar Texto Revisado",
                                            data=texto_revisado,
                                            file_name=f"texto_revisado_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                            mime="text/plain",
                                            key="download_texto_revisado"
                                        )
                                
                                with col_dl3:
                                    # Extrair apenas as explica√ß√µes se dispon√≠vel
                                    if "## üîç PRINCIPAIS ALTERA√á√ïES REALIZADAS" in resultado:
                                        explicacoes_start = resultado.find("## üîç PRINCIPAIS ALTERA√á√ïES REALIZADAS")
                                        explicacoes_end = resultado.find("##", explicacoes_start + 1)
                                        explicacoes = resultado[explicacoes_start:explicacoes_end] if explicacoes_end != -1 else resultado[explicacoes_start:]
                                        
                                        st.download_button(
                                            "üìù Baixar Explica√ß√µes",
                                            data=explicacoes,
                                            file_name=f"explicacoes_revisao_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                                            mime="text/plain",
                                            key="download_explicacoes"
                                        )
                                
                            except Exception as e:
                                st.error(f"‚ùå Erro ao realizar revis√£o: {str(e)}")
        
        with tab_arquivo:
            st.subheader("üìé Upload de Arquivos para Revis√£o")
            
            # Upload de m√∫ltiplos arquivos
            arquivos_upload = st.file_uploader(
                "Selecione arquivos para revis√£o:",
                type=['pdf', 'pptx', 'txt', 'docx'],
                accept_multiple_files=True,
                help="Arquivos ser√£o convertidos para texto e revisados ortograficamente",
                key="arquivos_revisao"
            )
            
            # Configura√ß√µes para arquivos
            with st.expander("‚öôÔ∏è Configura√ß√µes da Revis√£o para Arquivos"):
                analise_por_slide = st.checkbox(
                    "An√°lise detalhada por slide/p√°gina",
                    value=True,
                    help="Analisar cada slide/p√°gina individualmente",
                    key="analise_por_slide"
                )
                
                revisao_estilo_arquivos = st.checkbox(
                    "Incluir revis√£o de estilo",
                    value=True,
                    help="Analisar clareza, coes√£o e adequa√ß√£o ao tom da marca",
                    key="revisao_estilo_arquivos"
                )
                
                explicar_alteracoes_arquivos = st.checkbox(
                    "Explicar altera√ß√µes principais",
                    value=True,
                    help="Incluir justificativa para as mudan√ßas mais importantes",
                    key="explicar_alteracoes_arquivos"
                )
            
            if arquivos_upload:
                st.success(f"‚úÖ {len(arquivos_upload)} arquivo(s) carregado(s)")
                
                # Mostrar preview dos arquivos
                with st.expander("üìã Visualizar Arquivos Carregados", expanded=False):
                    for i, arquivo in enumerate(arquivos_upload):
                        st.write(f"**{arquivo.name}** ({arquivo.size} bytes)")
                
                if st.button("üîç Revisar Todos os Arquivos", type="primary", key="revisar_arquivos"):
                    resultados_completos = []
                    
                    for arquivo in arquivos_upload:
                        with st.spinner(f"Processando {arquivo.name}..."):
                            try:
                                # Extrair texto do arquivo
                                texto_extraido = ""
                                slides_info = []
                                
                                if arquivo.type == "application/pdf":
                                    texto_extraido, slides_info = extract_text_from_pdf_com_slides(arquivo)
                                elif arquivo.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                                    texto_extraido, slides_info = extract_text_from_pptx_com_slides(arquivo)
                                elif arquivo.type == "text/plain":
                                    texto_extraido = extrair_texto_arquivo(arquivo)
                                elif arquivo.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                                    texto_extraido = extrair_texto_arquivo(arquivo)
                                else:
                                    st.warning(f"Tipo de arquivo n√£o suportado: {arquivo.name}")
                                    continue
                                
                                if texto_extraido and len(texto_extraido.strip()) > 0:
                                    doc_info = {
                                        'nome': arquivo.name,
                                        'conteudo': texto_extraido,
                                        'slides': slides_info,
                                        'tipo': arquivo.type
                                    }
                                    
                                    # Escolher o m√©todo de revis√£o baseado nas configura√ß√µes
                                    if analise_por_slide and slides_info:
                                        # Revis√£o detalhada por slide
                                        resultado = revisar_documento_por_slides(
                                            doc_info,
                                            agente,
                                            segmentos_revisao,
                                            revisao_estilo_arquivos,
                                            explicar_alteracoes_arquivos
                                        )
                                    else:
                                        # Revis√£o geral do documento
                                        resultado = revisar_texto_ortografia(
                                            texto=texto_extraido,
                                            agente=agente,
                                            segmentos_selecionados=segmentos_revisao,
                                            revisao_estilo=revisao_estilo_arquivos,
                                            manter_estrutura=True,
                                            explicar_alteracoes=explicar_alteracoes_arquivos
                                        )
                                    
                                    resultados_completos.append({
                                        'nome': arquivo.name,
                                        'texto_original': texto_extraido,
                                        'resultado': resultado,
                                        'tipo': 'por_slide' if (analise_por_slide and slides_info) else 'geral'
                                    })
                                    
                                    # Exibir resultado individual
                                    with st.expander(f"üìÑ Resultado - {arquivo.name}", expanded=False):
                                        st.markdown(resultado)
                                        
                                        # Estat√≠sticas do arquivo processado
                                        palavras_orig = len(texto_extraido.split())
                                        st.info(f"üìä Arquivo original: {palavras_orig} palavras")
                                        if slides_info:
                                            st.info(f"üìë {len(slides_info)} slides/p√°ginas processados")
                                        
                                else:
                                    st.warning(f"‚ùå N√£o foi poss√≠vel extrair texto do arquivo: {arquivo.name}")
                                
                            except Exception as e:
                                st.error(f"‚ùå Erro ao processar {arquivo.name}: {str(e)}")
                    
                    # Bot√£o para download de todos os resultados
                    if resultados_completos:
                        st.markdown("---")
                        st.subheader("üì¶ Download de Todos os Resultados")
                        
                        # Criar relat√≥rio consolidado
                        relatorio_consolidado = f"# RELAT√ìRIO DE REVIS√ÉO ORTOGR√ÅFICA\n\n"
                        relatorio_consolidado += f"**Data:** {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}\n"
                        relatorio_consolidado += f"**Agente:** {agente['nome']}\n"
                        relatorio_consolidado += f"**Total de Arquivos:** {len(resultados_completos)}\n\n"
                        
                        for resultado in resultados_completos:
                            relatorio_consolidado += f"## üìÑ {resultado['nome']}\n\n"
                            relatorio_consolidado += f"{resultado['resultado']}\n\n"
                            relatorio_consolidado += "---\n\n"
                        
                        st.download_button(
                            "üíæ Baixar Relat√≥rio Consolidado",
                            data=relatorio_consolidado,
                            file_name=f"relatorio_revisao_arquivos_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
                            mime="text/plain",
                            key="download_consolidado"
                        )
            
            else:
                st.info("""
                **üìé Como usar o upload de arquivos:**
                
                1. Selecione um ou mais arquivos (PDF, PPTX, TXT, DOCX)
                2. Configure as op√ß√µes de revis√£o
                3. Clique em **"Revisar Todos os Arquivos"**
                
                **üìã Formatos suportados:**
                - PDF (documentos, apresenta√ß√µes) - com an√°lise por p√°gina
                - PPTX (apresenta√ß√µes PowerPoint) - com an√°lise por slide
                - TXT (arquivos de texto)
                - DOCX (documentos Word)
                
                **üîç An√°lise por Slide/P√°gina:**
                - Identifica slides/p√°ginas espec√≠ficos com problemas
                - Revis√£o detalhada de cada se√ß√£o
                - Facilita a localiza√ß√£o e corre√ß√£o de erros
                """)
        
        # Se√ß√£o informativa
        with st.expander("‚ÑπÔ∏è Sobre a Revis√£o Ortogr√°fica"):
            st.markdown("""
            ### üéØ O que √© Analisado
            
            **üî§ Ortografia:**
            - Erros de grafia e acentua√ß√£o
            - Uso correto de mai√∫sculas e min√∫sculas
            - Escrita de n√∫meros e datas
            - Concord√¢ncia nominal e verbal
            
            **üìñ Gram√°tica:**
            - Estrutura sint√°tica das frases
            - Uso adequado de preposi√ß√µes
            - Coloca√ß√£o pronominal
            - Reg√™ncia verbal e nominal
            
            **üî† Pontua√ß√£o:**
            - Uso de v√≠rgulas, pontos, dois-pontos
            - Aplica√ß√£o de travess√µes e par√™nteses
            - Pontua√ß√£o de cita√ß√µes e di√°logos
            
            **üìù Estilo e Clareza:**
            - Coes√£o e coer√™ncia textual
            - Adequa√ß√£o ao tom da marca
            - Clareza na comunica√ß√£o
            - Elimina√ß√£o de v√≠cios de linguagem
            
            ### üí° Dicas para Melhor Revis√£o
            
            1. **Texto Completo**: Cole o texto integral para an√°lise detalhada
            2. **Segmentos Relevantes**: Selecione as bases de conhecimento apropriadas
            3. **Contexto Espec√≠fico**: Use agentes especializados para cada tipo de conte√∫do
            4. **Implementa√ß√£o**: Aplique as sugest√µes sistematicamente
            
            ### üé® Benef√≠cios da Revis√£o Contextual
            
            - **Consist√™ncia da Marca**: Mant√©m o tom e estilo adequados
            - **Qualidade Profissional**: Elimina erros que prejudicam a credibilidade
            - **Otimiza√ß√£o de Conte√∫do**: Melhora a clareza e impacto da comunica√ß√£o
            - **Efici√™ncia**: Reduz tempo de revis√£o manual
            """)

# --- ABA: MONITORAMENTO DE REDES ---
with tab_mapping["Monitoramento de Redes"]:
    st.header("ü§ñ Agente de Monitoramento")
    st.markdown("**Especialista que fala como gente**")

    def gerar_resposta_agente(pergunta_usuario: str, historico: List[Dict] = None, agente_monitoramento=None) -> str:
        """Gera resposta do agente usando RAG e base do agente de monitoramento"""
        
        # Configura√ß√£o do agente - usa base do agente selecionado ou padr√£o
        if agente_monitoramento and agente_monitoramento.get('base_conhecimento'):
            system_prompt = agente_monitoramento['base_conhecimento']
        else:
            # Fallback para prompt padr√£o se n√£o houver agente selecionado
            system_prompt = """
            PERSONALIDADE: Especialista t√©cnico do agroneg√≥cio com habilidade social - "Especialista que fala como gente"

            TOM DE VOZ:
            - T√©cnico, confi√°vel e seguro, mas acess√≠vel
            - Evita exageros e promessas vazias
            - Sempre embasado em fatos e ci√™ncia
            - Frases curtas e diretas, mais simp√°ticas
            - Toque de leveza e ironia pontual quando o contexto permite

            DIRETRIZES:
            - N√ÉO inventar informa√ß√µes t√©cnicas
            - Sempre basear respostas em fatos
            - Manter tom profissional mas acess√≠vel
            - Adaptar resposta ao tipo de pergunta
            """
        
        # Constr√≥i o prompt final
        prompt_final = f"""
        {system_prompt}
        
        
        PERGUNTA DO USU√ÅRIO:
        {pergunta_usuario}
        
        HIST√ìRICO DA CONVERSA (se aplic√°vel):
        {historico if historico else "Nenhum hist√≥rico anterior"}
        
        INSTRU√á√ïES FINAIS:
        Adapte seu tom ao tipo de pergunta:
        - Perguntas t√©cnicas: seja preciso e did√°tico
        - Perguntas sociais: seja leve e engajador  
        - Cr√≠ticas ou problemas: seja construtivo e proativo
        
        Sua resposta deve refletir a personalidade do "especialista que fala como gente".
        """
        
        try:
            resposta = modelo_texto.generate_content(prompt_final)
            return resposta.text
        except Exception as e:
            return f"Erro ao gerar resposta: {str(e)}"

    # SELE√á√ÉO DE AGENTE DE MONITORAMENTO
    st.header("üîß Configura√ß√£o do Agente de Monitoramento")
    
    # Carregar apenas agentes de monitoramento
    agentes_monitoramento = [agente for agente in listar_agentes() if agente.get('categoria') == 'Monitoramento']
    
    col_sel1, col_sel2 = st.columns([3, 1])
    
    with col_sel1:
        if agentes_monitoramento:
            # Criar op√ß√µes para selectbox
            opcoes_agentes = {f"{agente['nome']}": agente for agente in agentes_monitoramento}
            
            agente_selecionado_nome = st.selectbox(
                "Selecione o agente de monitoramento:",
                list(opcoes_agentes.keys()),
                key="seletor_monitoramento"
            )
            
            agente_monitoramento = opcoes_agentes[agente_selecionado_nome]
            
            # Mostrar informa√ß√µes do agente selecionado
            with st.expander("üìã Informa√ß√µes do Agente Selecionado", expanded=False):
                if agente_monitoramento.get('base_conhecimento'):
                    st.text_area(
                        "Base de Conhecimento:",
                        value=agente_monitoramento['base_conhecimento'],
                        height=200,
                        disabled=True
                    )
                else:
                    st.warning("‚ö†Ô∏è Este agente n√£o possui base de conhecimento configurada")
                
                st.write(f"**Criado em:** {agente_monitoramento['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                # Mostrar propriet√°rio se for admin
                if get_current_user() == "admin" and agente_monitoramento.get('criado_por'):
                    st.write(f"**üë§ Propriet√°rio:** {agente_monitoramento['criado_por']}")
        
        else:
            st.error("‚ùå Nenhum agente de monitoramento encontrado.")
            st.info("üí° Crie um agente de monitoramento na aba 'Gerenciar Agentes' primeiro.")
            agente_monitoramento = None
    
    with col_sel2:
        if st.button("üîÑ Atualizar Lista", key="atualizar_monitoramento"):
            st.rerun()

    # Sidebar com informa√ß√µes
    with st.sidebar:
        st.header("‚ÑπÔ∏è Sobre o Monitoramento")
        
        if agente_monitoramento:
            st.success(f"**Agente Ativo:** {agente_monitoramento['nome']}")
        else:
            st.warning("‚ö†Ô∏è Nenhum agente selecionado")
        
        st.markdown("""
        **Personalidade:**
        - üéØ T√©cnico mas acess√≠vel
        - üí¨ Direto mas simp√°tico
        - üå± Conhece o campo e a internet
        - üî¨ Baseado em ci√™ncia
        
        **Capacidades:**
        - Respostas t√©cnicas baseadas em RAG
        - Engajamento em redes sociais
        - Suporte a produtores
        - Esclarecimento de d√∫vidas
        """)

        
        if st.button("üîÑ Reiniciar Conversa", key="reiniciar_monitoramento"):
            if "messages_monitoramento" in st.session_state:
                st.session_state.messages_monitoramento = []
            st.rerun()

        # Status da conex√£o
        
        if os.getenv('OPENAI_API_KEY'):
            st.success("‚úÖ OpenAI: Configurado")
        else:
            st.warning("‚ö†Ô∏è OpenAI: N√£o configurado")

    # Inicializar hist√≥rico de mensagens espec√≠fico para monitoramento
    if "messages_monitoramento" not in st.session_state:
        st.session_state.messages_monitoramento = []

    # √Årea de chat principal
    st.header("üí¨ Simulador de Respostas do Agente")

    # Exemplos de perguntas r√°pidas
    st.subheader("üéØ Exemplos para testar:")
    col1, col2, col3 = st.columns(3)

    with col1:
        if st.button("‚ùì D√∫vida T√©cnica", use_container_width=True, key="exemplo_tecnico"):
            st.session_state.messages_monitoramento.append({"role": "user", "content": "Esse produto serve pra todas as culturas?"})

    with col2:
        if st.button("üòä Coment√°rio Social", use_container_width=True, key="exemplo_social"):
            st.session_state.messages_monitoramento.append({"role": "user", "content": "O campo t√° lindo demais!"})

    with col3:
        if st.button("‚ö†Ô∏è Cr√≠tica/Problema", use_container_width=True, key="exemplo_critica"):
            st.session_state.messages_monitoramento.append({"role": "user", "content": "Usei e n√£o funcionou."})

    # Exibir hist√≥rico de mensagens
    for message in st.session_state.messages_monitoramento:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Input do usu√°rio
    if prompt := st.chat_input("Digite sua mensagem ou pergunta...", key="chat_monitoramento"):
        # Adicionar mensagem do usu√°rio
        st.session_state.messages_monitoramento.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Gerar resposta do agente
        with st.chat_message("assistant"):
            with st.spinner("üå± Consultando base t√©cnica..."):
                resposta = gerar_resposta_agente(
                    prompt, 
                    st.session_state.messages_monitoramento,
                    agente_monitoramento
                )
                st.markdown(resposta)
                
                # Adicionar ao hist√≥rico
                st.session_state.messages_monitoramento.append({"role": "assistant", "content": resposta})

    # Se√ß√£o de an√°lise de performance
    st.markdown("---")
    st.header("üìä An√°lise da Resposta")

    if st.session_state.messages_monitoramento:
        ultima_resposta = st.session_state.messages_monitoramento[-1]["content"] if st.session_state.messages_monitoramento[-1]["role"] == "assistant" else ""
        
        if ultima_resposta:
            col_analise1, col_analise2, col_analise3 = st.columns(3)
            
            with col_analise1:
                # An√°lise de tom
                if "üòä" in ultima_resposta or "üòç" in ultima_resposta:
                    st.metric("Tom Identificado", "Social/Engajador", delta="Leve")
                elif "üî¨" in ultima_resposta or "üìä" in ultima_resposta:
                    st.metric("Tom Identificado", "T√©cnico", delta="Preciso")
                else:
                    st.metric("Tom Identificado", "Balanceado", delta="Adaptado")
            
            with col_analise2:
                # Comprimento da resposta
                palavras = len(ultima_resposta.split())
                st.metric("Tamanho", f"{palavras} palavras")
            
            with col_analise3:
                # Uso de emojis
                emojis = sum(1 for char in ultima_resposta if char in "üòÄüòÉüòÑüòÅüòÜüòÖüòÇü§£‚ò∫Ô∏èüòäüòáüôÇüôÉüòâüòåüòçü•∞üòòüòóüòôüòöüòãüòõüòùüòúü§™ü§®üßêü§ìüòéü§©ü•≥üòèüòíüòûüòîüòüüòïüôÅ‚òπÔ∏èüò£üòñüò´üò©ü•∫üò¢üò≠üò§üò†üò°ü§¨ü§Øüò≥ü•µü•∂üò±üò®üò∞üò•üòìü§óü§îü§≠ü§´ü§•üò∂üòêüòëüò¨üôÑüòØüò¶üòßüòÆüò≤ü•±üò¥ü§§üò™üòµü§êü•¥ü§¢ü§Æü§ßüò∑ü§íü§ïü§ëü§†üòàüëøüëπüë∫ü§°üí©üëªüíÄ‚ò†Ô∏èüëΩüëæü§ñüéÉüò∫üò∏üòπüòªüòºüòΩüôÄüòøüòæ")
                st.metric("Emojis", emojis, delta="Moderado" if emojis <= 2 else "Alto")

    # Se√ß√£o de exemplos de uso
    with st.expander("üìã Exemplos de Respostas do Agente"):
        st.markdown("""
        **üéØ PERGUNTA T√âCNICA:**
        *Usu√°rio:* "Qual a diferen√ßa entre os nematoides de galha e de cisto na soja?"
        
        **ü§ñ AGENTE:** "Boa pergunta! Os nematoides de galha (Meloidogyne) formam aquelas 'incha√ß√µes' nas ra√≠zes, enquanto os de cisto (Heterodera) ficam mais externos. Ambos roubam nutrientes, mas o manejo pode ser diferente. Temos solu√ß√µes espec√≠ficas para cada caso! üå±"
        
        **üéØ COMENT√ÅRIO SOCIAL:**
        *Usu√°rio:* "Adorei ver as fotos da lavoura no stories!"
        
        **ü§ñ AGENTE:** "A gente tamb√©m ama compartilhar esses momentos! Quando a tecnologia encontra o cuidado certo, o campo fica ainda mais bonito üòç Compartilhe suas fotos tamb√©m!"
        
        **üéØ CR√çTICA/PROBLEMA:**
        *Usu√°rio:* "A aplica√ß√£o n√£o deu o resultado esperado"
        
        **ü§ñ AGENTE:** "Poxa, que pena saber disso! Vamos entender melhor o que aconteceu. Pode me contar sobre as condi√ß√µes de aplica√ß√£o? Assim conseguimos te orientar melhor da pr√≥xima vez. A equipe t√©cnica tamb√©m est√° √† disposi√ß√£o! üìû"
        """)

# --- Fun√ß√µes auxiliares para busca web ---
def buscar_perplexity(pergunta: str, contexto_agente: str = None) -> str:
    """Realiza busca na web usando API do Perplexity"""
    try:
        headers = {
            "Authorization": f"Bearer {perp_api_key}",
            "Content-Type": "application/json"
        }
        
        # Construir o conte√∫do da mensagem
        messages = []
        
        if contexto_agente:
            messages.append({
                "role": "system",
                "content": f"Contexto do agente: {contexto_agente}"
            })
        
        messages.append({
            "role": "user",
            "content": pergunta
        })
        
        data = {
            "model": "sonar-medium-online",
            "messages": messages,
            "max_tokens": 2000,
            "temperature": 0.1
        }
        
        response = requests.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers,
            json=data,
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            return f"‚ùå Erro na busca: {response.status_code} - {response.text}"
            
    except Exception as e:
        return f"‚ùå Erro ao conectar com Perplexity: {str(e)}"

def analisar_urls_perplexity(urls: List[str], pergunta: str, contexto_agente: str = None) -> str:
    """Analisa URLs espec√≠ficas usando Perplexity"""
    try:
        headers = {
            "Authorization": f"Bearer {perp_api_key}",
            "Content-Type": "application/json"
        }
        
        # Construir contexto com URLs
        urls_contexto = "\n".join([f"- {url}" for url in urls])
        
        messages = []
        
        if contexto_agente:
            messages.append({
                "role": "system",
                "content": f"Contexto do agente: {contexto_agente}"
            })
        
        messages.append({
            "role": "user",
            "content": f"""Analise as seguintes URLs e responda √† pergunta:

URLs para an√°lise:
{urls_contexto}

Pergunta: {pergunta}

Forne√ßa uma an√°lise detalhada baseada no conte√∫do dessas URLs."""
        })
        
        data = {
            "model": "sonar-medium-online",
            "messages": messages,
            "max_tokens": 3000,
            "temperature": 0.1
        }
        
        response = requests.post(
            "https://api.perplexity.ai/chat/completions",
            headers=headers,
            json=data,
            timeout=45
        )
        
        if response.status_code == 200:
            result = response.json()
            return result['choices'][0]['message']['content']
        else:
            return f"‚ùå Erro na an√°lise: {response.status_code} - {response.text}"
            
    except Exception as e:
        return f"‚ùå Erro ao analisar URLs: {str(e)}"

def transcrever_audio_video(arquivo, tipo):
    """Fun√ß√£o placeholder para transcri√ß√£o de √°udio/v√≠deo"""
    return f"Transcri√ß√£o do {tipo} {arquivo.name} - Esta funcionalidade requer configura√ß√£o adicional de APIs de transcri√ß√£o."

# --- Estiliza√ß√£o ---
st.markdown("""
<style>
    .stChatMessage {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    [data-testid="stChatMessageContent"] {
        font-size: 1rem;
    }
    .stChatInput {
        bottom: 20px;
        position: fixed;
        width: calc(100% - 5rem);
    }
    div[data-testid="stTabs"] {
        margin-top: -30px;
    }
    div[data-testid="stVerticalBlock"] > div:has(>.stTextArea) {
        border-left: 3px solid #4CAF50;
        padding-left: 1rem;
    }
    .segment-indicator {
        background-color: #f0f2f6;
        padding: 0.5rem;
        border-radius: 0.5rem;
        margin: 0.5rem 0;
        border-left: 4px solid #4CAF50;
    }
    .video-analysis-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .inheritance-badge {
        background-color: #e3f2fd;
        color: #1976d2;
        padding: 0.2rem 0.5rem;
        border-radius: 12px;
        font-size: 0.8rem;
        margin-left: 0.5rem;
    }
    .web-search-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .seo-analysis-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .spelling-review-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .validation-unified-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .user-indicator {
        background-color: #e8f5e8;
        padding: 0.3rem 0.8rem;
        border-radius: 15px;
        font-size: 0.8rem;
        color: #2e7d32;
        border: 1px solid #c8e6c9;
        margin-left: 0.5rem;
    }
</style>
""", unsafe_allow_html=True)

# --- Informa√ß√µes do sistema na sidebar ---
with st.sidebar:
    st.markdown("---")
    st.subheader("üîê Sistema de Isolamento")
    
    current_user = get_current_user()
    if current_user == "admin":
        st.success("üëë **Modo Administrador**")
        st.info("Visualizando e gerenciando TODOS os agentes do sistema")
    else:
        st.success(f"üë§ **Usu√°rio: {current_user}**")
        st.info("Visualizando e gerenciando apenas SEUS agentes")
    
    # Estat√≠sticas r√°pidas
    agentes_usuario = listar_agentes()
    if agentes_usuario:
        categorias_count = {}
        for agente in agentes_usuario:
            cat = agente.get('categoria', 'Social')
            categorias_count[cat] = categorias_count.get(cat, 0) + 1
        
        st.markdown("### üìä Seus Agentes")
        for categoria, count in categorias_count.items():
            st.write(f"- **{categoria}:** {count} agente(s)")
        
        st.write(f"**Total:** {len(agentes_usuario)} agente(s)")

# --- Rodap√© ---
st.markdown("---")
st.caption(f"ü§ñ Agente Social v2.0 | Usu√°rio: {get_current_user()} | {datetime.datetime.now().strftime('%d/%m/%Y %H:%M')}")
